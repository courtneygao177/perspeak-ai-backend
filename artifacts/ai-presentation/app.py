import os
import json
import random
import base64
import io
import re
import traceback
import uuid
import concurrent.futures
try:
    from json_repair import repair_json as _repair_json
    _HAS_JSON_REPAIR = True
except ImportError:
    _HAS_JSON_REPAIR = False
from flask import Flask, request, jsonify, render_template, session, redirect, url_for
from werkzeug.utils import secure_filename
import audio_engine
from config.defense_knowledge_base import (
    DEFENSE_QUESTION_BANK,
    INTERRUPT_STRATEGIES,
    DEFENSE_STRATEGY_BY_TYPE,
)

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "ai-presentation-dev-secret-2024")
app.config["UPLOAD_FOLDER"] = os.path.join(os.path.dirname(__file__), "uploads")

# ── Server-side report store (avoids 4KB cookie overflow) ─────────────────────
# Large blobs (evaluation, qa_bank, answers) are written to disk as JSON files
# keyed by a UUID. Only the 36-char key is stored in the session cookie.
_REPORT_DIR = os.path.join(os.path.dirname(__file__), "uploads", "reports")

def _save_report(evaluation, qa_bank, answers):
    """Persist report blobs to disk; return key to store in session."""
    os.makedirs(_REPORT_DIR, exist_ok=True)
    key = str(uuid.uuid4())
    payload = {"evaluation": evaluation, "qa_bank": qa_bank, "answers": answers}
    with open(os.path.join(_REPORT_DIR, f"{key}.json"), "w") as fh:
        json.dump(payload, fh)
    return key

def _load_report(key):
    """Load report blobs from disk; returns None if missing/corrupt."""
    if not key:
        return None
    path = os.path.join(_REPORT_DIR, f"{key}.json")
    try:
        with open(path) as fh:
            return json.load(fh)
    except (FileNotFoundError, json.JSONDecodeError):
        return None

# ── Server-side slides store (avoids 4KB cookie overflow on large decks) ──────
_SLIDES_DIR = os.path.join(os.path.dirname(__file__), "uploads", "slide_store")

def _save_slides(slides):
    """Write slides list to disk; return the UUID key to store in session."""
    os.makedirs(_SLIDES_DIR, exist_ok=True)
    key = str(uuid.uuid4())
    with open(os.path.join(_SLIDES_DIR, f"{key}.json"), "w", encoding="utf-8") as fh:
        json.dump(slides, fh, ensure_ascii=False)
    return key

def _load_slides(session_obj):
    """Load slides from disk using session's slide_key. Falls back to MOCK_SLIDES."""
    key = session_obj.get("slide_key", "")
    if not key:
        return MOCK_SLIDES
    path = os.path.join(_SLIDES_DIR, f"{key}.json")
    try:
        with open(path, encoding="utf-8") as fh:
            return json.load(fh)
    except (FileNotFoundError, json.JSONDecodeError):
        return MOCK_SLIDES

app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024

# ── CORS: allow all origins so the proxied iframe can reach Flask ──────────────
@app.after_request
def add_cors_headers(response):
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, OPTIONS"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization"
    return response

@app.route("/x/<path:p>", methods=["OPTIONS"])
@app.route("/<path:p>", methods=["OPTIONS"])
def options_handler(p=""):
    return "", 204

ALLOWED_EXTENSIONS = {"pdf", "ppt", "pptx"}

# ─────────────────────────────────────────────
# OPENAI-COMPATIBLE UNIFIED API CLIENT SETUP
# Works with any OpenAI-compatible relay/proxy (中转站)
# Set UNIFIED_API_KEY and UNIFIED_BASE_URL in Replit Secrets
# ─────────────────────────────────────────────
try:
    from openai import OpenAI as _OpenAI
    _UNIFIED_KEY = os.environ.get("UNIFIED_API_KEY", "").strip()
    _UNIFIED_URL = os.environ.get("UNIFIED_BASE_URL", "").strip()

    # Normalise base_url: strip whitespace, then strip any trailing
    # /chat/completions or /completions so the SDK can append correctly.
    # e.g. "https://api.ohmygpt.com/v1/chat/completions" → "https://api.ohmygpt.com/v1"
    if _UNIFIED_URL:
        for _suffix in ["/chat/completions", "/completions"]:
            if _UNIFIED_URL.rstrip("/ ").endswith(_suffix):
                _UNIFIED_URL = _UNIFIED_URL.rstrip("/ ")[: -len(_suffix)]
                break
        _UNIFIED_URL = _UNIFIED_URL.rstrip("/ ")

    if _UNIFIED_KEY and _UNIFIED_URL:
        # timeout=120s — Gemini 1.5 Pro on 6-slide + Q&A eval can take >60s via relay
        _ai_client = _OpenAI(api_key=_UNIFIED_KEY, base_url=_UNIFIED_URL, timeout=120.0)
        AI_ENABLED = True
    else:
        _ai_client = None
        AI_ENABLED = False
except Exception as _e:
    _ai_client = None
    AI_ENABLED = False

# ── Multi-model routing ───────────────────────────────────────────────────────
# Each step is routed to the best-fit model via the same relay endpoint.
# Override any model via env vars in Replit Secrets.
VISION_MODEL = os.environ.get("VISION_MODEL", "claude-sonnet-5")           # Step 1: Vision
TEXT_MODEL   = os.environ.get("TEXT_MODEL",   "gpt-4o")                      # Steps 3,5,6: reasoning
EVAL_MODEL   = os.environ.get("EVAL_MODEL",   "gemini-2.5-flash")            # Steps 8,9: long-context eval
MODEL        = os.environ.get("UNIFIED_MODEL", TEXT_MODEL)                   # legacy fallback
MAX_TOKENS   = 8192


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


# ─────────────────────────────────────────────
# CHALLENGE TYPE DICTIONARIES
# ─────────────────────────────────────────────
ACADEMIC_CHALLENGE_TYPES = {
    "Methodology Weakness": "Sample size too small or research design flawed",
    "Unsupported Claim": "Assertion made without sufficient evidence or citation",
    "Causality Issue": "Correlation presented as causation",
    "Literature Gap": "Missing key prior work or theoretical framework",
    "Overgeneralization": "Conclusion extrapolated far beyond the data",
}

BUSINESS_CHALLENGE_TYPES = {
    "Market Size Doubt": "Total addressable market is unclear or overstated",
    "Monetization Issue": "Revenue model is vague or economically unsound",
    "Defensibility": "Competitive moat is weak or nonexistent",
    "Unrealistic Assumption": "Growth or adoption projections lack grounding",
    "Competitor Challenge": "Existing players are stronger than acknowledged",
}

INTERRUPTABLE_SCENARIOS = {"MBA Case Pitch"}   # Thesis Defense → post-session Q&A only

# ─────────────────────────────────────────────
# AUDIENCE PERSONA DEFINITIONS
# ─────────────────────────────────────────────
AUDIENCE_PERSONA = {
    "Professor": (
        "You are a rigorous academic professor. You speak in precise, formal academic English. "
        "You favor long, multi-clause sentences that probe methodology, logic, and evidence. "
        "You are not hostile — but you are relentless in demanding intellectual rigor. "
        "You always cite potential flaws in research design, sample validity, or causal claims."
    ),
    "VC": (
        "You are a seasoned venture capitalist with a cold, results-oriented demeanor. "
        "You speak in short, sharp sentences. You care only about: market size, revenue model, "
        "competitive moat, and capital efficiency. You are openly skeptical and press hard on "
        "unit economics, customer acquisition cost, and defensibility. You do not give compliments."
    ),
    "Classmates": (
        "You are a curious and friendly peer student. You speak informally and ask genuine, "
        "practical questions. You want to understand how this applies to real life, why it matters "
        "to ordinary people, and whether you could explain it to someone outside the field."
    ),
}

# ─────────────────────────────────────────────
# AUDIENCE-BASED QUESTION LANGUAGE LEVEL
# Controls the English proficiency level/register of the QUESTION TEXT itself
# (not the feedback). Applied to every AI-generated question prompt.
# ─────────────────────────────────────────────
AUDIENCE_LANGUAGE_LEVEL = {
    "Classmates": (
        "Phrase the question in casual, friendly, everyday English at an IELTS 5.5-6.0 level. "
        "Use simple, common vocabulary and short, straightforward sentence structures — the way a "
        "warm, approachable classmate would casually ask something. Avoid academic jargon or formal register."
    ),
    "Professor": (
        "Phrase the question in professional, academic English at an IELTS 6.5 level. "
        "Use precise scholarly vocabulary and a formal academic register appropriate for a professor, "
        "but keep sentences clear and readable — avoid needlessly obscure words, overly convoluted syntax, "
        "or excessive jargon that would make the question hard to parse."
    ),
}


def audience_language_directive(audience):
    """Return the language-level instruction for a given audience, or '' if none defined (e.g. VC)."""
    return AUDIENCE_LANGUAGE_LEVEL.get(audience, "")


DIFFICULTY_INSTRUCTIONS = {
    "Easy": (
        "You are SUPPORTIVE. Ask only 1 gentle clarifying question per slide. "
        "Acknowledge what the presenter did well before asking. Use encouraging language. "
        "Maximum follow-up rounds: 1."
    ),
    "Medium": (
        "You are BALANCED. Challenge weak points directly but professionally. "
        "Ask 1–2 follow-up rounds on a substantive vulnerability. "
        "Maximum follow-up rounds: 2."
    ),
    "Hard": (
        "You are AGGRESSIVELY SKEPTICAL. Immediately attack the weakest point. "
        "Do not compliment. Demand concrete evidence, numbers, and mechanisms. "
        "Press relentlessly with 2+ follow-up rounds. Never accept a vague answer. "
        "Maximum follow-up rounds: 2."
    ),
}

MAX_FOLLOWUP_ROUNDS = {"Easy": 1, "Medium": 2, "Hard": 2}

# ─────────────────────────────────────────────
# FALLBACK MOCK DATA (used when AI is unavailable)
# ─────────────────────────────────────────────
MOCK_SLIDES = [
    {
        "page": 1,
        "title": "Introduction",
        "content": "Our solution addresses a critical gap in the EdTech market. "
                   "78% of non-native English speakers struggle with academic presentations. "
                   "Our AI-powered platform simulates real presentation environments.",
    },
    {
        "page": 2,
        "title": "Core Data & Market Validation",
        "content": "User research with 50 university students across 3 campuses. "
                   "Results: 43% improvement in presentation confidence after 4 sessions. "
                   "Target market: 1.8 million international students in the US alone.",
    },
    {
        "page": 3,
        "title": "Conclusion & Roadmap",
        "content": "Phase 1: Core simulation engine. Phase 2: LMS integration. "
                   "Phase 3: Multilingual support. Seeking $500K seed funding.",
    },
]

MOCK_CHALLENGE = {
    "trigger_page": 2,
    "challenge_type": "Methodology Weakness",
    "initial_challenge": (
        "According to your slide 2, your user test sample size is only 50 people. "
        "How can you prove this solution is scalable and reliable for the wider market?"
    ),
}

MOCK_FOLLOWUP_POOL = [
    "But 50 people cannot represent global diversity. What is your concrete strategy to mitigate this sampling bias?",
    "Your 43% improvement metric — what was the baseline? Was there a control group?",
    "You cite 1.8 million students but only tested on 3 campuses. How do you extrapolate market fit?",
    "The EdTech space is littered with failed AI tutoring products. What makes your moat defensible?",
]

MOCK_QA_BANK = [
    {"id": 1, "question": "How does your platform handle different English accents in speech recognition?",
     "category": "Technical Feasibility", "difficulty": "Hard", "challenge_type": "Methodology Weakness"},
    {"id": 2, "question": "What is your customer acquisition cost and projected LTV at scale?",
     "category": "Business Model", "difficulty": "Medium", "challenge_type": "Monetization Issue"},
    {"id": 3, "question": "How do you protect student data privacy under FERPA?",
     "category": "Legal & Compliance", "difficulty": "Medium", "challenge_type": "Unsupported Claim"},
]

# ─────────────────────────────────────────────
# CLASS PRESENTATION: TED Q&A MATRIX
# Audience-pinned dual-track question bank.
# No random.choice — track is decided by the user-selected audience.
# ─────────────────────────────────────────────
TED_QA_MATRIX = {
    # ── Dimension 1: Twitter Test / Message Map ────────────────────────────────
    "twitter_headline": {
        "logic": "Tests ability to summarise the core idea in ≤140 characters (Message Map / Twitter Test)",
        "classmate": (
            "You shared a lot of cool details. "
            "If you had to tweet your main point in less than 140 characters right now, what would it be? "
            "Don't expand on details — give me one punchy, social-media-style sentence. "
            "(引导：Twitter化核心提炼。不要展开细节，用一句像社交媒体文案一样短小精悍的话总结)"
        ),
        "professor": (
            "Your presentation covered various facts, but I want to test your message clarity. "
            "Can you compress the absolute core argument of your research into one single tweet-length sentence — "
            "no technical terms, just the essence? "
            "(引导：结论先行。不要背完整段落，用社交媒体标题的方式告诉我你的核心发现是什么)"
        ),
    },
    # ── Dimension 2: New Knowledge Feeling ────────────────────────────────────
    "novelty_challenge": {
        "logic": "Tests whether the talk delivers a fresh perspective — new knowledge feeling over Google-level facts",
        "classmate": (
            "Honestly, I could just Google most of this. "
            "What's the one surprising thing your project found that I wouldn't expect?"
        ),
        "professor": (
            "Much of this information is readily accessible through a basic literature search. "
            "What is the most novel or unexpected insight that your research specifically contributes?"
        ),
    },
    # ── Dimension 3: Rule of Three ─────────────────────────────────────────────
    "rule_of_three": {
        "logic": "Tests whether speaker can distil the talk to three memorable chunks",
        "classmate": (
            "If our classmates wake up tomorrow morning and can only remember three key things "
            "from your presentation, which three things do you hope they will stick to?"
        ),
        "professor": (
            "Human working memory is limited, so focus is essential. "
            "If I am grading you based on your three most important arguments, how would you define them right now?"
        ),
    },
    # ── Dimension 4: 5-W Detail Excavation ────────────────────────────────────
    "five_w_detail": {
        "logic": "Tests 5-W specificity — WHEN/WHO/WHERE on case studies (Carnegie vivid-example principle)",
        "classmate": (
            "About the specific case study you mentioned, don't just give us abstract theories. "
            "Can you pinpoint the concrete details — specifically WHO was involved and WHEN this happened? "
            "(引导：Illustrative Support。不要讲空话，讲讲具体的时间、人物和发生的事情)"
        ),
        "professor": (
            "About the specific case study you mentioned, don't just give us abstract theories. "
            "Can you pinpoint the concrete details — specifically WHO was involved and WHEN this happened? "
            "Vague examples weaken academic arguments considerably. "
            "(引导：Illustrative Support。不要讲空话，给我具体的时间、人物和可验证的事实)"
        ),
    },
    # ── Dimension 5: Audience Benefit Connection ───────────────────────────────
    "audience_benefit": {
        "logic": "Tests whether the speaker connects the topic to the audience's operational benefit",
        "classmate": (
            "How does this actually affect us as students? "
            "Can you break it down into just THREE simple points so it's easy to follow? "
            "(引导：Rule of Three。不要长篇大论，强制用【第一、第二、第三】的结构来数数回答)"
        ),
        "professor": (
            "You covered the technical aspects well. But from the audience's perspective, "
            "what is the tangible, practical benefit of this research — "
            "and can you give me exactly three concrete takeaways the broader community can act on? "
            "(引导：Rule of Three + 受众利益。不要只讲学术价值，给我三个普通人能理解的实际影响)"
        ),
    },
    # ── Dimension 6: Consensus Game ────────────────────────────────────────────
    "consensus_game": {
        "logic": "Tests diplomatic persuasion — finding common ground before making a case (Carnegie consensus-building)",
        "classmate": (
            "I have reservations about your stance. "
            "Before you convince me, where do you think we can find a common ground?"
        ),
        "professor": (
            "I have reservations about the conclusions you have drawn. "
            "Before you attempt to convince me further, where do you believe our positions actually agree?"
        ),
    },
}

# Baseline fallback pool for Class Presentation (audience-neutral)
CLASS_PRES_QA_POOL = [
    {"id": 10, "question": "Can you give a real-life example that supports your main point?",
     "category": "Class Presentation", "difficulty": "Easy", "challenge_type": "Unsupported Claim"},
    {"id": 11, "question": "How does this topic connect to what we have been learning in class?",
     "category": "Class Presentation", "difficulty": "Easy", "challenge_type": "Relevance"},
    {"id": 12, "question": "What would change if your main assumption turned out to be wrong?",
     "category": "Class Presentation", "difficulty": "Medium", "challenge_type": "Causality Issue"},
    {"id": 13, "question": "Which part of your presentation do you think needed more evidence, and why?",
     "category": "Class Presentation", "difficulty": "Medium", "challenge_type": "Clarity"},
    {"id": 14, "question": "If you had one more minute, what extra detail would you add and where?",
     "category": "Class Presentation", "difficulty": "Hard", "challenge_type": "Depth"},
]

# ─────────────────────────────────────────────
# ANCHOR QUESTION POOL — Dual-Track Q&A System
# ─────────────────────────────────────────────
# When pressure_level (difficulty) >= Medium (≥2 questions), Q2 is always
# drawn from this pool. It carries a 大白话 scaffold hint embedded in the
# question text and triggers scene-specific dimension scoring in the CQ engine.
# Exactly ONE anchor per session ("有且仅有一个标准锚点题").
ANCHOR_QUESTION_POOL = {
    "class_presentation": [
        {
            "id": "anchor_cp_twitter",
            "anchor_type": "Twitter Headline",
            "question_by_audience": {
                "classmate": (
                    "Okay, cool story, but if you had to tweet your main point in like one line, "
                    "what would it be? "
                    "(引导：Twitter化核心提炼。不要展开细节，用一句像社交媒体文案一样短小精悍的话总结)"
                ),
                "professor": (
                    "You shared a lot of detail there. "
                    "If you had to condense your main point into a single tweet-length sentence, what would it be? "
                    "(引导：Twitter化核心提炼。不要展开细节，用一句像社交媒体文案一样短小精悍的话总结)"
                ),
            },
            "question_type": "anchor",
            "category": "Class Presentation — Anchor",
            "challenge_type": "Twitter Headline",
            "target_dim": "rule_of_three",
            "scaffold_signal": "Twitter condensation: does the user produce ONE short punchy sentence vs. expanding into details?",
        },
        {
            "id": "anchor_cp_rule3",
            "anchor_type": "Audience Benefit",
            "question_by_audience": {
                "classmate": (
                    "How does this actually affect us as students? "
                    "Can you break it down into just THREE simple points so it's easy to follow? "
                    "(引导：Rule of Three。不要长篇大论，强制用【第一、第二、第三】的结构来数数回答)"
                ),
                "professor": (
                    "From the audience's perspective, how does this topic connect to broader relevance? "
                    "Could you summarize the core benefits in exactly THREE distinct points? "
                    "(引导：Rule of Three。不要长篇大论，强制用【第一、第二、第三】的结构来数数回答)"
                ),
            },
            "question_type": "anchor",
            "category": "Class Presentation — Anchor",
            "challenge_type": "Audience Benefit",
            "target_dim": "rule_of_three",
            "scaffold_signal": "Rule of Three: does the user explicitly enumerate first / second / third?",
        },
        {
            "id": "anchor_cp_5w",
            "anchor_type": "5-W Detail",
            "question_by_audience": {
                "classmate": (
                    "Wait, can you give a real example? "
                    "Like, who was actually involved and when did this happen? "
                    "(引导：Illustrative Support。不要讲空话，讲讲具体的时间、人物和发生的事情)"
                ),
                "professor": (
                    "Regarding the case study you mentioned, rather than abstract theory, "
                    "can you specify the concrete details — specifically WHO was involved and WHEN this occurred? "
                    "(引导：Illustrative Support。不要讲空话，讲讲具体的时间、人物和发生的事情)"
                ),
            },
            "question_type": "anchor",
            "category": "Class Presentation — Anchor",
            "challenge_type": "5-W Detail",
            "target_dim": "illustrative_support",
            "scaffold_signal": "5-W specificity: does the user name a specific WHO and WHEN with concrete detail?",
        },
    ],
    "thesis_defense": [
        {
            "id": "anchor_td_convsense",
            "anchor_type": "Conversational Readiness",
            "question": (
                "Wait, that theoretical model is way too abstract. "
                "Can you explain the core logic to me without using any academic jargon, "
                "perhaps using a real-world analogy to show your conversational readiness? "
                "(引导：Conversational Sense。别背书，用喝咖啡聊天时的大白话或者打比方来解释)"
            ),
            "question_type": "anchor",
            "category": "Thesis Defense — Anchor",
            "challenge_type": "Conversational Readiness",
            "target_dim": "tact",
            "scaffold_signal": "Plain-language analogy: does the user drop academic jargon and use an everyday comparison or metaphor?",
        },
        {
            "id": "anchor_td_directness",
            "anchor_type": "Directness",
            "question": (
                "You stated a key causal relationship in your research. "
                "Please give me a direct, straight-to-the-point response: "
                "is this causal relationship solid, or is it just a correlation? "
                "Give your clear stance first. "
                "(引导：Directness 结论先行。不要绕弯子，第一句立刻说【是】或【不是】，然后再解释为什么)"
            ),
            "question_type": "anchor",
            "category": "Thesis Defense — Anchor",
            "challenge_type": "Directness",
            "target_dim": "directness",
            "scaffold_signal": "Directness: does sentence 1 begin with Yes/No/My finding is/I argue/The causal link is?",
        },
        {
            "id": "anchor_td_tact",
            "anchor_type": "Tact",
            "question": (
                "Your sample size seems heavily limited. "
                "Before you aggressively defend it, can you find a common ground with my concern "
                "and then explain your validation steps? "
                "(引导：Tact 学术外交手腕。不要一被质疑就急着反驳，先高情商地承认对方提得对，然后再解释)"
            ),
            "question_type": "anchor",
            "category": "Thesis Defense — Anchor",
            "challenge_type": "Tact",
            "target_dim": "tact",
            "scaffold_signal": "Tact: does the user open with a Carnegie acknowledgment phrase before defending?",
        },
    ],
    "case_pitch": [
        {
            "id": "anchor_cp_elevator",
            "anchor_type": "Elevator Test",
            "question": (
                "Sorry to interrupt, but imagine we only have 30 seconds left in the elevator. "
                "Give me your conclusion first: what are the top three key drivers of your business model? "
                "(引导：Conclusion First + Rule of Three。先说结论，然后立刻用三个数字或者业务驱动力来撑住它)"
            ),
            "question_type": "anchor",
            "category": "MBA Case Pitch — Anchor",
            "challenge_type": "Elevator Test",
            "target_dim": "conclusion_first",
            "scaffold_signal": "Conclusion first + Rule of Three: does user state a verdict in sentence 1, then enumerate 3 drivers?",
        },
        {
            "id": "anchor_cp_pathos",
            "anchor_type": "Pathos Story",
            "question": (
                "We have seen dozens of similar startup teams this week. "
                "Why exactly should we invest in YOU? "
                "Inject some passion and tell us a 10-second real user story that proves you are irreplaceable. "
                "(引导：Persuasion Mix / Pathos。别干巴巴地念财务数据，讲一个10秒钟的真实用户故事来打动我)"
            ),
            "question_type": "anchor",
            "category": "MBA Case Pitch — Anchor",
            "challenge_type": "Pathos Story",
            "target_dim": "persuasion_mix",
            "scaffold_signal": "Pathos story: does the user tell a specific user/customer story with emotional stakes?",
        },
    ],
}

# ─────────────────────────────────────────────
# THESIS DEFENSE: SCENE-SPECIFIC CHALLENGE POOL
# Persona locked: Professor (rigorous, high-pressure, academic).
# Used as fallback for generate_followup_question when AI is off,
# and as a reference bank for the CQ evaluation engine.
# ─────────────────────────────────────────────
THESIS_DEFENSE_CHALLENGE_POOL = [
    {
        "id": 201,
        "question": (
            "How exactly do you justify your chosen methodology over alternative frameworks, "
            "and what is your specific contribution to knowledge?"
        ),
        "category": "Thesis Defense", "difficulty": "Hard",
        "challenge_type": "Methodology Weakness", "questioner": "Professor",
    },
    {
        "id": 202,
        "question": (
            "Wait, that theoretical model is way too abstract. "
            "Can you explain the core logic to me without using any academic jargon, "
            "perhaps using a real-world analogy to show your conversational readiness? "
            "(引导：Conversational Sense。别背书，用喝咖啡聊天时的大白话或者打比方来解释)"
        ),
        "category": "Thesis Defense", "difficulty": "Hard",
        "challenge_type": "Clarity", "questioner": "Professor",
    },
    {
        "id": 203,
        "question": (
            "In hindsight, if you were given a $500,000 grant to do this research all over again, "
            "what different decisions would you make?"
        ),
        "category": "Thesis Defense", "difficulty": "Medium",
        "challenge_type": "Research Design", "questioner": "Professor",
    },
    {
        "id": 204,
        "question": (
            "On page 125, you stated that X leads to Y. "
            "Please give me a direct, straight-to-the-point response: "
            "is this causal relationship solid, or is it just a correlation? Give your clear stance first. "
            "(引导：Directness 结论先行。不要绕弯子，第一句立刻说【是】或【不是】，然后再解释为什么)"
        ),
        "category": "Thesis Defense", "difficulty": "Hard",
        "challenge_type": "Causality Issue", "questioner": "Professor",
    },
    {
        "id": 205,
        "question": (
            "Your sample size seems heavily limited. "
            "Before you aggressively defend it, can you find a common ground with my concern "
            "and then explain your validation steps? "
            "(引导：Tact 学术外交手腕。不要一被质疑就急着反驳，先高情商地承认对方提得对，然后再解释)"
        ),
        "category": "Thesis Defense", "difficulty": "Hard",
        "challenge_type": "Methodology Weakness", "questioner": "Professor",
    },
]

# ─────────────────────────────────────────────
# MBA CASE PITCH: SCENE-SPECIFIC CHALLENGE POOL
# Persona locked: VC / Investor (cold, results-oriented, capital-efficient).
# Used as fallback for generate_followup_question when AI is off,
# and as a reference bank for the CQ evaluation engine.
# ─────────────────────────────────────────────
CASE_PITCH_CHALLENGE_POOL = [
    {
        "id": 301,
        "question": (
            "Sorry to interrupt, but imagine we only have 30 seconds left in the elevator. "
            "Give me your conclusion first: what are the top three key drivers of your business model? "
            "(引导：Conclusion First + Rule of Three。先说结论，然后立刻用三个数字或者业务驱动力来撑住它)"
        ),
        "category": "Case Pitch", "difficulty": "Hard",
        "challenge_type": "Clarity", "questioner": "VC",
    },
    {
        "id": 302,
        "question": (
            "According to the 80/20 rule, you are spending too much time on minor details. "
            "What is the minimum necessary facts supporting your core argument?"
        ),
        "category": "Case Pitch", "difficulty": "Hard",
        "challenge_type": "Structure", "questioner": "VC",
    },
    {
        "id": 303,
        "question": (
            "I don't care about your idealistic revenue projections. "
            "Explain the exact parameter settings of your discount rate right now."
        ),
        "category": "Case Pitch", "difficulty": "Hard",
        "challenge_type": "Financial Feasibility", "questioner": "VC",
    },
    {
        "id": 304,
        "question": (
            "We have seen dozens of similar startup teams this week. "
            "Why exactly should we invest in YOU? "
            "Inject some passion and tell us a 10-second real user story that proves you are irreplaceable. "
            "(引导：The Persuasion Mix / Pathos。别干巴巴地念财务数据，讲一个10秒钟的真实用户故事来打动我)"
        ),
        "category": "Case Pitch", "difficulty": "Hard",
        "challenge_type": "Team Competence", "questioner": "VC",
    },
]


# ─────────────────────────────────────────────
# PDF → BASE64 IMAGES PIPELINE (Step 1)
# ─────────────────────────────────────────────
def extract_pdf_images_as_base64(filepath, max_pages=20):
    """Convert PDF pages to base64-encoded PNG images for Vision API."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        images = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            mat = fitz.Matrix(1.5, 1.5)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
            images.append({"page": i + 1, "base64": b64, "media_type": "image/png"})
        doc.close()
        return images
    except Exception as e:
        app.logger.warning(f"PDF extraction failed: {e}")
        return []


def extract_ppt_images_as_base64(filepath, max_pages=20):
    """For PPT/PPTX files — extract text per slide."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides):
            if i >= max_pages:
                break
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides.append({"page": i + 1, "text": "\n".join(texts)})
        return slides
    except Exception:
        return []


def extract_pdf_text_slides(filepath, max_pages=20):
    """
    Fast text extraction from PDF pages — no AI, no base64 encoding.
    Used to populate session["slides"] immediately after upload so the config
    page shows the REAL page count and titles instead of the 3-page mock.
    Falls back to MOCK_SLIDES only if the file cannot be opened at all.
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        slides = []
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            text = page.get_text("text").strip()
            lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
            # First non-empty line becomes the title; next few lines become content
            title   = lines[0][:80] if lines else f"Slide {i + 1}"
            content = " ".join(lines[1:6]) if len(lines) > 1 else "(No text detected on this page)"
            slides.append({
                "page":       i + 1,
                "title":      title,
                "content":    content,
                "key_claims": [],
            })
        doc.close()
        return slides if slides else MOCK_SLIDES
    except Exception as e:
        app.logger.warning(f"extract_pdf_text_slides failed: {e}")
        return MOCK_SLIDES


def pptx_to_slides(pptx_data):
    """Convert extract_ppt_images_as_base64 output to the standard slides format."""
    result = []
    for s in pptx_data:
        lines   = [ln.strip() for ln in s["text"].split("\n") if ln.strip()]
        title   = lines[0][:80] if lines else f"Slide {s['page']}"
        content = " ".join(lines[1:6]) if len(lines) > 1 else ""
        result.append({
            "page":       s["page"],
            "title":      title,
            "content":    content,
            "key_claims": [],
        })
    return result if result else MOCK_SLIDES


# ─────────────────────────────────────────────
# VISION: ANALYZE SLIDES (OpenAI-compatible multimodal)
# ─────────────────────────────────────────────
def analyze_slides_with_claude(images_b64, filename):
    """
    Send slide images to the AI Vision endpoint (OpenAI-compatible).
    Falls back to mock data on any error.
    """
    if not AI_ENABLED or not images_b64:
        return MOCK_SLIDES, False

    try:
        # Build OpenAI-compatible multimodal content list
        content_parts = [
            {
                "type": "text",
                "text": (
                    "You are an expert presentation analyst. I am uploading slide images from a presentation. "
                    "For each slide, extract: the slide title, a concise summary of the main content (2-4 sentences), "
                    "and the key claims or data points made.\n\n"
                    "Return a JSON array ONLY (no other text) with this exact structure:\n"
                    '[{"page": 1, "title": "...", "content": "...", "key_claims": ["claim1", "claim2"]}]\n\n'
                    "Be concise and factual. Extract exactly what is on the slides."
                ),
            }
        ]
        for img in images_b64:
            content_parts.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:{img['media_type']};base64,{img['base64']}"
                },
            })

        response = _ai_client.chat.completions.create(
            model=VISION_MODEL,
            max_tokens=MAX_TOKENS,
            messages=[{"role": "user", "content": content_parts}],
        )

        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        slides = json.loads(raw)
        for i, s in enumerate(slides):
            s.setdefault("page", i + 1)
            s.setdefault("key_claims", [])
        return slides, True

    except Exception as e:
        app.logger.error(f"AI slide analysis failed: {e}")
        return MOCK_SLIDES, False


# ─────────────────────────────────────────────
# CLAUDE: MASTER ENGINE — BUILD CHALLENGE SEED + QA BANK (Step 3)
# ─────────────────────────────────────────────
def build_master_engine(slides, audience, scenario, difficulty):
    """
    The Master Engine: Claude reads slides + config and returns:
    - challenge_seed (for interruptable scenarios)
    - static_qa_bank (for Class Presentation)
    Uses the full challenge classification dictionary in the system prompt.
    Falls back to mock data on error.
    """
    if not AI_ENABLED:
        return MOCK_CHALLENGE, MOCK_QA_BANK

    # Select challenge type dictionary based on scenario
    if scenario == "MBA Case Pitch":
        challenge_dict = BUSINESS_CHALLENGE_TYPES
        challenge_dict_name = "Business"
    else:
        challenge_dict = ACADEMIC_CHALLENGE_TYPES
        challenge_dict_name = "Academic"

    # Build scenario-specific instructions
    if scenario in ("Class Presentation", "Academic Presentation"):
        task_instruction = (
            "This is a CLASS PRESENTATION (no mid-session interruptions). "
            "The presenter is a university or international student giving a general topic presentation "
            "to classmates or a professor — NOT a research thesis defense. "
            "Generate exactly 3 questions a real classmate or professor might ask in a seminar setting: "
            "focus on clarity of explanation, quality of examples, and relevance to the audience. "
            "Each question must name which challenge type it targets. "
            "Return them as 'static_qa_bank' in the JSON.\n"
            "Set 'challenge_seed' to null."
        )
    else:
        trigger_page = min(2, len(slides))
        task_instruction = (
            f"This scenario has MID-SESSION INTERRUPTIONS. "
            f"Identify the single most devastating vulnerability in slide {trigger_page}. "
            f"Map it to exactly one challenge type from the dictionary above. "
            f"Generate the sharpest possible opening challenge question in the voice of the {audience} persona. "
            f"Set 'trigger_page' to {trigger_page}.\n"
            "Set 'static_qa_bank' to an empty array []."
        )

    slide_text = "\n\n".join(
        f"--- SLIDE {s['page']}: {s.get('title','')}\n{s.get('content','')}"
        for s in slides
    )

    challenge_dict_text = "\n".join(
        f"  - {k}: {v}" for k, v in challenge_dict.items()
    )

    audience_persona_text = AUDIENCE_PERSONA.get(audience, AUDIENCE_PERSONA["Professor"])
    difficulty_text = DIFFICULTY_INSTRUCTIONS.get(difficulty, DIFFICULTY_INSTRUCTIONS["Medium"])
    _lang_directive = audience_language_directive(audience)
    language_block = f"\nQUESTION LANGUAGE LEVEL:\n{_lang_directive}\n" if _lang_directive else ""

    prompt = f"""You are a world-class presentation examiner AI.

AUDIENCE PERSONA:
{audience_persona_text}

DIFFICULTY MODE:
{difficulty_text}
{language_block}
{challenge_dict_name.upper()} CHALLENGE TYPE DICTIONARY (you MUST pick from these):
{challenge_dict_text}

PRESENTATION SLIDES:
{slide_text}

TASK:
{task_instruction}

Return ONLY valid JSON in this exact structure (no other text, no markdown):
{{
  "challenge_seed": {{
    "trigger_page": <int>,
    "challenge_type": "<one of the challenge type keys above>",
    "initial_challenge": "<the exact question to ask the presenter, in the persona's voice>"
  }},
  "static_qa_bank": [
    {{
      "id": 1,
      "question": "<deep question>",
      "challenge_type": "<challenge type key>",
      "category": "<Academic/Business/Technical>",
      "difficulty": "<Easy/Medium/Hard>"
    }}
  ]
}}"""

    try:
        response = _ai_client.chat.completions.create(
            model=TEXT_MODEL,
            max_tokens=MAX_TOKENS,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        data = json.loads(raw)

        challenge_seed = data.get("challenge_seed") or MOCK_CHALLENGE
        static_qa_bank = data.get("static_qa_bank") or []

        for i, q in enumerate(static_qa_bank):
            q["id"] = i + 1

        return challenge_seed, static_qa_bank

    except Exception as e:
        app.logger.error(f"Master engine failed: {e}")
        return MOCK_CHALLENGE, MOCK_QA_BANK


# ─────────────────────────────────────────────
# THESIS DEFENSE: DYNAMIC QUESTION FUSION ENGINE (Step 7)
# Rewrites DEFENSE_QUESTION_BANK templates by fusing them with the user's
# actual thesis topic / slide content, so the AI examiner never reads the
# raw template verbatim. answering_strategy stays anchored to the template
# question's underlying logic — only the wording of the question changes.
# ─────────────────────────────────────────────
def _build_thesis_context(slides):
    """Compress parsed slides into a compact topic + keyword block for prompt injection."""
    if not slides:
        return "No slide content available — the user's PPT/PDF could not be parsed."
    titles = [s.get("title", "") for s in slides if s.get("title")]
    topic = titles[0] if titles else "Untitled Thesis"
    lines = [f"Thesis Topic (inferred from slide 1): {topic}"]
    for s in slides[:8]:
        claims = ", ".join(s.get("key_claims", []) or [])
        line = f"Slide {s.get('page')}: {s.get('title', '')} — {s.get('content', '')[:180]}"
        if claims:
            line += f" | Key claims: {claims}"
        lines.append(line)
    return "\n".join(lines)


def generate_custom_defense_question(base_question_obj, thesis_context):
    """
    Rewrite a single DEFENSE_QUESTION_BANK template question by deeply embedding
    the user's actual thesis topic, keywords, or slide content. Falls back to the
    original template question text on any AI failure or when AI is disabled.
    """
    base_text = base_question_obj.get("question", "")
    if not AI_ENABLED:
        return base_text

    system_prompt = (
        "You are an expert professor conducting a formal thesis defense. Your task is to "
        "take a question from the template database and rewrite it by deeply embedding "
        "the user's actual thesis topic, keywords, or presentation slides.\n\n"
        "[Rules for Generating the Question]:\n"
        "- DO NOT just read the template question verbatim. You MUST mention specific terms, "
        "methods, or markets from the user's thesis context.\n"
        f"- {audience_language_directive('Professor')}\n"
        "- Ensure the user can easily recognize that you have genuinely reviewed their specific slides.\n"
        "- Preserve the template question's underlying intent/challenge type — change only how it "
        "is phrased, never what is fundamentally being asked.\n"
        "- Return ONLY the rewritten question text — no labels, no quotes, no explanations."
    )
    user_prompt = (
        f"Base Question Template: {base_text}\n"
        f"Answering Strategy for reference only (do not repeat verbatim): "
        f"{base_question_obj.get('answering_strategy', '')}\n\n"
        f"User's Thesis Topic/PPT Content:\n{thesis_context}\n\n"
        "Task: Rewrite the Base Question Template into a customized question. Inject the "
        "user's specific topic details naturally. Keep the professional academic tone intact."
    )
    try:
        resp = _ai_client.chat.completions.create(
            model=TEXT_MODEL,
            max_tokens=200,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        )
        text = (resp.choices[0].message.content or "").strip().strip('"')
        return text if text else base_text
    except Exception as e:
        app.logger.error(f"[DEFENSE QA] Custom question generation failed for {base_question_obj.get('id')}: {e}")
        return base_text


def _build_presentation_context(slides):
    """
    Compress parsed slides into a compact topic + keyword block for prompt injection
    in generate_custom_anchor_question(). Parallel to _build_thesis_context().
    """
    if not slides:
        return "No slide content available — the user's PPT/PDF could not be parsed."
    titles = [s.get("title", "") for s in slides if s.get("title")]
    topic = titles[0] if titles else "Untitled Presentation"
    lines = [f"Presentation Topic (inferred from slide 1): {topic}"]
    for s in slides[:6]:
        line = f"Slide {s.get('page')}: {s.get('title', '')} — {s.get('content', '')[:180]}"
        lines.append(line)
    return "\n".join(lines)


def generate_custom_anchor_question(anchor_template, slides, audience):
    """
    Rewrite a class_presentation anchor question so that the question BODY
    references the user's actual slide content/topic, while preserving the
    (引导：...) scaffold hint verbatim — that hint is what the CQ engine scores.

    Falls back to the static template text on AI failure or when AI is off.
    """
    import re as _re
    track = "professor" if audience.lower() == "professor" else "classmate"
    by_audience = anchor_template.get("question_by_audience", {})
    base_text = (
        by_audience.get(track)
        or by_audience.get("professor")
        or anchor_template.get("question", "")
    )

    if not AI_ENABLED or not slides:
        return base_text

    # Peel off the 引导 hint so we can re-attach it unchanged after rewriting.
    hint_match = _re.search(r'(\(引导[：:][^)]+\))', base_text)
    scaffold_hint = hint_match.group(1) if hint_match else ""
    body_only     = base_text[:hint_match.start()].strip() if hint_match else base_text

    anchor_type          = anchor_template.get("anchor_type", "")
    presentation_context = _build_presentation_context(slides)
    track_persona        = "a professor" if track == "professor" else "a classmate"

    prompt = (
        f"You are {track_persona} in a university classroom. A student just finished presenting.\n\n"
        f"Presentation content:\n{presentation_context}\n\n"
        f"Question challenge type: {anchor_type}\n"
        f"Original question template: {body_only}\n\n"
        "Task: Rewrite the original question template so it references specific content from the "
        "presentation above — mention the actual topic, a specific finding, slide detail, or key claim. "
        "Keep the same challenge intent and tone as the original template. "
        "Return ONLY the rewritten question body — no labels, no quotes, no explanations, "
        "and do NOT include any 引导 hint text (that will be appended separately)."
    )
    try:
        resp = _ai_client.chat.completions.create(
            model=TEXT_MODEL,
            max_tokens=160,
            messages=[{"role": "user", "content": prompt}],
        )
        new_body = (resp.choices[0].message.content or "").strip().strip('"')
        if not new_body:
            return base_text
        return f"{new_body} {scaffold_hint}".strip() if scaffold_hint else new_body
    except Exception as e:
        app.logger.error(
            f"[ANCHOR QA] Custom anchor generation failed for {anchor_template.get('id')}: {e}"
        )
        return base_text


def customize_defense_qa_bank(pool, slides):
    """
    Apply generate_custom_defense_question() to every sampled template question,
    in parallel (bounded thread pool) to keep /config latency reasonable even at
    Hard difficulty (8 questions). id / challenge_type / category / answering_strategy
    are copied through unchanged — only 'question' is replaced with the fused text.
    The original template text is preserved as 'template_question' for debugging/logging.
    """
    import concurrent.futures as _cf

    results = [None] * len(pool)

    if not AI_ENABLED or not pool:
        for i, q in enumerate(pool):
            q2 = dict(q)
            q2["template_question"] = q.get("question", "")
            results[i] = q2
        return results

    thesis_context = _build_thesis_context(slides)

    def _work(i, q):
        q2 = dict(q)
        q2["template_question"] = q.get("question", "")
        q2["question"] = generate_custom_defense_question(q, thesis_context)
        return i, q2

    with _cf.ThreadPoolExecutor(max_workers=min(6, len(pool))) as ex:
        futures = [ex.submit(_work, i, q) for i, q in enumerate(pool)]
        for fut in _cf.as_completed(futures):
            i, q2 = fut.result()
            results[i] = q2

    return results


# ─────────────────────────────────────────────
# CLAUDE: DYNAMIC FOLLOW-UP GENERATOR (Step 6)
# ─────────────────────────────────────────────
def generate_followup_question(
    slides, audience, difficulty, challenge_type, chat_history, user_answer, current_page
):
    """
    Generate a context-aware follow-up question given the user's answer
    and the full conversation history. Falls back to mock pool on error.
    """
    if not AI_ENABLED:
        used = [h["content"] for h in chat_history if h["role"] == "assistant"]
        # Scene-specific fallback pools for persona fidelity
        if audience == "VC":
            pool = [q["question"] for q in CASE_PITCH_CHALLENGE_POOL]
        elif audience == "Professor":
            pool = [q["question"] for q in THESIS_DEFENSE_CHALLENGE_POOL]
        else:
            pool = list(MOCK_FOLLOWUP_POOL)
        available = [q for q in pool if q not in used]
        return random.choice(available) if available else pool[0]

    current_slide = next((s for s in slides if s["page"] == current_page), slides[0])
    persona = AUDIENCE_PERSONA.get(audience, AUDIENCE_PERSONA["Professor"])
    difficulty_inst = DIFFICULTY_INSTRUCTIONS.get(difficulty, DIFFICULTY_INSTRUCTIONS["Medium"])
    _lang_directive = audience_language_directive(audience) or (
        "Phrase the question in clear, professional English appropriate for the persona."
    )

    history_text = "\n".join(
        f"{'EXAMINER' if h['role'] == 'assistant' else 'PRESENTER'}: {h['content']}"
        for h in chat_history
    )

    # Scene-aware scaffolding instructions: every AI follow-up must embed an explicit coaching hint
    _scaffolding_hint = {
        "VC": (
            "SCAFFOLDING RULE: Your question MUST end with a parenthetical coaching hint in Chinese "
            "that tells the presenter exactly how to structure their answer. "
            "Examples: (引导：结论先行，再用三个数字撑住它) / "
            "(引导：别念数据，用一个10秒真实用户故事打动我) / "
            "(引导：先说是还是不是，然后再给理由)"
        ),
        "Professor": (
            "SCAFFOLDING RULE: Your question MUST end with a parenthetical coaching hint in Chinese "
            "that tells the presenter exactly how to answer. "
            "Examples: (引导：直接说结论，第一句先给立场) / "
            "(引导：别用术语，用大白话或打比方来解释) / "
            "(引导：先高情商地承认对方有道理，再解释你的做法)"
        ),
        "Classmates": (
            "SCAFFOLDING RULE: Your question MUST end with a parenthetical coaching hint in Chinese "
            "that tells the presenter how to frame their answer. "
            "Examples: (引导：用第一、第二、第三的结构来数数回答) / "
            "(引导：讲一个具体的时间和人物，不要讲空话) / "
            "(引导：一句话总结，像发朋友圈那样短小精悍)"
        ),
    }.get(audience, (
        "SCAFFOLDING RULE: Your question MUST end with a parenthetical coaching hint in Chinese "
        "that tells the presenter how to structure their answer."
    ))

    prompt = f"""You are a {audience} examiner running a high-stakes presentation challenge.

YOUR PERSONA: {persona}

DIFFICULTY: {difficulty_inst}

SLIDE BEING DISCUSSED:
Title: {current_slide.get('title', '')}
Content: {current_slide.get('content', '')}

CHALLENGE TYPE BEING PRESSED: {challenge_type}

CONVERSATION SO FAR:
{history_text}
PRESENTER: {user_answer}

TASK: Generate your NEXT follow-up question. You must:
1. Directly reference something specific from the presenter's answer above
2. Press harder on the same "{challenge_type}" vulnerability — do not change topic
3. Stay strictly in character as the {audience} persona
4. Be {difficulty} in intensity
5. {_scaffolding_hint}
6. {_lang_directive}

Return ONLY the question text with the coaching hint appended. No preamble, no labels, no markdown."""

    try:
        response = _ai_client.chat.completions.create(
            model=TEXT_MODEL,
            max_tokens=512,
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        app.logger.error(f"Follow-up generation failed: {e}")
        used = [h["content"] for h in chat_history if h["role"] == "assistant"]
        available = [q for q in MOCK_FOLLOWUP_POOL if q not in used]
        return random.choice(available) if available else MOCK_FOLLOWUP_POOL[0]


# ─────────────────────────────────────────────
# CLASS PRESENTATION Q&A BUILDER
# ─────────────────────────────────────────────
def build_class_presentation_qa(audience, difficulty):
    """
    Build an audience-pinned Q&A bank for the Class Presentation scenario.

    - 'Professor' audience  → pulls professor-track questions from TED_QA_MATRIX
    - 'Classmates' audience → pulls classmate-track questions from TED_QA_MATRIX
    - No random.choice on questioner role — strictly follows the user-selected audience.
    - Baseline CLASS_PRES_QA_POOL items are appended as fallback depth.
    """
    track = "professor" if audience.lower() == "professor" else "classmate"
    qa_count = {"Easy": 1, "Medium": 2, "Hard": 3}.get(difficulty, 2)

    ted_questions = []
    for idx, (dim_key, dim_data) in enumerate(TED_QA_MATRIX.items()):
        ted_questions.append({
            "id":            100 + idx,
            "question":      dim_data[track],
            "category":      "Class Presentation — TED Framework",
            "difficulty":    difficulty,
            "challenge_type": dim_key.replace("_", " ").title(),
            "ted_dimension": dim_key,
            "questioner":    audience,
        })

    # Annotate baseline pool with questioner identity
    baseline = []
    for q in CLASS_PRES_QA_POOL:
        item = dict(q)
        item["questioner"] = audience
        baseline.append(item)

    # TED questions first; baseline as extra depth; trim to qa_count
    combined = (ted_questions + baseline)[:qa_count]
    return combined


# ─────────────────────────────────────────────
# DUAL-TRACK Q&A BUILDER
# ─────────────────────────────────────────────
def generate_free_qa_question(slides, audience, scene_slug):
    """
    Generate a single context-aware free question using AI + slide content.
    No scaffold hint — purely contextual. Returns a question dict with question_type='free'.
    """
    SCENE_PERSONA = {
        "class_presentation": f"a curious {audience.lower()} in a university classroom",
        "thesis_defense":      "a thesis committee examiner probing the research",
        "case_pitch":          "a VC investor challenging the business pitch",
    }
    persona    = SCENE_PERSONA.get(scene_slug, f"a {audience.lower()}")
    slide_text = "\n".join(
        f"Slide {s['page']}: {s.get('title','')} — {s.get('content','')[:200]}"
        for s in (slides or [])[:5]
    )
    _fallbacks = {
        "class_presentation": "What do you think is the most surprising insight from your presentation, and why does it matter to your audience?",
        "thesis_defense":     "How does your methodology address the most significant limitation you identified in your literature review?",
        "case_pitch":         "What is the single biggest risk to your business model right now, and what is your specific mitigation plan?",
    }
    fallback_q = _fallbacks.get(scene_slug, "Can you walk us through the most important takeaway from your presentation?")

    if not AI_ENABLED or not slides:
        return {
            "id": "free_q_fallback", "question": fallback_q,
            "question_type": "free", "category": f"{scene_slug.replace('_',' ').title()} — AI Free Question",
            "difficulty": "Medium", "challenge_type": "Contextual",
        }
    # thesis_defense is always the Professor persona; case_pitch (VC) has no mandated IELTS level yet.
    _effective_audience = "Professor" if scene_slug == "thesis_defense" else audience
    _lang_directive = audience_language_directive(_effective_audience)
    _lang_line = f"- {_lang_directive}\n" if _lang_directive else ""
    prompt = (
        f"You are {persona}.\n\n"
        f"Presentation content:\n{slide_text}\n\n"
        "Generate ONE sharp, context-specific question that:\n"
        "- Directly references something specific from the slides above\n"
        "- Is appropriate for the setting\n"
        "- Has NO scaffolding hints, guidance notes, or Chinese text\n"
        f"{_lang_line}"
        "- Is 1-2 sentences maximum\n\n"
        "Return ONLY the question text — no labels, no explanations."
    )
    try:
        resp = _ai_client.chat.completions.create(
            model=VISION_MODEL, max_tokens=120,
            messages=[{"role": "user", "content": prompt}],
        )
        q_text = resp.choices[0].message.content.strip().strip('"')
        return {
            "id": "free_q_ai", "question": q_text,
            "question_type": "free", "category": f"{scene_slug.replace('_',' ').title()} — AI Free Question",
            "difficulty": "Medium", "challenge_type": "Contextual",
        }
    except Exception as e:
        app.logger.error(f"[FREE QA] AI generation failed: {e}")
        return {
            "id": "free_q_fallback", "question": fallback_q,
            "question_type": "free", "category": f"{scene_slug.replace('_',' ').title()} — AI Free Question",
            "difficulty": "Medium", "challenge_type": "Contextual",
        }


def build_dual_track_qa(slides, audience, scene_slug, difficulty):
    """
    Build the dual-track Q&A bank for post-session Q&A.

    Track structure by difficulty / question count:
      Easy  (1 q): Q1 — AI free question only
      Medium (2 q): Q1 — AI free  +  Q2 — anchor (scaffolded, randomly selected)
      Hard  (3 q): Q1 — AI free  +  Q2 — anchor  +  Q3 — baseline fallback free

    Exactly ONE anchor per session ("有且仅有一个标准锚点题").
    """
    qa_count = {"Easy": 1, "Medium": 2, "Hard": 3}.get(difficulty, 2)

    q1 = generate_free_qa_question(slides, audience, scene_slug)
    result = [q1]

    if qa_count >= 2:
        anchor_pool = ANCHOR_QUESTION_POOL.get(scene_slug, [])
        if anchor_pool:
            anchor_template = random.choice(anchor_pool)
            q2 = dict(anchor_template)
            by_audience = q2.pop("question_by_audience", None)
            if by_audience:
                track = "professor" if audience.lower() == "professor" else "classmate"
                q2["question"] = by_audience.get(track, by_audience.get("professor"))
            q2["questioner"] = audience
            # Class Presentation: dynamically rewrite the anchor question body to
            # reference the user's actual slide content, while preserving the
            # (引导：...) scaffold hint verbatim for CQ scoring.
            if scene_slug == "class_presentation" and slides:
                q2["template_question"] = q2["question"]   # keep original for debugging
                q2["question"] = generate_custom_anchor_question(anchor_template, slides, audience)
            result.append(q2)
        else:
            track = "professor" if audience.lower() == "professor" else "classmate"
            fallback_dims = list(TED_QA_MATRIX.values())
            if fallback_dims:
                dim_data = fallback_dims[0]
                result.append({
                    "id": "anchor_fallback", "question": dim_data[track],
                    "question_type": "anchor", "category": "Anchor Fallback",
                    "challenge_type": "Structural", "target_dim": "rule_of_three",
                    "scaffold_signal": "structured enumeration",
                    "anchor_type": "Structural", "questioner": audience,
                })

    if qa_count >= 3:
        baseline_pool = [q for q in CLASS_PRES_QA_POOL if q.get("difficulty") != "Easy"]
        if baseline_pool:
            q3 = dict(random.choice(baseline_pool))
            q3["question_type"] = "free"
            q3["questioner"]    = audience
            result.append(q3)

    return result


# ═══════════════════════════════════════════════════════════════════════════════
# COMMUNICATION QUALITY (CQ) ENGINE
# Scene-targeted evaluation of Q&A and interrupt exchanges only.
# Three scene branches: thesis_defense / case_pitch / class_presentation
# Each branch has 3 weighted sub-dimensions totalling 100 pts.
# ═══════════════════════════════════════════════════════════════════════════════

def _cq_heuristic_universal(qa_texts):
    """
    Universal heuristic for FREE questions (no scene anchor).
    Returns (directness_logic, conversational_resonance, evidence_substantiation) 0-100.
    Dims match Module 1: Directness & Logic (40%) / Conv Resonance (30%) / Evidence (30%).
    """
    all_text = " ".join(qa_texts).lower()

    # Dim 1: Directness & Logic
    DIRECT_OPEN = (
        r'^(?:yes[,\s]|no[,\s]|the (?:key|main|core|answer|point|issue)|'
        r'i (?:believe|argue|think|found|see)|we (?:found|show|believe)|'
        r'in short[,\s]|to answer[,\s]|my (?:view|point|answer|stance))'
    )
    DODGE_OPEN = (
        r'^(?:as (?:i|we) mentioned|to give (?:some )?context|'
        r'that\'?s (?:a )?(?:great|good|interesting)|thank you for|'
        r'so[,\s](?:as|before|first[,\s])|first[,\s](?:let me|i want to))'
    )
    directness = 55
    for txt in qa_texts:
        first = txt.strip()[:200].lower()
        if re.search(DIRECT_OPEN, first, re.MULTILINE):
            directness = min(100, directness + 18)
        if re.search(DODGE_OPEN, first, re.MULTILINE):
            directness = max(0, directness - 12)
    directness = max(0, min(100, directness))

    # Dim 2: Conversational Resonance
    ORAL_RE = (
        r'\b(?:you know|think about (?:it|this)|the thing is|basically|'
        r'i mean|right\?|imagine|let me (?:explain|tell you)|'
        r'what i\'m saying is|it\'s like|here\'s the (?:thing|deal)|'
        r'actually|you see)\b'
    )
    SCRIPT_RE = r'\b(?:furthermore|nevertheless|in conclusion|to summarize|accordingly|thus|hence|moreover)\b'
    oral_hits   = len(re.findall(ORAL_RE,   all_text))
    script_hits = len(re.findall(SCRIPT_RE, all_text))
    if   oral_hits >= 4 and script_hits == 0: conv_resonance = 92
    elif oral_hits >= 3 and script_hits <= 1: conv_resonance = 80
    elif oral_hits >= 2:                       conv_resonance = 70
    elif oral_hits >= 1:                       conv_resonance = 57
    elif script_hits >= 3:                     conv_resonance = 35
    else:                                      conv_resonance = 50

    # Dim 3: Evidence & Substantiation
    EVIDENCE_RE = (
        r'\b(?:for example|for instance|such as|specifically|'
        r'research (?:shows?|found|indicates?)|studies? (?:show|found)|'
        r'data (?:shows?|suggests?)|according to|'
        r'\d+[\.,]?\d*\s*(?:%|percent|million|billion|users?|customers?)|'
        r'in (?:20\d\d|19\d\d)|'
        r'a (?:friend|colleague|client|customer|student|professor) (?:told|said|shared|mentioned|found))\b'
    )
    ev_hits = len(re.findall(EVIDENCE_RE, all_text))
    if   ev_hits >= 3: evidence = 90
    elif ev_hits >= 2: evidence = 75
    elif ev_hits >= 1: evidence = 58
    else:              evidence = 35

    return directness, conv_resonance, evidence


def _check_anchor_compliance(anchor_text, target_dim):
    """
    Heuristic: did the user follow the anchor question's scaffold hint?
    Returns (passed: bool, score: int 0-100).
    """
    if not anchor_text:
        return False, 20
    text_l = anchor_text.lower()
    _checks = {
        "rule_of_three":        (r'\b(first(?:ly)?[,\s]|second(?:ly)?[,\s]|third(?:ly)?[,\s]|there are (?:three|3)|number (?:one|two|three))\b', 2),
        "illustrative_support": (r'\b(who|when|where|for example|specifically|in (?:20\d\d|19\d\d)|a (?:person|student|teacher|professor|company|friend)|result was)\b', 2),
        "directness":           (r'^(?:yes[,\s]|no[,\s]|the answer is|i believe|i argue|my finding|it is|the key|the main|to answer)', 1),
        "tact":                 (r'\b(valid concern|good point|you(?:\'?re| are) right|i understand|fair point|that\'?s (?:a )?valid|i agree|i see your)\b', 1),
        "conclusion_first":     (r'^(?:yes[,\s]|no[,\s]|the (?:key|main|top|answer)|we believe|i (?:believe|argue)|the (?:three|first))', 1),
        "persuasion_mix":       (r'\b(user|customer|client|story|told us|experienced|felt|a (?:real|specific) (?:user|customer|person|example))\b', 1),
    }
    pattern, min_hits = _checks.get(target_dim, (r'\b\w+\b', 5))
    hits  = len(re.findall(pattern, text_l, re.MULTILINE | re.IGNORECASE))
    if hits >= min_hits:
        return True,  min(100, 65 + hits * 8)
    return False, max(20, 35 + hits * 10)


# Universal dim fallback fix phrases (used in mock dual-track result)
_UNIVERSAL_FIX_PHRASES = {
    "Directness & Logic":        "Lead with your stance immediately. Say this instead: 'To answer directly: [your core point]. The reason is [one key logic].' — never re-state the question first.",
    "Conversational Resonance":  "Sound like you're talking, not reading. Say this instead: 'Think about it this way — [your core point]. The key insight is [idea].' — add connectors like you know / the thing is.",
    "Evidence & Substantiation": "Add one concrete fact or story. Say this instead: 'For example, in [year], [who] at [place] did [what], and the result was [outcome].' — one specific case beats ten adjectives.",
}


def _cq_heuristic_thesis_defense(qa_texts):
    """Branch A: Thesis Defense — returns (directness, defensibility, tact) 0-100."""
    all_text = " ".join(qa_texts).lower()

    YES_RESP_RE = (
        r'\b(that is (?:a )?(?:valid|fair|good|important)|you(?:\'?re| are) right|'
        r'i (?:agree|see your point)|good point|fair point|that\'s a (?:valid|fair|good))\b'
    )
    DIRECT_OPEN_RE = (
        r'^(?:yes|no|the answer is|our approach|we (?:believe|found|show|demonstrated|used)|'
        r'i (?:believe|argue|found|conducted)|the (?:key|main|central|critical)|'
        r'in short|to answer|our (?:study|research|data|method)|the (?:result|finding))'
    )
    BG_PADDING_RE = (
        r'^(?:as (?:i mentioned|we discussed)|to (?:give|provide) (?:some )?context|'
        r'in (?:our|the) (?:study|research|paper)|the (?:study|research|paper) (?:investigated|examined))'
    )

    directness = 55
    for txt in qa_texts:
        first = txt.strip()[:180].lower()
        if re.search(DIRECT_OPEN_RE, first):
            directness = min(100, directness + 15)
        if re.search(BG_PADDING_RE, first):
            directness = max(0, directness - 15)

    DATA_RE = (
        r'\b(?:\d+(?:\.\d+)?%|\d[\d,]*\s*(?:participants|cases|subjects|samples|years)|'
        r'p\s*[<>=]\s*0\.\d+|statistically|significant(?:ly)?|evidence|empirical|'
        r'validated|finding|result|literature|citation|coefficient|regression|n\s*=\s*\d+)\b'
    )
    SUBJ_RE = r'\b(?:very|really|quite|somewhat|rather|pretty|fairly|maybe|perhaps|probably|i feel|in my opinion)\b'
    data_hits = len(re.findall(DATA_RE, all_text))
    subj_hits = len(re.findall(SUBJ_RE, all_text))

    if   data_hits >= 5: defensibility = 92
    elif data_hits >= 3: defensibility = 78
    elif data_hits >= 1: defensibility = 60
    else:                defensibility = 40
    if subj_hits > data_hits + 2:
        defensibility = max(0, defensibility - 15)

    tact_hits  = len(re.findall(YES_RESP_RE, all_text))
    DIPLO_RE = (
        r'\b(?:i (?:can see|understand) (?:your|why)|while (?:you|that|this) (?:raise|point|note)|'
        r'building on (?:that|your)|that is (?:indeed|certainly) a|you raise a (?:valid|good|fair))\b'
    )
    diplo_hits = len(re.findall(DIPLO_RE, all_text))
    if   tact_hits >= 2 or diplo_hits >= 2: tact = 92
    elif tact_hits >= 1 or diplo_hits >= 1: tact = 74
    else:                                   tact = 44

    return directness, defensibility, tact


def _cq_heuristic_case_pitch(qa_texts, total_qa_seconds=0):
    """Branch B: MBA Case Pitch — returns (conclusion_first, persuasion_mix, command_presence)."""
    all_text = " ".join(qa_texts).lower()

    ANS_FIRST_RE = (
        r'^(?:yes[,.!]|no[,.!]|we (?:are|can|do|will|have|plan)|'
        r'the (?:answer|key|solution|bottom line) is|our (?:approach|strategy|model|product)|'
        r'i (?:recommend|propose|believe)|absolutely|definitely|the market)'
    )
    LOGIC_RE = r'\b(?:because|the reason (?:is|being)|due to|driven by|which means|the rationale|specifically)\b'

    con_scores = []
    for txt in qa_texts:
        first = txt.strip()[:200].lower().lstrip()
        has_af = bool(re.search(ANS_FIRST_RE, first))
        has_l  = bool(re.search(LOGIC_RE, first))
        if   has_af and has_l: con_scores.append(95)
        elif has_af:           con_scores.append(75)
        elif has_l:            con_scores.append(60)
        else:                  con_scores.append(40)
    conclusion_first = int(sum(con_scores) / len(con_scores)) if con_scores else 50

    LOGOS_RE = (
        r'\b(?:\d+(?:\.\d+)?%|\$[\d,]+|[\d,]+\s*(?:users|customers|revenue|million|billion)|'
        r'metric|kpi|roi|growth|retention|churn|arpu|ltv|cac)\b'
    )
    ETHOS_RE = (
        r'\b(?:we(?:\'ve| have) (?:built|launched|deployed|partnered|validated)|our team|'
        r'years of experience|track record|we (?:work with|serve|partner)|our (?:portfolio|clients))\b'
    )
    PATHOS_RE = (
        r'\b(?:imagine|picture|story|one of our (?:users|customers)|a (?:user|customer|client) '
        r'(?:told|said|shared)|real (?:example|case)|transform|empower|struggling|pain|challenge)\b'
    )
    logos  = len(re.findall(LOGOS_RE,  all_text))
    ethos  = len(re.findall(ETHOS_RE,  all_text))
    pathos = len(re.findall(PATHOS_RE, all_text))

    if   logos > 0 and ethos > 0 and pathos > 0:  persuasion_mix = 90
    elif logos > 0 and (ethos > 0 or pathos > 0): persuasion_mix = 74
    elif logos > 3 and pathos == 0:               persuasion_mix = 52
    elif (logos + ethos + pathos) == 0:           persuasion_mix = 38
    else:                                          persuasion_mix = 62

    HEDGE_RE   = r'\b(?:maybe|perhaps|possibly|i think it might|i\'m not sure but|i guess)\b'
    hedge_hits = len(re.findall(HEDGE_RE, all_text))
    total_qa_words = sum(len(t.split()) for t in qa_texts)
    if total_qa_seconds > 30 and total_qa_words > 10:
        qa_wpm = round(total_qa_words / (total_qa_seconds / 60))
        if   160 <= qa_wpm <= 190:   command = 95
        elif 140 <= qa_wpm <= 220:   command = 76
        elif qa_wpm > 220:           command = 52
        else:                        command = 60
    else:
        command = 72 - min(20, hedge_hits * 8)

    return conclusion_first, persuasion_mix, max(0, command)


def _cq_heuristic_class_presentation(qa_texts):
    """Branch C: Class Presentation — returns (rule_of_three, conv_sense, illustrative_support)."""
    all_text = " ".join(qa_texts).lower()

    ORDINAL_RE    = r'\b(?:first(?:ly)?[,\s]|second(?:ly)?[,\s]|third(?:ly)?[,\s]|firstly|secondly|thirdly|number one|number two|number three)\b'
    THREE_FRAME_RE = r'\b(?:three (?:things|points|key|reasons|aspects|ways|parts)|3 (?:things|points|key))\b'
    ord_hits   = len(re.findall(ORDINAL_RE,    all_text))
    frame_hits = len(re.findall(THREE_FRAME_RE, all_text))

    if   frame_hits >= 1 and ord_hits >= 3: rule_of_three = 95
    elif ord_hits >= 3:                     rule_of_three = 80
    elif ord_hits >= 2:                     rule_of_three = 65
    elif ord_hits >= 1:                     rule_of_three = 50
    else:                                   rule_of_three = 35

    ORAL_RE = (
        r'\b(?:you know|think about|imagine|right\?|the thing is|what i mean is|'
        r'for example|in other words|actually|let me explain|you see|it\'s like)\b'
    )
    SCRIPT_RE = r'\b(?:furthermore|nevertheless|in conclusion|to summarize|accordingly|thus|hence|moreover)\b'
    oral_hits   = len(re.findall(ORAL_RE,   all_text))
    script_hits = len(re.findall(SCRIPT_RE, all_text))

    if   oral_hits >= 4 and script_hits <= 2: conv_sense = 90
    elif oral_hits >= 2:                      conv_sense = 74
    elif oral_hits >= 1:                      conv_sense = 60
    elif script_hits >= 3:                    conv_sense = 38
    else:                                     conv_sense = 50

    EXAMPLE_RE = (
        r'\b(?:for example|for instance|such as|like when|one (?:example|case|time|instance)|'
        r'i (?:remember|recall|once saw)|when i (?:was|worked|studied|saw)|'
        r'a (?:friend|colleague|classmate|professor|student) (?:told|said|shared))\b'
    )
    FIVE_W_RE = (
        r'\b(?:in (?:china|the us|beijing|shanghai|london|new york|tokyo)|'
        r'at (?:harvard|stanford|mit|oxford|cambridge|our university)|'
        r'last (?:year|month|week|semester)|professor \w+)\b'
    )
    ex_hits     = len(re.findall(EXAMPLE_RE, all_text))
    five_w_hits = len(re.findall(FIVE_W_RE,  all_text))

    if   ex_hits >= 2 and five_w_hits >= 1: illustrative = 92
    elif ex_hits >= 2:                      illustrative = 78
    elif ex_hits >= 1:                      illustrative = 62
    else:                                   illustrative = 38

    return rule_of_three, conv_sense, illustrative


def _cq_no_data_result(scene_slug):
    """Return placeholder when no Q&A transcripts exist."""
    _labels = {
        "thesis_defense":    "Thesis Defense",
        "case_pitch":        "Case Pitch",
        "class_presentation":"Class Presentation",
    }
    return {
        "has_data":              False,
        "scene_slug":            scene_slug,
        "scene_label":           _labels.get(scene_slug, "Communication"),
        "cq_total":              0,
        "cq_scores":             {},
        "dim_names":             [],
        "weights":               [],
        "what_i_did_good":       [],
        "areas_for_improvement": [],
        "exchange_count":        0,
        "communication_quality_report": {
            "overall_cq_score":      0,
            "per_question_analysis": [],
        },
    }


def _repair_truncated_json(raw):
    """
    Attempt to close a Gemini response truncated mid-JSON.
    Returns parsed dict or None.
    """
    chunk = raw.strip()
    # Walk character by character tracking string/depth state
    in_string   = False
    escape_next = False
    depth_brace   = 0
    depth_bracket = 0
    last_safe_pos = 0          # position of last safely-closed top-level value

    for i, c in enumerate(chunk):
        if escape_next:
            escape_next = False
            continue
        if c == "\\" and in_string:
            escape_next = True
            continue
        if c == '"':
            in_string = not in_string
            continue
        if not in_string:
            if   c == '{': depth_brace   += 1
            elif c == '}': depth_brace   -= 1
            elif c == '[': depth_bracket += 1
            elif c == ']': depth_bracket -= 1
            if depth_brace == 1 and depth_bracket == 0:
                last_safe_pos = i

    # If we're mid-string, truncate to last safe position
    if in_string and last_safe_pos > 0:
        chunk = chunk[:last_safe_pos + 1]
        # recount after truncation
        depth_brace   = chunk.count('{') - chunk.count('}')
        depth_bracket = chunk.count('[') - chunk.count(']')

    # Strip trailing commas before we add closing tokens
    chunk = re.sub(r',\s*$', '', chunk.rstrip())

    # Add missing closing tokens
    chunk += ']' * max(0, depth_bracket) + '}' * max(0, depth_brace)

    try:
        return json.loads(chunk)
    except Exception:
        return None


def _cq_mock_result(scene_slug, heuristic_scores, cq_total, dim_names,
                    exchange_count=0, ai_scores=None,
                    good_override=None, fix_override=None,
                    per_question_override=None):
    """
    Score-aware, scene-specific heuristic fallback.
    Used when Gemini is unavailable or returns corrupted JSON.
    Generates coaching text calibrated to the Python heuristic anchor scores —
    NOT generic placeholder strings.
    """
    _labels = {
        "thesis_defense":    "Thesis Defense",
        "case_pitch":        "Case Pitch",
        "class_presentation":"Class Presentation",
    }
    scene_label = _labels.get(scene_slug, "Communication")

    # Use AI-recovered scores when available (partial JSON repair), else heuristic
    scores_to_use = ai_scores if ai_scores else heuristic_scores
    scores        = [scores_to_use.get(d, 60) for d in dim_names]

    # ── Scene coaching tables (dim, high_praise, low_issue, how_to_fix) ───────
    COACHING = {
        "thesis_defense": [
            ("Directness",
             "You stated your position right away and answered the examiner's challenge directly, with no unnecessary lead-in — a core defense skill.",
             "You added too much background before making your point. In a defense, your core stance should appear in sentence 1 or 2.",
             "Start with your core argument first. Say this instead: 'My research demonstrates [X]. The evidence for this is [specific data].'"),
            ("Defensibility",
             "You backed your argument with concrete data, statistics, or citations, which strengthened your defense.",
             "Your answer relied on general reasoning without specific data. Academic examiners expect every claim to have evidence.",
             "Anchor every claim with data. Say this instead: 'Our n=[X] sample showed p<0.05 significance, confirming [conclusion].'"),
            ("Tact",
             "You handled the challenge diplomatically — acknowledging the examiner's concern before your rebuttal, a strong Carnegie 'Yes-Response'.",
             "You moved straight into rebuttal without finding common ground first. A Carnegie 'Yes-Response' builds rapport before you defend.",
             "Open with acknowledgment. Say this instead: 'That is a valid concern. I can see why you raise this — however, our data shows [answer].'"),
        ],
        "case_pitch": [
            ("Conclusion First",
             "You gave your conclusion before the supporting details — exactly the McKinsey Pyramid structure investors expect.",
             "You built context before reaching your conclusion. Investors lose interest within 10 seconds — your stance must come in the first sentence.",
             "Lead with your conclusion. Say this instead: 'Yes, [direct verdict]. The reason is [one key driver]. Here is the evidence: [metric].'"),
            ("Persuasion Mix",
             "Your answer balanced credibility (Ethos), data (Logos), and a customer story (Pathos) — a complete Aristotle rhetorical triangle.",
             "Your answer leaned too heavily on data (Logos), with little customer story (Pathos) or credibility signal (Ethos).",
             "Add one real customer story. Say this instead: 'One of our beta customers saw a [X%] improvement in [outcome] after [timeframe].'"),
            ("Command Presence",
             "You kept a confident, authoritative tone throughout the challenge — no hedging, no seeking approval.",
             "Under pressure, your language showed uncertainty — words like 'maybe', 'I think', 'perhaps'. Investors read this as a lack of conviction.",
             "Replace hedges with clear assertions. Say this instead: '[The answer] is [X]. We know this because [data]. There is no ambiguity.'"),
        ],
        "class_presentation": [
            ("Rule of Three",
             "You organized your answer into three clear parts, which made it easy to follow and remember.",
             "Your answer listed points but had no three-part structure. The Rule of Three helps your audience remember your main ideas.",
             "Structure your answer into three points. Say this instead: 'There are three key things: first, [X]; second, [Y]; and third, [Z].'"),
            ("Conversational Sense",
             "Your answer sounded like a natural conversation rather than a recited script, with oral connectors and a natural rhythm throughout.",
             "Your answer sounded formal or scripted. TED speakers talk like they are having a conversation, not giving a lecture.",
             "Add conversational language markers. Say this instead: 'Think about it this way — [your core point]. The key insight is [idea].'"),
            ("Illustrative Support",
             "You supported your point with a specific, real example, including clear details like time, place, and people — this made the idea vivid.",
             "You gave a conceptual explanation without a concrete 5-W example. A real story always beats abstract reasoning.",
             "Jump straight into a 5-W example. Say this instead: 'Let me give you a real case: in [year], [who] at [place] did [what], with result [outcome].'"),
        ],
    }

    coaching = COACHING.get(scene_slug, COACHING["thesis_defense"])
    # Pad if dim_names is shorter than coaching table
    while len(dim_names) < len(coaching):
        dim_names.append(coaching[len(dim_names)][0])

    what_i_did_good = []
    areas_for_improvement = []

    for i, (dim, high_praise, low_issue, fix_phrase) in enumerate(coaching):
        sc  = scores[i] if i < len(scores) else 60
        dim_label = dim_names[i] if i < len(dim_names) else dim

        if sc >= 70:
            what_i_did_good.append(
                f"[{scene_label}] {dim_label} ({sc}/100): {high_praise}"
            )
            areas_for_improvement.append({
                "dimension": dim_label,
                "issue":     f"[{scene_label}] {dim_label}: Good foundation ({sc}/100). Keep refining to push above 85.",
                "example":   f"Heuristic score: {sc}/100.",
                "how_to_fix": "Keep applying this technique, and focus on even more precise wording next time.",
            })
        else:
            what_i_did_good.append(
                f"[{scene_label}] {dim_label} ({sc}/100): You attempted this dimension. Focused practice will meaningfully improve your score."
            )
            areas_for_improvement.append({
                "dimension": dim_label,
                "issue":     f"[{scene_label}] {low_issue}",
                "example":   f"Heuristic score: {sc}/100. A full AI evaluation session provides exact-quote feedback.",
                "how_to_fix": fix_phrase,
            })

    return {
        "has_data":             True,
        "scene_slug":           scene_slug,
        "scene_label":          scene_label,
        "cq_total":             cq_total,
        "cq_scores":            scores_to_use,
        "dim_names":            dim_names,
        "weights":              [],
        "what_i_did_good":      good_override  if good_override  else what_i_did_good,
        "areas_for_improvement": fix_override  if fix_override   else areas_for_improvement,
        "exchange_count":       exchange_count,
        "communication_quality_report": {
            "overall_cq_score":      cq_total,
            "per_question_analysis": per_question_override if per_question_override is not None else [],
        },
    }


def _extract_dimension_quote(qa_texts, dim_slug, max_chars=165):
    """
    Extract the 1-2 most relevant sentences from qa_texts for ONE specific CQ dimension.

    Prevents the "same 180-char blob for all 3 cards" bug by scoring each sentence
    individually against dimension-specific linguistic signals.

    dim_slug values:
      thesis_defense:     "directness" | "defensibility" | "tact"
      case_pitch:         "conclusion_first" | "persuasion_mix" | "command_presence"
      class_presentation: "rule_of_three" | "conversational_sense" | "illustrative_support"
    """
    def _sentences(text):
        parts = re.split(r'(?<=[.!?])\s+', text.strip())
        return [p.strip() for p in parts if len(p.strip()) > 8]

    # Dimension → positive-signal patterns (high match = more relevant for this dim)
    _signals = {
        "rule_of_three": [
            r'\b(first(?:ly)?|second(?:ly)?|third(?:ly)?)\b',
            r'\bthere are (three|3)\b',
            r'\b(one[,\s]|two[,\s]|three[,\s])\b',
        ],
        "conversational_sense": [
            r'\b(you know|think about it|the thing is|imagine|basically|i mean|right\?)\b',
        ],
        "illustrative_support": [
            r'\b(for example|for instance|such as|specifically|like when)\b',
            r'\bin (20\d\d|19\d\d)\b',
            r'\b(result was|outcome|who was|at the time)\b',
        ],
        "directness": [
            r'\b(i found|my research|i argue|i believe|i show|i demonstrate|we found|our study|our data)\b',
            r'^(yes[,\s]|no[,\s])',
        ],
        "defensibility": [
            r'\d+\.?\d*\s*(%|percent)',
            r'\b(p\s*[<>=]\s*0\.\d+|n\s*=\s*\d+|participants|sample size|citation|reference)\b',
        ],
        "tact": [
            r'\b(valid concern|good point|you raise|that\'s true|i understand|fair point|that is a valid|i agree|i see your point)\b',
        ],
        "persuasion_mix": [
            r'\d+\.?\d*\s*(%|percent|growth|revenue)',
            r'\b(user|customer|client|story|told us|experienced|felt|reported|bought)\b',
        ],
        "command_presence": [
            r'\b(maybe|perhaps|kind of|sort of|might be|i think|i guess|not sure|possibly)\b',
        ],
        # conclusion_first: first sentence of any answer is always the most relevant signal
        "conclusion_first": None,
    }

    # conclusion_first: first sentence wins by definition
    if dim_slug == "conclusion_first":
        for answer in qa_texts:
            sents = _sentences(answer)
            if sents:
                s = sents[0]
                return f'You said: "{s[:max_chars]}{"…" if len(s) > max_chars else ""}"'
        fallback = " ".join(qa_texts).strip()
        return f'You said: "{fallback[:max_chars]}{"…" if len(fallback) > max_chars else ""}"'

    patterns = _signals.get(dim_slug, [])
    best_sent  = None
    best_score = -1

    for answer in qa_texts:
        for sent in _sentences(answer):
            score = 0
            sent_l = sent.lower()
            for pat in (patterns or []):
                score += len(re.findall(pat, sent_l, re.IGNORECASE))
            # Slight preference for medium-length sentences (signal-rich, not too short)
            wc = len(sent.split())
            if 6 <= wc <= 45:
                score += 0.3
            if score > best_score:
                best_score = score
                best_sent  = sent

    if best_sent:
        q = best_sent[:max_chars] + ("…" if len(best_sent) > max_chars else "")
        return f'You said: "{q}"'

    # Fallback: first sentence of first non-empty answer
    for answer in qa_texts:
        sents = _sentences(answer)
        if sents:
            s = sents[0]
            return f'You said: "{s[:max_chars]}{"…" if len(s) > max_chars else ""}"'

    all_text = " ".join(qa_texts).strip()
    return f'You said: "{all_text[:max_chars]}{"…" if len(all_text) > max_chars else ""}"'


def _build_cq_coaching_cards(qa_texts, scene_slug, scene_label, dim_names, scores):
    """
    Build CQ coaching cards that include the user's ACTUAL words in the
    'What you said' example field. Used when Gemini text generation fails or
    is truncated.

    Pattern-matches the real qa_texts against rubric criteria to produce
    specific (non-generic) feedback rather than filler placeholders.
    """
    all_text   = " ".join(qa_texts).strip()
    text_lower = all_text.lower()
    # user_quote is now extracted per-dimension in the loop via _extract_dimension_quote()

    # ── Pattern flags from actual user words ─────────────────────────────────
    has_first_person = any(p in text_lower for p in
        ["i found", "my research", "i argue", "i believe", "we found",
         "our study", "i show", "i demonstrate", "my analysis"])
    has_numbers      = bool(re.search(
        r'\d+\.?\d*\s*(%|percent|p\s*<|n\s*=|participants|samples)', text_lower))
    has_yes_resp     = any(p in text_lower for p in
        ["valid concern", "good point", "you raise", "that's true",
         "i understand", "i see your", "fair point", "that is a"])
    has_hedges       = any(p in text_lower for p in
        ["maybe", "perhaps", "kind of", "sort of", "might be",
         "i think", "i guess", "not sure"])
    has_structure    = any(p in text_lower for p in
        ["first", "second", "third", "firstly", "secondly", "thirdly",
         "one,", "two,", "three,"])
    has_example      = any(p in text_lower for p in
        ["for example", "for instance", "such as", "like when",
         "in 20", "in 19", "specifically"])

    # ── Scene-specific coaching tables ────────────────────────────────────────
    COACHING = {
        "thesis_defense": [
            (dim_names[0],
             has_first_person or scores.get(dim_names[0], 60) >= 65,
             "You used active voice and a first-person stance, showing the academic confidence of an independent researcher.",
             "You added too much background before stating your position. In a defense, your core stance should appear in sentence 1 or 2.",
             "Start with your core argument first. Say this instead: 'My research demonstrates [X]. The evidence is [specific data].'"),
            (dim_names[1],
             has_numbers or scores.get(dim_names[1], 60) >= 65,
             "You backed your argument with data, citations, or quantitative reasoning, which held up under the examiner's scrutiny.",
             "Your answer relied on general reasoning without specific data. Examiners expect every claim to be backed by numbers or citations.",
             "Anchor every claim with data. Say this instead: 'Our n=[X] sample showed [result], confirmed by [citation].'"),
            (dim_names[2],
             has_yes_resp or scores.get(dim_names[2], 60) >= 65,
             "You showed diplomatic skill — acknowledging the examiner's concern before your rebuttal.",
             "You moved straight into rebuttal without finding common ground first. A Carnegie 'Yes-Response' builds trust before you defend.",
             "Open with acknowledgment. Say this instead: 'That is a valid concern. Our data also shows [answer].'"),
        ],
        "case_pitch": [
            (dim_names[0],
             not has_hedges and scores.get(dim_names[0], 60) >= 65,
             "You gave your conclusion before the supporting details — exactly the McKinsey Pyramid structure investors expect.",
             "You built context before reaching your conclusion. Investors lose interest within 10 seconds — your stance must come in the first sentence.",
             "Lead with your conclusion. Say this instead: 'Yes, [direct verdict]. The reason: [one key driver].'"),
            (dim_names[1],
             (has_numbers or has_example) and scores.get(dim_names[1], 60) >= 65,
             "Your answer blended credibility, data, and a story — a strong Aristotle rhetorical triangle.",
             "Your answer leaned on logic, with little customer story or credibility signal.",
             "Add one real customer story. Say this instead: 'One customer saw [X%] improvement after [timeframe].'"),
            (dim_names[2],
             not has_hedges and scores.get(dim_names[2], 60) >= 65,
             "You kept a confident, authoritative tone throughout — no hedging, no seeking approval.",
             "Under pressure, your language showed uncertainty (maybe, I think, perhaps). Investors read this as a lack of conviction.",
             "Replace hedges with clear assertions. Say this instead: 'This is [X]. We know because [data].'"),
        ],
        "class_presentation": [
            (dim_names[0],
             has_structure and scores.get(dim_names[0], 60) >= 65,
             "You organized your answer into three clear parts, which made it easy to follow and remember.",
             "Your answer listed points but had no three-part structure. The Rule of Three improves audience retention.",
             "Structure your answer into three points. Say this instead: 'There are three things: first [X]; second [Y]; third [Z].'"),
            (dim_names[1],
             scores.get(dim_names[1], 60) >= 65,
             "Your answer sounded like a natural conversation rather than a recited script.",
             "Your answer sounded formal or scripted. TED speakers talk like they are having a conversation, not giving a lecture.",
             "Add conversational language markers. Say this instead: 'Think about it this way — [core point]. The key is [insight].'"),
            (dim_names[2],
             has_example and scores.get(dim_names[2], 60) >= 65,
             "You supported your point with a specific, real example, including clear details like time, place, and people.",
             "You gave a conceptual explanation without a concrete 5-W example. A real story always beats abstract reasoning.",
             "Jump straight into a 5-W example. Say this instead: 'In [year], [who] at [place] did [what], result: [outcome].'"),
        ],
    }

    coaching   = COACHING.get(scene_slug, COACHING["thesis_defense"])
    score_list = [scores.get(d, 60) for d in dim_names]

    # Map each coaching card index → its dimension slug for targeted quote extraction
    _dim_slug_map = {
        "thesis_defense":    ["directness",       "defensibility",    "tact"],
        "case_pitch":        ["conclusion_first",  "persuasion_mix",   "command_presence"],
        "class_presentation":["rule_of_three",     "conversational_sense", "illustrative_support"],
    }
    dim_slugs = _dim_slug_map.get(scene_slug, ["directness", "defensibility", "tact"])

    what_i_did_good      = []
    areas_for_improvement = []

    for i, (dim, is_strong, good_msg, bad_msg, fix) in enumerate(coaching):
        sc         = score_list[i] if i < len(score_list) else 60
        dim_label  = dim_names[i]  if i < len(dim_names)  else dim
        dim_slug_i = dim_slugs[i]  if i < len(dim_slugs)  else "directness"

        # Per-dimension targeted quote — extracts the 1-2 sentences most
        # relevant to THIS specific dimension, not the same blob for all 3 cards
        if qa_texts:
            per_dim_quote = _extract_dimension_quote(qa_texts, dim_slug_i)
        else:
            per_dim_quote = "No Q&A response was recorded for this session."

        if is_strong:
            what_i_did_good.append(f"[{scene_label}] {dim_label} ({sc}/100): {good_msg}")
        else:
            what_i_did_good.append(
                f"[{scene_label}] {dim_label} ({sc}/100): "
                "You attempted this dimension. Focused practice will improve your score."
            )

        areas_for_improvement.append({
            "dimension": dim_label,
            "issue": (
                f"[{scene_label}] {bad_msg}"
                if sc < 70 else
                f"[{scene_label}] {dim_label}: Good foundation ({sc}/100). Keep refining."
            ),
            "example":    per_dim_quote,   # ← targeted quote for THIS dimension only
            "how_to_fix": (
                fix if sc < 70
                else "Continue applying the same technique. Aim for consistency above 80."
            ),
        })

    return what_i_did_good, areas_for_improvement


def _nearest_slide_for_text(slides, *texts):
    """
    Heuristic slide-matcher: picks the slide whose title/content shares the
    most keyword overlap with the given question/answer text(s). Used by the
    non-AI fallback path to ground feedback in real PPT content instead of
    generic advice, even when the LLM is unavailable.
    """
    if not slides:
        return None
    combined = " ".join(t or "" for t in texts).lower()
    words = set(re.findall(r"[a-z]{4,}", combined))
    if not words:
        return slides[0]
    best, best_score = None, -1
    for s in slides:
        slide_words = set(re.findall(r"[a-z]{4,}", (s.get("title", "") + " " + s.get("content", "")).lower()))
        score = len(words & slide_words)
        if score > best_score:
            best, best_score = s, score
    return best or slides[0]


def _dedupe_pqa(pqa):
    """
    Some LLM responses emit one per_question_analysis entry per exchange AND
    per dimension pass, producing duplicate cards for the same question in
    the report UI. Keep only the first entry per question_id/question_text.
    """
    if not isinstance(pqa, list):
        return pqa
    seen = set()
    deduped = []
    for item in pqa:
        if not isinstance(item, dict):
            continue
        key = item.get("question_id") or item.get("question_text", "")
        if key in seen:
            continue
        seen.add(key)
        deduped.append(item)
    return deduped


def _cq_signal_scan(answer):
    """
    Extract lightweight, scene-agnostic linguistic signals from a single answer,
    used to ground the local per-question fallback analysis in the ACTUAL text
    of THIS specific answer — never a fixed template. Each signal also records
    which sentence (if any) triggered it, so feedback can quote the exact point
    in the answer that demonstrates (or fails to demonstrate) a dimension.
    """
    sentences = re.split(r'(?<=[.!?])\s+', answer.strip())
    sentences = [s.strip() for s in sentences if s.strip()] or [answer.strip()]
    first = sentences[0]

    def find(pattern, in_sentences=None):
        for s in (in_sentences or sentences):
            if re.search(pattern, s, re.I):
                return s
        return None

    return {
        "sentences": sentences,
        "first": first,
        "first_person_or_data_open": find(r"^\s*(i|we|our|my|the data|the results?|this (study|research))\b", [first]),
        "verdict_open":              find(r"^\s*(yes|no|we believe|absolutely|definitely|our (view|answer) is)\b", [first]),
        "evidence":                  find(r"\d|%|percent|p\s*<\s*0|p-value|significant|according to|studies? show|citation"),
        "acknowledgment":            find(r"\b(valid concern|fair point|i see why|good question|you raise|that'?s true|that is a (valid|fair))\b"),
        "hedge":                     find(r"\b(maybe|perhaps|i think|i guess|kind of|sort of|probably|might|not sure)\b"),
        "customer_or_metric":        find(r"\d|%|customer|user|client|beta|case study|revenue"),
        "oral_marker":               find(r"\b(you know|think about it|the thing is|basically|honestly)\b"),
        "example_marker":            find(r"\b(for example|for instance|such as|in \d{4}|last (year|month)|e\.g\.)\b"),
        "three_point_structure":     bool(re.search(r"\bfirst(ly)?\b.{0,150}\bsecond(ly)?\b.{0,200}\b(third(ly)?|finally|lastly)\b", answer, re.I | re.S)),
    }


def _dim_hit_and_sentence(sig, key):
    """Resolve one dimension's (hit, quotable_sentence) from a pre-computed signal scan."""
    if key == "no_hedge":
        s = sig["hedge"]
        # When no hedge phrase exists, quote the actual first sentence as the
        # concrete example of confident language — never a bare "no hedge found".
        return (s is None), (s if s is not None else sig["first"])
    if key == "three_point_structure":
        return sig["three_point_structure"], None
    s = sig.get(key)
    return (s is not None), s


_STOPWORDS_FOR_KEYWORD_OVERLAP = {
    "the", "a", "an", "is", "are", "was", "were", "do", "does", "did", "you", "your",
    "how", "what", "why", "which", "that", "this", "of", "to", "in", "on", "for", "and",
    "or", "with", "about", "it", "as", "will", "would", "could", "can", "i", "we", "our",
}


def _keyword_overlap_highlight(question, sentences):
    """
    When an answer doesn't trip any scene dimension signal, still ground the
    feedback in a specific point rather than a generic template: find a
    substantive keyword from THIS question that the user actually echoed in
    THIS answer, and return that keyword plus the sentence carrying it. Both
    the keyword and the sentence differ per question/answer pair.
    """
    q_keywords = [
        w.strip(".,?!'\"").lower() for w in question.split()
        if len(w.strip(".,?!'\"")) > 3 and w.strip(".,?!'\"").lower() not in _STOPWORDS_FOR_KEYWORD_OVERLAP
    ]
    for kw in q_keywords:
        for s in sentences:
            if kw in s.lower():
                return kw, s
    return None, None


# Scene-specific content-signal config used ONLY to pick which real point in the
# user's own answer to highlight — the rubric dimension names themselves are never
# printed in the generated text (feedback must read as a plain observation about
# what was actually said, not a citation of the scored framework).
_FALLBACK_DIM_CONFIG = {
    # Each rule is (internal_key, signal_key, good_desc, missing_desc, fix_hint).
    # internal_key is ONLY used to pick between rules and is never surfaced in the
    # generated text — feedback must read as a direct observation about what the
    # user said, not as a label pulled from the scored rubric.
    "thesis_defense": [
        ("directness", "first_person_or_data_open",
         "you led straight in with your own stance instead of padding with background",
         "your first sentence didn't open with a direct claim — it drifted into context before answering",
         "open with your actual position in the first sentence"),
        ("defensibility", "evidence",
         "you backed your claim with a concrete number or fact instead of a vague adjective",
         "your answer leaned on subjective words instead of a number, percentage, or citation",
         "add one hard number or fact to back the claim"),
        ("tact", "acknowledgment",
         "you acknowledged the concern before pushing back on it",
         "you responded to the challenge without first acknowledging the concern raised",
         "acknowledge the concern first, then defend your position"),
    ],
    "case_pitch": [
        ("conclusion_first", "verdict_open",
         "you gave your verdict right away instead of building up to it",
         "your first sentence didn't open with a direct yes/no/verdict — it built up to the point instead "
         "of leading with it",
         "state your verdict in the very first sentence"),
        ("persuasion_mix", "customer_or_metric",
         "you backed your point with a data point or a customer reference, giving it real weight",
         "your answer had no data point and no customer story to back up the pitch",
         "add a data point or a short customer story"),
        ("command_presence", "no_hedge",
         "you answered with confident, assertive language and no hedging",
         "your answer used hedging language, which weakens your authority under challenge",
         "cut the hedging words and state it plainly"),
    ],
    "class_presentation": [
        ("rule_of_three", "three_point_structure",
         "you structured your answer into clear first/second/third-style points",
         "your answer wasn't structured into clear points (first / second / third)",
         "break your answer into three clearly labeled points"),
        ("conversational_sense", "oral_marker",
         "you spoke naturally instead of sounding like you were reading a script",
         "your answer sounded formal or scripted, without natural spoken connectors",
         "add a natural spoken connector like 'you know' or 'the thing is'"),
        ("illustrative_support", "example_marker",
         "you grounded your point in a concrete example instead of an abstract statement",
         "your answer gave a conceptual explanation with no concrete example",
         "add a specific example with a who/when/where detail"),
    ],
}


def _fallback_per_question_analysis(comm_transcripts, slides=None, scene_slug=None):
    """
    Local, English, IELTS 5.5-6.0-level per-question analysis used whenever the
    LLM call fails, times out, or returns malformed JSON. Guarantees the
    frontend always receives a non-empty communication_quality_report so the
    per-question card layout never breaks.

    IMPORTANT: this is NOT a single canned template reused for every question.
    Each answer is scanned for content signals (see _FALLBACK_DIM_CONFIG) and
    the specific point in THIS answer that satisfies or misses one of those
    signals is what gets quoted and explained in plain language — never by
    naming the underlying rubric dimension — so two different answers produce
    genuinely different, content-grounded feedback.
    """
    dim_config = _FALLBACK_DIM_CONFIG.get(scene_slug, _FALLBACK_DIM_CONFIG["class_presentation"])
    items = []
    for t in comm_transcripts:
        answer = t.get("answer", "").strip()
        if not answer:
            continue
        question = t.get("question", "") or "(question not recorded)"
        words = answer.split()
        short_answer = len(words) <= 4
        # Always quote the user's COMPLETE answer for THIS question, never a
        # truncated fragment — `answer` comes from a single isolated
        # comm_transcripts entry (one question -> one answer), so this quote
        # can never bleed in words from a different question.
        quoted_fragment = answer

        sig = _cq_signal_scan(answer)
        rules_evaluated = []
        for key, signal_key, good_desc, missing_desc, fix_hint in dim_config:
            hit, sentence = _dim_hit_and_sentence(sig, signal_key)
            rules_evaluated.append({
                "key": key, "hit": hit, "sentence": sentence,
                "good_desc": good_desc, "missing_desc": missing_desc, "fix_hint": fix_hint,
            })

        hit_rule = next((d for d in rules_evaluated if d["hit"]), None)
        miss_rule = next(
            (d for d in rules_evaluated if not d["hit"] and (not hit_rule or d["key"] != hit_rule["key"])),
            None,
        )
        if miss_rule is None:
            miss_rule = next((d for d in rules_evaluated if not d["hit"]), None)

        if short_answer:
            what_good = (
                f"You said '{quoted_fragment}', which correctly names the core keyword the question "
                "was asking about. That shows you understood what was being asked, though there is not "
                "enough here yet to judge the answer in depth."
            )
        elif hit_rule:
            quote_bit = f" — you said '{hit_rule['sentence'].strip()}'" if hit_rule["sentence"] else ""
            what_good = f"{hit_rule['good_desc'][0].upper()}{hit_rule['good_desc'][1:]}{quote_bit}."
        else:
            # No content signal fired at all — still ground the praise in a specific
            # point from THIS answer instead of a repeated generic line, by finding a
            # keyword from THIS question that the user actually echoed.
            kw, kw_sentence = _keyword_overlap_highlight(question, sig["sentences"])
            if kw and kw_sentence:
                what_good = (
                    f"You directly engaged with the question's focus on '{kw}' by saying "
                    f"'{kw_sentence.strip()}', which kept your answer anchored to what was actually asked."
                )
            else:
                what_good = (
                    f"When you said '{quoted_fragment}', you gave a real, on-topic answer instead of "
                    "avoiding the question."
                )

        slide = _nearest_slide_for_text(slides, question, answer)
        slide_title   = (slide or {}).get("title", "")
        slide_content = (slide or {}).get("content", "")

        if miss_rule:
            # Always anchor the gap in a concrete sentence from THIS answer — either the
            # exact sentence that failed the check, or (if the gap is structural, i.e. the
            # signal is absent everywhere) the longest sentence in the answer, so the same
            # kind of gap across different answers still points at different real text.
            gap_quote = miss_rule["sentence"] or max(sig["sentences"], key=len)
            gap_line = (
                f"{miss_rule['missing_desc'][0].upper()}{miss_rule['missing_desc'][1:]}. For example, "
                f"you said '{gap_quote.strip()}', which still doesn't show it."
            )
        else:
            gap_line = "Your answer already covers the key things this question was looking for."
        fix_hint = (miss_rule or hit_rule or {}).get("fix_hint", "make your core point more concrete")

        if slide_title or slide_content:
            areas_improve = (
                f"{gap_line} On top of that, your slide on '{slide_title}' has more specific content "
                f"({slide_content[:120].strip()}{'…' if len(slide_content) > 120 else ''}) that you did "
                "not use — naming it would make your answer even stronger."
            )
            # NOTE: how_to_fix must be a clean, ready-to-use rewritten answer — never a
            # quote of the user's own (possibly garbled or incomplete) original words.
            how_fix = (
                f"Say this instead: 'To {fix_hint}, {slide_content[:160].strip()}"
                f"{'…' if len(slide_content) > 160 else ''} — that is why this matters for my presentation.'"
            )
        else:
            areas_improve = (
                f"{gap_line} Your answer is also too short to fully cover the question — add one "
                "concrete number, name, or fact to back it up."
            )
            how_fix = (
                f"Say this instead: 'To {fix_hint}, [state your core point clearly], "
                "and to be specific, the key evidence is [name one exact fact, number, or example from "
                "your own presentation].'"
            )

        items.append({
            "question_id":            t.get("question_id", ""),
            "question_text":          question,
            "user_actual_answer":     answer,
            "what_i_did_good":        what_good,
            "areas_for_improvement":  areas_improve,
            "how_to_fix":             how_fix,
        })
    return items


def run_communication_quality_evaluation(qa_answers, config, fe_qa_history=None,
                                          scene_slug=None, total_qa_seconds=0,
                                          slides=None, qa_bank=None):
    """
    Communication Quality (CQ) evaluation engine.
    Evaluates ONLY Q&A + interrupt exchanges — NOT the presentation narration.
    Returns a CQ result dict with 3 scene-specific dimension scores + Gemini feedback.
    """
    audience   = config.get("audience",   "Professor")
    difficulty = config.get("difficulty", "Medium")
    scenario   = config.get("scenario",   "Class Presentation")

    if not scene_slug:
        _slug_map = {
            "Thesis Defense":       "thesis_defense",       # THESIS_DEFENSE_CHALLENGE_POOL → Directness/Defensibility/Tact
            "MBA Case Pitch":       "case_pitch",           # CASE_PITCH_CHALLENGE_POOL → Conclusion First/Persuasion Mix/Command Presence
            "Class Presentation":   "class_presentation",   # TED_QA_MATRIX → Rule of Three/Conversational/Illustrative
            "Academic Presentation":"class_presentation",   # backward compat: old sessions stored this name
        }
        scene_slug = _slug_map.get(scenario, "thesis_defense")

    _labels = {
        "thesis_defense":    "Thesis Defense",
        "case_pitch":        "Case Pitch",
        "class_presentation":"Class Presentation",
    }
    # Preserve the user's original scenario name as the display label so that
    # "Academic Presentation" is NOT relabelled to "Thesis Defense" in the report,
    # even though both share the same internal rubric slug.
    scene_label = scenario if scenario else _labels.get(scene_slug, "Communication")

    # ── Collect Q&A exchange transcripts ──────────────────────────────────────
    # Build from frontend history (interrupt + academic Q&A captured in JS)
    fe_transcripts = []
    if fe_qa_history:
        ai_qs   = [h for h in fe_qa_history if h.get("role") == "ai"]
        user_as = [h for h in fe_qa_history if h.get("role") == "user"]
        for i, ua in enumerate(user_as):
            ai_item = ai_qs[i] if i < len(ai_qs) else {}
            fe_transcripts.append({
                "question":           ai_item.get("text", ""),
                "answer":             ua.get("text", ""),
                "type":               ua.get("type", "qa_answer"),
                "answering_strategy": ai_item.get("strategy", ""),
                "question_id":        f"interrupt_{i + 1}",
            })

    # Build from session answers (server-side Q&A records from submit-answer /
    # submit-academic-qa routes — may have richer text than the frontend capture)
    # Also carry question_type / anchor metadata for dual-track detection.
    session_transcripts = []
    for a in qa_answers:
        if a.get("type") in ("qa_answer", "academic_qa"):
            session_transcripts.append({
                "question":           a.get("question",        ""),
                "answer":             a.get("text",            ""),
                "type":               a.get("type",            ""),
                "question_type":      a.get("question_type",   "free"),
                "anchor_type":        a.get("anchor_type",     ""),
                "target_dim":         a.get("target_dim",      ""),
                "scaffold_signal":    a.get("scaffold_signal", ""),
                "question_id":        a.get("question_id",     ""),
                "answering_strategy": a.get("answering_strategy", ""),
            })

    # Use whichever source has more substantive text content
    fe_words   = sum(len(t["answer"].split()) for t in fe_transcripts)
    sess_words = sum(len(t["answer"].split()) for t in session_transcripts)
    app.logger.info(
        f"[CQ SRC] fe_exchanges={len(fe_transcripts)} fe_words={fe_words} | "
        f"sess_exchanges={len(session_transcripts)} sess_words={sess_words}"
    )

    if sess_words > fe_words:
        comm_transcripts = [t for t in session_transcripts if t["answer"].strip()]
    else:
        comm_transcripts = [t for t in fe_transcripts if t["answer"].strip()]

    # Merge any session entries the frontend didn't capture. Guard against
    # re-appending entries already present in comm_transcripts (this happened
    # when comm_transcripts was seeded from session_transcripts itself and
    # fe_transcripts was empty/smaller — every session entry looked "not yet
    # captured by the frontend" and got appended a second time, doubling the
    # exchange count and duplicating per-question analysis cards).
    existing_answers_set = {t["answer"].strip() for t in comm_transcripts}
    for t in session_transcripts:
        ans = t["answer"].strip()
        if ans and ans not in existing_answers_set:
            comm_transcripts.append(t)
            existing_answers_set.add(ans)

    # ── Annotate comm_transcripts with question_type from session metadata ──────
    # session_transcripts carry question_type; fe_transcripts may not.
    # Match by stripped answer text so question_type propagates correctly.
    _sess_meta = {
        t["answer"].strip(): t
        for t in session_transcripts if t["answer"].strip()
    }
    for t in comm_transcripts:
        ans  = t["answer"].strip()
        meta = _sess_meta.get(ans, {})
        t.setdefault("question_type",      meta.get("question_type",      "free"))
        t.setdefault("anchor_type",        meta.get("anchor_type",        ""))
        t.setdefault("target_dim",         meta.get("target_dim",         ""))
        t.setdefault("scaffold_signal",    meta.get("scaffold_signal",    ""))
        t.setdefault("question_id",        meta.get("question_id",        ""))
        t.setdefault("answering_strategy", meta.get("answering_strategy", ""))

    # ── Backfill missing question_id / answering_strategy from qa_bank ────────
    _qa_bank_by_id = {q.get("id"): q for q in (qa_bank or []) if q.get("id")}
    for i, t in enumerate(comm_transcripts):
        if not t.get("question_id"):
            t["question_id"] = f"q{i + 1}"
        if not t.get("answering_strategy"):
            _qb = _qa_bank_by_id.get(t["question_id"])
            if _qb:
                t["answering_strategy"] = _qb.get("answering_strategy", "")

    if not comm_transcripts or all(not t["answer"].strip() for t in comm_transcripts):
        app.logger.info("[CQ] No Q&A transcript data — returning no_data placeholder.")
        return _cq_no_data_result(scene_slug)

    qa_texts       = [t["answer"].strip() for t in comm_transcripts if t["answer"].strip()]
    exchange_count = len(comm_transcripts)

    # ── Dual-track routing: anchor question present → use Module 3 evaluator ──
    _anchor_ts = [t for t in comm_transcripts if t.get("question_type") == "anchor" and t["answer"].strip()]
    _free_ts   = [t for t in comm_transcripts if t.get("question_type") != "anchor"  and t["answer"].strip()]
    if _anchor_ts and _free_ts:
        app.logger.info(
            f"[CQ DUAL-TRACK] {len(_free_ts)} free + {len(_anchor_ts)} anchor exchange(s) detected. "
            "Routing to dual-track evaluator."
        )
        return _run_dual_track_cq_evaluation(
            _free_ts, _anchor_ts, scene_slug, scene_label,
            audience, difficulty, exchange_count, slides=slides
        )
    # No anchor detected → fall through to single-track evaluation (backward compat)

    # ── Python heuristic pre-scoring ──────────────────────────────────────────
    if scene_slug == "thesis_defense":
        s1, s2, s3 = _cq_heuristic_thesis_defense(qa_texts)
        dim_names   = ["Directness", "Defensibility", "Tact"]
        weights     = [0.40, 0.40, 0.20]
    elif scene_slug == "case_pitch":
        s1, s2, s3 = _cq_heuristic_case_pitch(qa_texts, total_qa_seconds)
        dim_names   = ["Conclusion First", "Persuasion Mix", "Command Presence"]
        weights     = [0.40, 0.30, 0.30]
    else:  # class_presentation
        s1, s2, s3 = _cq_heuristic_class_presentation(qa_texts)
        dim_names   = ["Rule of Three", "Conversational Sense", "Illustrative Support"]
        weights     = [0.40, 0.30, 0.30]

    heuristic_scores      = {dim_names[0]: int(s1), dim_names[1]: int(s2), dim_names[2]: int(s3)}
    cq_total_heuristic    = int(round(s1 * weights[0] + s2 * weights[1] + s3 * weights[2]))

    if not AI_ENABLED:
        return _cq_mock_result(scene_slug, heuristic_scores, cq_total_heuristic, dim_names, exchange_count)

    # ── Format exchanges for prompt ───────────────────────────────────────────
    # Sentences within each answer are labelled [1] [2] [3]… so the LLM can
    # pinpoint the exact sentence relevant to each dimension without repeating
    # the full answer blob across all three example fields.
    def _label_sentences(text):
        parts = re.split(r'(?<=[.!?])\s+', text.strip())
        labeled = [f"[{j+1}] {p.strip()}" for j, p in enumerate(parts) if p.strip()]
        return " ".join(labeled) if labeled else text

    exchanges_text = "\n\n".join(
        f"Q{i+1} [id={t.get('question_id', f'q{i+1}')}] [{t.get('type','qa')}]: {t['question']}\n"
        f"A{i+1}: {_label_sentences(t['answer'])}"
        for i, t in enumerate(comm_transcripts) if t["answer"].strip()
    ) or "(No Q&A answers recorded)"

    # ── Scene-specific rubric + feedback template strings ─────────────────────
    if scene_slug == "thesis_defense":
        rubric_text = (
            "RUBRIC — THESIS DEFENSE (0-100 each)\n\n"
            f"DIM 1 — {dim_names[0]} (weight 40%)\n"
            "90-100: First sentence directly addresses the challenge. Strong verbs (We found / The data shows). No dodging.\n"
            "70-89:  Mostly direct but 1 filler sentence before the real answer.\n"
            "50-69:  Opens with 2+ background sentences before addressing the question.\n"
            "0-49:   Consistently avoids question — background padding instead of stance.\n"
            "PENALTY: 2 consecutive background sentences = Question Dodging (−15 pts).\n\n"
            f"DIM 2 — {dim_names[1]} (weight 40%)\n"
            "90-100: Objective evidence — numbers, percentages, p-values, literature citations.\n"
            "70-89:  Mostly evidence-based, some unsupported assertions.\n"
            "50-69:  Primarily subjective language (I think / very / quite).\n"
            "0-49:   All subjective adjectives, no quantitative defense.\n\n"
            f"DIM 3 — {dim_names[2]} (weight 20%)\n"
            "90-100: Carnegie Yes-Response before every rebuttal (That is a valid concern...).\n"
            "70-89:  Acknowledgment in some responses.\n"
            "50-69:  Neutral — neither diplomatic nor combative.\n"
            "0-49:   Directly contradicts without any acknowledgment.\n\n"
            f"HEURISTIC ANCHORS (adjust ±15 based on actual text):\n"
            f"- {dim_names[0]}: {s1}/100\n"
            f"- {dim_names[1]}: {s2}/100\n"
            f"- {dim_names[2]}: {s3}/100"
        )
        good_templates = (
            "what_i_did_good — EXACTLY 3 strings. Insert EXACT words from A1/A2/A3 in every [bracket]:\n"
            f'1. "[Thesis Defense] Academic Voice & Directness: You used the first person and active voice efficiently, saying: \'[EXACT first-person quote like I found / My research shows / I argue from A1/A2/A3]\'. This projected an independent researcher image."\n'
            f'2. "[Thesis Defense] Defensibility: You backed your claim with hard evidence when you said: \'[EXACT quote containing data/p-values/citation/n=X from A1/A2/A3]\'. This showed scientific rigor and withstood the examiner\'s scrutiny."\n'
            f'3. "[Thesis Defense] Directness & Tact — Yes-Response: When challenged, you used a strong Yes-Response by saying: \'[EXACT Carnegie acknowledgment from A1/A2/A3, e.g. That is a valid concern / You raise an important point]\'. You built consensus before defending."\n\n'
            "areas_for_improvement — EXACTLY 3 objects with keys: dimension, issue, example, how_to_fix:\n"
            f'1. dimension="{dim_names[0]}", issue="[Thesis Defense] Question Dodging: When asked about methodology limitations, you gave background padding instead of a direct stance.", '
            f'example="You said: \'[EXACT dodgy opening words from A1/A2/A3 that avoid a direct answer — copy verbatim]\'", '
            f'how_to_fix="Address the flaw directly. Say this instead: \'That is a critical point. While our cohort has limitations, our statistical power analysis indicates...\'"\n'
            f'2. dimension="{dim_names[1]}", issue="[Thesis Defense] Low Persuasiveness: Your response relied on subjective language with no hard evidence.", '
            f'example="You said: \'[EXACT quote using very / I think / quite / somewhat instead of data from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Inject one hard fact. Say this instead: \'Our data shows [specific number/%, p-value] which indicates [direct conclusion]...\'"\n'
            f'3. dimension="{dim_names[2]}", issue="[Thesis Defense] Blunt Rebuttal: You responded to the challenge without acknowledging the examiner\'s concern first.", '
            f'example="You said: \'[EXACT blunt opener from A1/A2/A3 that jumps straight to rebuttal — copy verbatim]\'", '
            f'how_to_fix="Open with a Carnegie Yes-Response. Say this instead: \'That is a critical point. I see why you raise this. However, our [evidence] demonstrates...\'"\n'
        )

    elif scene_slug == "case_pitch":
        rubric_text = (
            "RUBRIC — MBA CASE PITCH / VENTURE PITCH (0-100 each)\n\n"
            f"DIM 1 — {dim_names[0]} (weight 40%)\n"
            "McKinsey Pyramid Principle: Answer → Reason → Evidence.\n"
            "90-100: Starts with a clear stance (Yes/No/We believe). Second sentence uses a logic pillar (Because / The reason is).\n"
            "70-89:  Usually answers first, occasionally buries the lede.\n"
            "50-69:  Leads with context before the answer.\n"
            "0-49:   Builds to the conclusion — VC loses interest.\n"
            "PENALTY: First sentence is context with no directional stance → −15 pts.\n\n"
            f"DIM 2 — {dim_names[1]} (weight 30%)\n"
            "Aristotle: Ethos (credibility) + Logos (data/metrics) + Pathos (customer story/pain).\n"
            "90-100: All three present.\n70-89: Two of three.\n50-69: Only one — e.g., all data, no story.\n0-49: No persuasion elements.\n\n"
            f"DIM 3 — {dim_names[2]} (weight 30%)\n"
            "90-100: Measured, confident responses. No excessive hedging (maybe/perhaps/I think it might).\n"
            "70-89:  Mostly authoritative, small hedges.\n50-69: Frequent hedges, hesitant language.\n0-49: Rambling or incoherent under challenge.\n\n"
            f"HEURISTIC ANCHORS:\n"
            f"- {dim_names[0]}: {s1}/100\n"
            f"- {dim_names[1]}: {s2}/100\n"
            f"- {dim_names[2]}: {s3}/100"
        )
        good_templates = (
            "what_i_did_good — EXACTLY 3 strings. Insert EXACT words from A1/A2/A3 in every [bracket]:\n"
            f'1. "[Case Pitch] Conclusion First & Prize Frame: You gave a clear stance and flipped the power dynamics by saying: \'[EXACT first sentence showing directional stance or prize-frame positioning from A1/A2]\'. Excellent corporate delivery."\n'
            f'2. "[Case Pitch] Persuasion Mix — Ethos + Logos + Pathos: You balanced data and human story. You cited \'[EXACT logos/data quote from A1/A2/A3]\' and humanized it with \'[EXACT pathos/customer-story quote from A1/A2/A3]\'. Aristotle\'s rhetorical triangle in action."\n'
            f'3. "[Case Pitch] Command Presence: When challenged, you held your ground without hedging. You said: \'[EXACT confident authoritative quote from A1/A2/A3]\'. No maybe, no I think — exactly what VCs respect."\n\n'
            "areas_for_improvement — EXACTLY 3 objects with keys: dimension, issue, example, how_to_fix:\n"
            f'1. dimension="{dim_names[0]}", issue="[Case Pitch] No Pyramid Structure: When asked about [topic], you built up to the conclusion instead of leading with your stance.", '
            f'example="You said: \'[EXACT bottom-up opening that buries the answer at the end from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Lead with your stance. Say this instead: \'Yes, [direct answer]. The reason is [one logic pillar]. Here is the evidence: [data point].\'"\n'
            f'2. dimension="{dim_names[1]}", issue="[Case Pitch] Low Pathos: Your response was purely Logos — raw data with no human story or customer pain point.", '
            f'example="You said: \'[EXACT data-only quote without any customer story or emotional hook from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Inject a real customer success story. Say this instead: \'To illustrate this, one of our beta users recently reported a 40% cost reduction by...\'"\n'
            f'3. dimension="{dim_names[2]}", issue="[Case Pitch] Hedging Under Pressure: When challenged, your language showed uncertainty or approval-seeking.", '
            f'example="You said: \'[EXACT hedging quote with maybe/I think/what do you think/is that okay from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Replace hedges with confident assertions. Say this instead: \'[Same idea restated with decisive command-presence language — no questions, no qualifiers].\'"\n'
        )

    else:  # class_presentation
        rubric_text = (
            "RUBRIC — CLASS PRESENTATION Q&A (0-100 each)\n\n"
            f"DIM 1 — {dim_names[0]} (weight 40%)\n"
            "90-100: Explicitly structures into 3 named buckets (there are three things: first...second...third).\n"
            "70-89:  Uses sequential markers (First / Then / Finally) without declaring the 3-point frame.\n"
            "50-69:  One or two ordinal markers, otherwise unstructured.\n"
            "0-49:   Stream of consciousness. No discernible structure.\n\n"
            f"DIM 2 — {dim_names[1]} (weight 30%)\n"
            "90-100: Natural conversation, not recited script. Uses oral markers (You know / Think about it / The thing is).\n"
            "70-89:  Some conversational elements.\n50-69: Mixed natural and scripted.\n0-49: Reads like a paper. Formal academic language only.\n\n"
            f"DIM 3 — {dim_names[2]} (weight 30%)\n"
            "Carnegie: Jump immediately into a 5-W example (Who, What, When, Where, Why/Outcome).\n"
            "90-100: Concrete example with specific details (name, place, year, outcome).\n"
            "70-89:  Examples provided but lack 5-W specificity.\n50-69: Generic ('a company might...'). 0-49: No examples — pure abstract concept.\n\n"
            f"HEURISTIC ANCHORS:\n"
            f"- {dim_names[0]}: {s1}/100\n"
            f"- {dim_names[1]}: {s2}/100\n"
            f"- {dim_names[2]}: {s3}/100"
        )
        good_templates = (
            "what_i_did_good — EXACTLY 3 strings. Insert EXACT words from A1/A2/A3 in every [bracket]:\n"
            f'1. "[Class Presentation] Rule of Three — Structure: You perfectly limited your core takeaways to three points in your opening, saying: \'[EXACT quote showing explicit 3-point framing from A1/A2/A3 — copy verbatim]\'. This kept your audience focused and prevented information overload."\n'
            f'2. "[Class Presentation] Conversational Sense — No Script Reading: Instead of reading from a paper, you spoke naturally and said: \'[EXACT oral-marker quote like You know / Think about it / The thing is from A1/A2/A3 — copy verbatim]\'. This gave your answer a TED-style conversational feel."\n'
            f'3. "[Class Presentation] Illustrative Support — 5-W Example: Instead of giving a dry explanation, you immediately jumped into a specific example when you said: \'[EXACT example quote with 5-W who/when/where detail from A1/A2/A3 — copy verbatim]\'. This made your answer vivid and memorable."\n\n'
            "areas_for_improvement — EXACTLY 3 objects with keys: dimension, issue, example, how_to_fix:\n"
            f'1. dimension="{dim_names[0]}", issue="[Class Presentation] Information Overload: Your answer was unstructured — no Rule of Three framing used.", '
            f'example="You said: \'[EXACT scattered quote without any structure from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Structure into three clear buckets. Say this instead: \'To answer your question, there are three key things: first, [point 1]; second, [point 2]; and third, [point 3].\'"\n'
            f'2. dimension="{dim_names[1]}", issue="[Class Presentation] Script Reading: Your tone sounded like reading from a paper rather than having a conversation.", '
            f'example="You literally read: \'[EXACT mechanical quote with no natural oral connectors from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Speak like a conversation. Say this instead: \'[Same idea rewritten with You know / Think about it / The thing is at the start].\'"\n'
            f'3. dimension="{dim_names[2]}", issue="[Class Presentation] No Specific Example: You gave a conceptual explanation with no real-world 5-W example.", '
            f'example="You said: \'[EXACT abstract explanation quote with no specific details from A1/A2/A3 — copy verbatim]\'", '
            f'how_to_fix="Jump into a 5-W example immediately. Say this instead: \'Let me give you a concrete example: [who] did [what] in [when/where], and the result was [outcome].\'"\n'
        )

    # ── Per-dimension signal labels used in RULE 6 ───────────────────────────
    # Tell the LLM exactly what linguistic signal to look for per dimension,
    # so each example field targets a different type of evidence.
    _quote_signal_labels = {
        "thesis_defense": [
            "direct stance marker — does the FIRST sentence lead with a claim or dodge with context?",
            "hard evidence — a number, percentage, p-value, n=X, or literature citation",
            "tact signal — an acknowledgment phrase (valid concern / fair point) OR a blunt rebuttal opener",
        ],
        "case_pitch": [
            "verdict-first opening — does sentence [1] of any answer lead with Yes/No/We believe?",
            "persuasion signal — a data metric (% / revenue / users) OR a customer story / emotional hook",
            "confidence signal — a hedge word (maybe/perhaps/I think/I guess) OR an assertive statement",
        ],
        "class_presentation": [
            "list structure — first / second / third ordinal markers, or total absence of any structure",
            "oral tone — natural connector (you know / think about it / the thing is) OR scripted/formal phrasing",
            "concrete example — for example / specifically / in [year] / 5-W who+when+where detail, or total abstraction",
        ],
    }

    # ── Compose Gemini prompt ─────────────────────────────────────────────────
    w0p = int(weights[0] * 100)
    w1p = int(weights[1] * 100)
    w2p = int(weights[2] * 100)

    # ── 3-way cross-analysis context: Strategies × User Answers × Slide Content ──
    _slide_content_text = ""
    if slides:
        _slide_content_text = "\n".join(
            f"  Slide {s.get('page','?')}: {s.get('title','')} — {s.get('content','')[:200]}"
            for s in (slides or [])
        )

    _strategy_analysis_text = ""
    for _i, _ex in enumerate(fe_transcripts[:5], 1):
        _strat = _ex.get("answering_strategy", "")
        if _strat:
            _strategy_analysis_text += (
                f"\n  Q{_i} Recommended Strategy: {_strat}\n"
                f"  Q{_i} User Answer: {_ex.get('answer','')[:300]}\n"
            )

    _three_way_section = ""
    if _strategy_analysis_text:
        _three_way_section = (
            "══════════════════════════════════════════════\n"
            "3-WAY CROSS-ANALYSIS (use this for coaching)\n"
            "══════════════════════════════════════════════\n"
            "For each answer, compare: (1) Recommended Strategy vs (2) User's Actual Answer vs "
            "(3) Slide Evidence below.\n"
            "Did the user follow the strategy structure? Did they cite slide data?\n"
            f"{_strategy_analysis_text}\n"
        )
    if _slide_content_text:
        _three_way_section += (
            "── Defence Slide Content (check if user cited this evidence) ──\n"
            f"{_slide_content_text}\n\n"
        )

    cq_prompt = (
        f"You are an expert Communication Quality (CQ) coach for {scene_label} scenarios. "
        f"Evaluate ONLY the Q&A exchange transcripts below — NOT the presentation narration.\n\n"
        f"Audience: {audience} | Difficulty: {difficulty} | Scene: {scene_label}\n\n"
        "══════════════════════════════════════════════\n"
        f"{rubric_text}\n\n"
        f"{_three_way_section}"
        "══════════════════════════════════════════════\n"
        "Q&A EXCHANGE TRANSCRIPTS (evaluate ONLY these)\n"
        "══════════════════════════════════════════════\n"
        f"{exchanges_text}\n\n"
        "══════════════════════════════════════════════\n"
        "OUTPUT RULES\n"
        "══════════════════════════════════════════════\n"
        "RULE 0 (LANGUAGE — MANDATORY): ALL text in what_i_did_good, areas_for_improvement, how_to_fix, "
        "and every field inside communication_quality_report MUST be written in PURE ENGLISH ONLY. "
        "Do NOT write any Chinese characters anywhere in the output — not even as translation aids.\n"
        "RULE 1: IELTS 5.5-6.0 vocabulary. Short, clear sentences for non-native speakers.\n"
        "RULE 2 (MOST IMPORTANT): Every item MUST quote EXACT words from A1/A2/A3 above. NEVER invent quotes.\n"
        "RULE 3: EXACTLY 3 items each in what_i_did_good and areas_for_improvement.\n"
        f"RULE 4 — TEMPLATES:\n{good_templates}\n"
        "RULE 5: Use heuristic anchors as starting points. Adjust ±15 max based on actual text.\n"
        f"RULE 6 (QUOTE ISOLATION — MANDATORY): Each dimension's \"example\" field MUST use a DIFFERENT sentence fragment.\n"
        "  Answers are pre-labelled [1] [2] [3]… — use these labels to pick the single most relevant sentence per dimension.\n"
        "  Do NOT include the [n] labels in your JSON output — copy only the raw sentence text.\n"
        "  Do NOT paste the same sentence into two different example fields.\n"
        f"  ▸ {dim_names[0]} example → find the sentence that best reveals [{_quote_signal_labels[scene_slug][0]}]\n"
        f"  ▸ {dim_names[1]} example → find the sentence that best reveals [{_quote_signal_labels[scene_slug][1]}]\n"
        f"  ▸ {dim_names[2]} example → find the sentence that best reveals [{_quote_signal_labels[scene_slug][2]}]\n"
        "  If a dimension's signal is completely absent, quote the nearest relevant fragment and note the absence.\n\n"
        "RULE 7 (PER-QUESTION TIMELINE — 3-AXIS FUSION — MANDATORY): In ADDITION to the dimension-level "
        "fields above, produce a communication_quality_report object with ONE entry per Q&A exchange "
        "listed in the transcripts (Q1, Q2, Q3…). For EACH exchange you MUST silently perform a 3-axis "
        "fusion BEFORE writing any field:\n"
        "  AXIS 1 — the REAL slide content relevant to this question (see slide list above / defence slide "
        "content section). Find the specific slide whose title/content matches this question's topic.\n"
        "  AXIS 2 — internally infer the ideal, complete answer this exact question deserves, built ONLY "
        "from concrete facts/numbers/named concepts found in that slide's real content (do not output this "
        "inferred answer as its own field — use it only to judge the gap in AXIS 3).\n"
        "  AXIS 3 — the user's ACTUAL answer text for this question, exactly as given, however short.\n"
        "  Now write the three output fields, GROUNDED in this fusion — never generic:\n"
        "  - question_id: copy the exact [id=...] value shown next to that Qn.\n"
        "  - question_text: copy the question text verbatim.\n"
        "  - user_actual_answer: copy the user's full answer verbatim (remove the [n] sentence labels).\n"
        "  - what_i_did_good: 1-2 sentences, pure English, IELTS 5.5-6.0. You MUST quote the user's COMPLETE "
        "AXIS 3 answer word-for-word in full — the ENTIRE text of user_actual_answer for THIS question only, "
        "never a truncated fragment, never a partial excerpt, and never words borrowed from a different "
        "question's answer. If the answer is only one or two words, quote those one or two words in full and "
        "praise that they captured the right keyword, noting it needs expansion — do not invent a longer "
        "quote or add words the user did not say.\n"
        "  - areas_for_improvement: 1-2 sentences, pure English, IELTS 5.5-6.0. Name the SPECIFIC fact, "
        "number, named concept, or term from the AXIS 1 slide content that the user's AXIS 3 answer left "
        "out. NEVER write generic advice such as 'add more evidence' or 'be clearer' — always name the "
        "exact missing slide detail. If no matching slide exists, name the specific gap in the user's own "
        "answer instead (still concrete, never generic).\n"
        "  - how_to_fix: MUST start with the exact phrase 'Say this instead:' followed by a FULL, "
        "ready-to-recite paragraph of 2-4 complete sentences that the student could literally read aloud. "
        "Write a CLEAN, OPTIMIZED answer built from the specific AXIS 1 slide vocabulary/data — in flat "
        "IELTS 5.5-6.0 sentence structure. Do NOT quote or reuse the user's own AXIS 3 wording here (their "
        "original words may be garbled, incomplete, or an incorrect answer); write a fresh, correct, "
        "well-formed answer as if speaking it for the first time. NEVER give generic methodology advice "
        "like 'state your conclusion first, then give evidence' — always write the actual answer content "
        "itself.\n"
        "  Every field must be pure English — no Chinese.\n\n"
        "Return ONLY valid JSON. No markdown fences. No text outside JSON.\n\n"
        "{\n"
        f'  "cq_scores": {{"{dim_names[0]}": <int 0-100>, "{dim_names[1]}": <int 0-100>, "{dim_names[2]}": <int 0-100>}},\n'
        f'  "cq_total": <int 0-100, weighted {dim_names[0]}x{w0p}% + {dim_names[1]}x{w1p}% + {dim_names[2]}x{w2p}%>,\n'
        '  "what_i_did_good": [\n'
        f'    "[{scene_label}] {dim_names[0]}: ...",\n'
        f'    "[{scene_label}] {dim_names[1]}: ...",\n'
        f'    "[{scene_label}] {dim_names[2]}: ..."\n'
        '  ],\n'
        '  "areas_for_improvement": [\n'
        f'    {{"dimension": "{dim_names[0]}", "issue": "...", "example": "You said: \'...\'", "how_to_fix": "Say this instead: \'...\'"}},\n'
        f'    {{"dimension": "{dim_names[1]}", "issue": "...", "example": "You said: \'...\'", "how_to_fix": "Say this instead: \'...\'"}},\n'
        f'    {{"dimension": "{dim_names[2]}", "issue": "...", "example": "You said: \'...\'", "how_to_fix": "Say this instead: \'...\'"}}\n'
        '  ],\n'
        '  "communication_quality_report": {\n'
        '    "overall_cq_score": <int 0-100>,\n'
        '    "per_question_analysis": [\n'
        '      {"question_id": "...", "question_text": "...", "user_actual_answer": "...", '
        '"what_i_did_good": "...", "areas_for_improvement": "...", "how_to_fix": "Say this instead: ..."}\n'
        '    ]\n'
        '  }\n'
        '}'
    )

    app.logger.info(
        f"[CQ EVAL] scene={scene_slug} | qa_exchanges={exchange_count} | "
        f"qa_words={sum(len(t.split()) for t in qa_texts)} | anchors={heuristic_scores}"
    )

    try:
        response = _ai_client.chat.completions.create(
            model=EVAL_MODEL,
            max_tokens=4096,
            messages=[{"role": "user", "content": cq_prompt}],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$",          "", raw)
        result = json.loads(raw)

        # Normalize scores
        raw_cq = result.get("cq_scores", {})
        result["cq_scores"] = {k: int(min(100, max(0, float(v)))) for k, v in raw_cq.items()}
        if not result["cq_scores"]:
            result["cq_scores"] = heuristic_scores

        # Recompute CQ total from normalised scores
        vals     = [result["cq_scores"].get(d, 60) for d in dim_names]
        ai_total = result.get("cq_total", 0)
        result["cq_total"] = (
            int(ai_total) if 0 < ai_total <= 100
            else int(round(sum(v * w for v, w in zip(vals, weights))))
        )

        if not result.get("what_i_did_good"):
            result["what_i_did_good"] = [f"[{scene_label}] Communication recorded."]
        if not result.get("areas_for_improvement"):
            result["areas_for_improvement"] = []

        # Normalize communication_quality_report — fall back to a local
        # English per-question analysis if the LLM omitted/mangled it.
        cqr = result.get("communication_quality_report") or {}
        pqa = cqr.get("per_question_analysis")
        if not pqa or not isinstance(pqa, list):
            pqa = _fallback_per_question_analysis(comm_transcripts, slides=slides, scene_slug=scene_slug)
        pqa = _dedupe_pqa(pqa)
        result["communication_quality_report"] = {
            "overall_cq_score":     int(cqr.get("overall_cq_score", result["cq_total"]) or result["cq_total"]),
            "per_question_analysis": pqa,
        }

        result["has_data"]       = True
        result["scene_slug"]     = scene_slug
        result["scene_label"]    = scene_label
        result["dim_names"]      = dim_names
        result["weights"]        = weights
        result["exchange_count"] = exchange_count
        return result

    except json.JSONDecodeError as e:
        app.logger.error(f"[CQ EVAL ERROR] JSONDecodeError: {str(e)[:200]} — attempting JSON repair")
        # ── JSON repair: Gemini truncated the response mid-string ─────────────
        try:
            repaired = _repair_truncated_json(raw)
            if repaired and isinstance(repaired, dict):
                app.logger.info("[CQ EVAL] JSON repair succeeded — using recovered partial result")
                # Normalise whatever scores were recovered
                raw_cq = repaired.get("cq_scores", {})
                repaired["cq_scores"] = {
                    k: int(min(100, max(0, float(v)))) for k, v in raw_cq.items()
                } if raw_cq else heuristic_scores

                # If coaching text arrays were truncated away, generate them locally
                # with the user's ACTUAL qa_texts injected as verbatim quotes.
                ai_scores_recovered = repaired["cq_scores"] if repaired["cq_scores"] else None
                if not repaired.get("what_i_did_good") or not repaired.get("areas_for_improvement"):
                    good_items, fix_items = _build_cq_coaching_cards(
                        qa_texts,
                        scene_slug, scene_label, dim_names,
                        ai_scores_recovered or heuristic_scores,
                    )
                    repaired.setdefault("what_i_did_good",       good_items)
                    repaired.setdefault("areas_for_improvement",  fix_items)

                vals     = [repaired["cq_scores"].get(d, 60) for d in dim_names]
                ai_total = repaired.get("cq_total", 0)
                repaired["cq_total"] = (
                    int(ai_total) if 0 < int(ai_total) <= 100
                    else int(round(sum(v * w for v, w in zip(vals, weights))))
                )
                cqr = repaired.get("communication_quality_report") or {}
                pqa = cqr.get("per_question_analysis")
                if not pqa or not isinstance(pqa, list):
                    pqa = _fallback_per_question_analysis(comm_transcripts, slides=slides, scene_slug=scene_slug)
                pqa = _dedupe_pqa(pqa)
                repaired["communication_quality_report"] = {
                    "overall_cq_score":      int(cqr.get("overall_cq_score", repaired["cq_total"]) or repaired["cq_total"]),
                    "per_question_analysis": pqa,
                }

                repaired["has_data"]       = True
                repaired["scene_slug"]     = scene_slug
                repaired["scene_label"]    = scene_label
                repaired["dim_names"]      = dim_names
                repaired["weights"]        = weights
                repaired["exchange_count"] = exchange_count
                return repaired
        except Exception as repair_err:
            app.logger.error(f"[CQ EVAL] JSON repair also failed: {repair_err}")

        # All repair attempts exhausted — build coaching cards with actual user words
        good_items, fix_items = _build_cq_coaching_cards(
            qa_texts, scene_slug, scene_label, dim_names, heuristic_scores
        )
        return _cq_mock_result(
            scene_slug, heuristic_scores, cq_total_heuristic, dim_names, exchange_count,
            good_override=good_items, fix_override=fix_items,
            per_question_override=_fallback_per_question_analysis(comm_transcripts, slides=slides, scene_slug=scene_slug),
        )

    except Exception as e:
        app.logger.error(f"[CQ EVAL ERROR] {type(e).__name__}: {str(e)[:300]}")
        traceback.print_exc()
        good_items, fix_items = _build_cq_coaching_cards(
            qa_texts, scene_slug, scene_label, dim_names, heuristic_scores
        )
        return _cq_mock_result(
            scene_slug, heuristic_scores, cq_total_heuristic, dim_names, exchange_count,
            good_override=good_items, fix_override=fix_items,
            per_question_override=_fallback_per_question_analysis(comm_transcripts, slides=slides, scene_slug=scene_slug),
        )


def _build_dual_track_mock_result(universal_heuristic, anchor_passed, anchor_score,
                                    anchor_type, anchor_text, scene_slug, scene_label,
                                    universal_total, exchange_count, free_texts, target_dim,
                                    per_question_override=None):
    """Fallback dual-track CQ result when Gemini is unavailable."""
    u_scores = dict(universal_heuristic)
    u_scores[f"Anchor — {anchor_type}"] = anchor_score
    combined_total = int(round(universal_total * 0.75 + anchor_score * 0.25))
    dim_names  = list(universal_heuristic.keys()) + [f"Anchor — {anchor_type}"]
    best_u_dim = max(universal_heuristic, key=universal_heuristic.get)
    mid_u_dim  = sorted(universal_heuristic, key=universal_heuristic.get)[1]

    good_items = [
        f"[Universal] {best_u_dim} ({universal_heuristic[best_u_dim]}/100): "
        f"Your response showed strength here. Keep applying this approach.",
        f"[Universal] {mid_u_dim} ({universal_heuristic[mid_u_dim]}/100): "
        f"Good baseline. Focused practice will push this above 75.",
    ]
    if anchor_passed:
        good_items.append(
            f"【ANCHOR PASS】[{scene_label}] {anchor_type} ({anchor_score}/100): "
            f"You followed the scaffold hint and demonstrated the target skill. Excellent!"
        )
    else:
        good_items.append(
            f"[{scene_label}] {anchor_type} ({anchor_score}/100): "
            f"You attempted the anchor question. Focus on following the scaffold hint exactly next time."
        )

    worst_u_dim = min(universal_heuristic, key=universal_heuristic.get)
    fix_items = [{
        "dimension": worst_u_dim,
        "issue":     f"[Universal] {worst_u_dim} needs improvement.",
        "example":   f'You said: "{free_texts[0][:120]}…"' if free_texts else "See transcript.",
        "how_to_fix": _UNIVERSAL_FIX_PHRASES.get(worst_u_dim, "Practice focused Q&A coaching."),
    }]
    if not anchor_passed:
        fix_items.append({
            "dimension": anchor_type,
            "issue":     (
                f"[{scene_label}] The scaffold hint was provided but the {anchor_type} technique "
                "was not fully demonstrated."
            ),
            "example":   f'You said: "{anchor_text[:120]}…"' if anchor_text else "No answer recorded.",
            "how_to_fix": (
                f"Say this instead: Demonstrate {anchor_type} in your very first sentence — "
                "follow the scaffold exactly and practice until it is automatic."
            ),
        })
    return {
        "has_data":              True,
        "scene_slug":            scene_slug,
        "scene_label":           scene_label,
        "cq_total":              combined_total,
        "cq_scores":             u_scores,
        "dim_names":             dim_names,
        "weights":               [],
        "what_i_did_good":       good_items[:3],
        "areas_for_improvement": fix_items[:3],
        "exchange_count":        exchange_count,
        "dual_track":            True,
        "anchor_passed":         anchor_passed,
        "anchor_score":          anchor_score,
        "communication_quality_report": {
            "overall_cq_score":      combined_total,
            "per_question_analysis": per_question_override if per_question_override is not None else [],
        },
    }


def _run_dual_track_cq_evaluation(free_transcripts, anchor_transcripts, scene_slug,
                                    scene_label, audience, difficulty, exchange_count,
                                    slides=None):
    """
    Dual-track CQ evaluation engine (Module 3).

    Section A — Free questions → 3 universal dims (Directness / Conv Resonance / Evidence).
    Section B — Anchor question → scaffold compliance + dynamic 'Say this instead' rewrite.

    CQ scores: 3 universal dims + 1 anchor dim (4 bars in report chart).
    Combined total: universal × 75% + anchor × 25%.
    """
    free_texts  = [t["answer"].strip() for t in free_transcripts if t["answer"].strip()]
    anchor_info = anchor_transcripts[0]   # exactly ONE anchor per session
    anchor_text = anchor_info.get("answer", "").strip()

    d_score, cr_score, es_score = _cq_heuristic_universal(free_texts) if free_texts else (55, 55, 55)
    universal_heuristic = {
        "Directness & Logic":        int(d_score),
        "Conversational Resonance":  int(cr_score),
        "Evidence & Substantiation": int(es_score),
    }
    universal_total_h = int(round(d_score * 0.40 + cr_score * 0.30 + es_score * 0.30))

    anchor_type     = anchor_info.get("anchor_type", "Anchor")
    scaffold_signal = anchor_info.get("scaffold_signal", "")
    target_dim      = anchor_info.get("target_dim", "")
    anchor_q        = anchor_info.get("question", "")
    anchor_passed_h, anchor_score_h = _check_anchor_compliance(anchor_text, target_dim)

    app.logger.info(
        f"[CQ DUAL-TRACK] scene={scene_slug} | free={len(free_transcripts)} | "
        f"anchor={anchor_type} | heuristic={'PASS' if anchor_passed_h else 'FAIL'}({anchor_score_h})"
    )

    if not AI_ENABLED:
        return _build_dual_track_mock_result(
            universal_heuristic, anchor_passed_h, anchor_score_h,
            anchor_type, anchor_text, scene_slug, scene_label,
            universal_total_h, exchange_count, free_texts, target_dim,
            per_question_override=_fallback_per_question_analysis(free_transcripts + anchor_transcripts, slides=slides, scene_slug=scene_slug),
        )

    def _label_sentences(text):
        parts   = re.split(r'(?<=[.!?])\s+', text.strip())
        labeled = [f"[{j+1}] {p.strip()}" for j, p in enumerate(parts) if p.strip()]
        return " ".join(labeled) if labeled else text

    free_exchanges_text = "\n\n".join(
        f"Q{i+1} [free]: {t['question']}\nA{i+1}: {_label_sentences(t['answer'])}"
        for i, t in enumerate(free_transcripts) if t["answer"].strip()
    ) or "(No free Q&A answers recorded)"

    anchor_exchange_text = (
        f"Q_anchor [{anchor_type}]: {anchor_q}\n"
        f"A_anchor: {_label_sentences(anchor_text)}"
    ) if anchor_text else "(No anchor answer recorded)"

    _dual_slide_content_text = ""
    if slides:
        _dual_slide_content_text = "\n".join(
            f"  Slide {s.get('page', '?')}: {s.get('title', '')} — {s.get('content', '')[:200]}"
            for s in slides
        )
    _dual_slide_section = (
        (
            "══════════════════════════════════════════════\n"
            "PRESENTATION SLIDE CONTENT (real data — use for per-question fusion below)\n"
            "══════════════════════════════════════════════\n"
            f"{_dual_slide_content_text}\n\n"
        ) if _dual_slide_content_text else ""
    )

    # Pre-compute strings containing quotes/backslashes (Python 3.11 f-string restriction)
    _scene_example_map = {
        "class_presentation": (
            "CLASS PRES example — Say this instead: "
            '"Think about it this way — when a taxi driver navigates the city, '
            "they are not just driving, they are scanning live patterns. "
            'That is exactly how our framework operates." '
            "(weave user domain into conversational analogy)"
        ),
        "thesis_defense": (
            "THESIS example — Say this instead: "
            '"That is a critical challenge on sample size. To address this directly: '
            "while the cohort is small, our post-hoc power analysis gives 0.85, "
            'confirming the validity of this correlation." '
            "(concede first, then hard academic counter)"
        ),
        "case_pitch": (
            "PITCH example — Say this instead: "
            '"Yes, this model is highly replicable — our unit economics rely on a '
            "40 percent organic referral rate, driven by two key operational factors: "
            'first..." '
            "(MECE conclusion-first, McKinsey style)"
        ),
    }
    _scene_example   = _scene_example_map.get(scene_slug, "")
    _anchor_pass_ln  = (
        f"    If passed: [ANCHOR PASS][Metric: {anchor_type}] {anchor_type} "
        "(N/100): high praise. You said: exact A_anchor sentence."
    )
    _anchor_fail_ln  = (
        f"    If not passed: [Metric: {anchor_type}] {anchor_type} "
        "(N/100): You attempted this. Follow the scaffold hint more closely next time."
    )
    _good_item3_tmpl = (
        f'    "[Metric: {anchor_type}] {anchor_type} (N/100): anchor feedback. '
        'You said: exact A_anchor sentence"'
    )
    _fix_anchor_tmpl = (
        f'{{"dimension": "{anchor_type}", "issue": "...", '
        '"example": "You said: exact A_anchor sentence", '
        '"how_to_fix": "Say this instead: dynamic rewrite using A_anchor domain content"}}'
    )

    cq_dual_prompt = (
        "You are a Communication Quality (CQ) coach evaluating two types of Q&A exchanges.\n\n"
        "══════════════════════════════════════════════\n"
        "SECTION A — FREE QUESTION (Universal Dimensions)\n"
        "══════════════════════════════════════════════\n"
        "DIM 1 — Directness & Logic (40%)\n"
        "90-100: Core stance in sentence 1-2. Strong openers: Yes/No/I believe/The key is.\n"
        "70-89:  Mostly direct, one filler before the real answer.\n"
        "50-69:  2+ background sentences before addressing the question.\n"
        "0-49:   Question Dodging — opens with compliments/thanks/re-stating.\n"
        "PENALTY: 'That's a great question' or rephrasing = -15 pts.\n\n"
        "DIM 2 — Conversational Resonance (30%)\n"
        "90-100: Live dialogue tone. Oral connectors: you know / think about it / the thing is.\n"
        "70-89:  Some conversational elements.\n  50-69: Mixed.\n"
        "0-49:   Reads like a formal essay: furthermore / therefore / in conclusion.\n\n"
        "DIM 3 — Evidence & Substantiation (30%)\n"
        "90-100: At least one concrete fact, data point, named case, year, or analogy.\n"
        "70-89:  General example.  50-69: Topic mentioned, no specifics.\n"
        "0-49:   Pure subjective (very/I feel) — zero evidence.\n\n"
        f"ANCHORS (adjust ±15): Directness={d_score} | Conv={cr_score} | Evidence={es_score}\n\n"
        f"FREE EXCHANGES:\n{free_exchanges_text}\n\n"
        "══════════════════════════════════════════════\n"
        "SECTION B — ANCHOR QUESTION (Scaffold Compliance)\n"
        "══════════════════════════════════════════════\n"
        "The anchor question included a 引导 (scaffold hint) in Chinese telling the user EXACTLY what to do.\n"
        f"Scene: {scene_label} | Target: {anchor_type}\n"
        f"Check: {scaffold_signal}\n"
        f"Heuristic estimate: {'PASSED' if anchor_passed_h else 'FAILED'} ({anchor_score_h}/100)\n\n"
        f"ANCHOR EXCHANGE:\n{anchor_exchange_text}\n\n"
        f"{_dual_slide_section}"
        "══════════════════════════════════════════════\n"
        "OUTPUT RULES\n"
        "══════════════════════════════════════════════\n"
        "RULE 0a (JSON SAFETY): Return ONLY valid JSON. Use DOUBLE QUOTES for ALL strings. NEVER single quotes.\n"
        "RULE 0b (LANGUAGE — MANDATORY): ALL text in what_i_did_good, areas_for_improvement, how_to_fix, and "
        "every field inside communication_quality_report MUST be written in PURE ENGLISH ONLY. Do NOT write "
        "any Chinese characters anywhere in the output.\n"
        "RULE 1: IELTS 5.5-6.0 vocabulary. Short sentences.\n"
        "RULE 2 (CRITICAL): Every example MUST quote EXACT words from the labelled transcripts. NEVER invent.\n"
        "  - FREE dims (Directness/Resonance/Evidence): quote ONLY from A1 sentences. NOT from A_anchor.\n"
        "  - ANCHOR dim: quote ONLY from A_anchor sentences. NOT from A1.\n"
        "RULE 3: what_i_did_good = EXACTLY 3 strings:\n"
        "  Item 1: [Metric: Universal] Directness & Logic (N/100): brief praise. You said: exact sentence [1] or [2] from A1.\n"
        "  Item 2: [Metric: Universal] Conversational Resonance (N/100): brief praise. You said: exact different sentence from A1.\n"
        f"  Item 3 (anchor):\n{_anchor_pass_ln}\n{_anchor_fail_ln}\n"
        "RULE 4: areas_for_improvement = 1-3 objects. QUOTE ISOLATION — each example uses a DIFFERENT sentence:\n"
        "  - Universal dim items quote ONLY from A1. Anchor dim items quote ONLY from A_anchor.\n"
        "  - If anchor_passed=false, include anchor item LAST. Its how_to_fix MUST begin with Say this instead:\n"
        "    followed by a DYNAMIC 2-3 sentence rewrite using the user ACTUAL vocabulary and domain content\n"
        "    from A_anchor (their topic, claims, specific words). NEVER write generic placeholder sentences.\n"
        f"    {_scene_example}\n"
        "RULE 5 (PER-QUESTION TIMELINE — 3-AXIS FUSION — MANDATORY): ALSO produce a communication_quality_report "
        "object with ONE entry per exchange in BOTH the FREE EXCHANGES and the ANCHOR EXCHANGE. For EACH "
        "exchange, silently perform a 3-axis fusion BEFORE writing any field:\n"
        "  AXIS 1 — the REAL slide content above whose title/content matches this question's topic.\n"
        "  AXIS 2 — internally infer the ideal, complete answer this question deserves, built ONLY from "
        "concrete facts/numbers/named concepts in that slide (do not output this as its own field).\n"
        "  AXIS 3 — the user's ACTUAL answer text for this exchange, exactly as given, however short.\n"
        "  Then write, GROUNDED in this fusion — never generic:\n"
        "  - question_id (copy the [id=...] if shown, else use 'anchor' for the anchor exchange)\n"
        "  - question_text (verbatim)\n"
        "  - user_actual_answer (verbatim, no [n] labels)\n"
        "  - what_i_did_good: 1-2 sentences, pure English, IELTS 5.5-6.0. Quote the user's COMPLETE AXIS 3 "
        "answer word-for-word in full — the ENTIRE text of user_actual_answer for THIS exchange only, never "
        "a truncated fragment, never a partial excerpt, and never words borrowed from a different exchange "
        "(free vs anchor). If the answer is only one or two words, quote those one or two words in full and "
        "praise the correct keyword, noting it needs expanding — never invent a longer quote or add words "
        "the user did not say.\n"
        "  - areas_for_improvement: 1-2 sentences, pure English, IELTS 5.5-6.0. Name the SPECIFIC fact, "
        "number, or named concept from the AXIS 1 slide that the AXIS 3 answer left out. NEVER write "
        "generic advice like 'add more evidence' — always name the exact missing slide detail. If no "
        "matching slide exists, name the specific gap in the user's own answer instead.\n"
        "  - how_to_fix: MUST start with 'Say this instead:' followed by a FULL, ready-to-recite paragraph "
        "of 2-4 complete sentences the student could literally read aloud. Write a CLEAN, OPTIMIZED answer "
        "built from the specific AXIS 1 slide vocabulary/data, in flat IELTS 5.5-6.0 sentence structure. Do "
        "NOT quote or reuse the user's own AXIS 3 wording here (their original words may be garbled, "
        "incomplete, or an incorrect answer); write a fresh, correct, well-formed answer as if speaking it "
        "for the first time. NEVER give generic methodology advice (e.g. 'state your conclusion first') — "
        "write the actual answer content itself.\n"
        "No Chinese anywhere.\n"
        "Return ONLY valid JSON. No markdown. No text outside JSON. Double quotes only:\n"
        "{\n"
        '  "universal_scores": {"Directness & Logic": int, "Conversational Resonance": int, "Evidence & Substantiation": int},\n'
        '  "universal_total": int,\n'
        '  "anchor_passed": bool,\n'
        '  "anchor_score": int,\n'
        '  "what_i_did_good": [\n'
        '    "[Metric: Universal] Directness & Logic (N/100): praise. You said: exact A1 sentence.",\n'
        '    "[Metric: Universal] Conversational Resonance (N/100): praise. You said: exact different A1 sentence.",\n'
        f'    {_good_item3_tmpl}\n'
        '  ],\n'
        '  "areas_for_improvement": [\n'
        '    {"dimension": "Evidence & Substantiation", "issue": "...", "example": "You said: exact A1 sentence.", "how_to_fix": "..."},\n'
        f'    {_fix_anchor_tmpl}\n'
        '  ],\n'
        '  "communication_quality_report": {\n'
        '    "overall_cq_score": <int 0-100>,\n'
        '    "per_question_analysis": [\n'
        '      {"question_id": "...", "question_text": "...", "user_actual_answer": "...", '
        '"what_i_did_good": "...", "areas_for_improvement": "...", "how_to_fix": "Say this instead: ..."}\n'
        '    ]\n'
        '  }\n'
        '}'
    )

    _DIM_NAMES_BASE = ["Directness & Logic", "Conversational Resonance", "Evidence & Substantiation"]

    try:
        response = _ai_client.chat.completions.create(
            model=EVAL_MODEL, max_tokens=2048,
            messages=[{"role": "user", "content": cq_dual_prompt}],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$",          "", raw)
        # Pre-process: fix single-quoted string values (Gemini outputs Python-style strings)
        try:
            result = json.loads(raw)
        except json.JSONDecodeError:
            _fixed = re.sub(
                r"(?<=[:\[,])\s*'([^'\n]{0,400})'",
                lambda m: ' "' + m.group(1).replace('"', '\\"') + '"',
                raw
            )
            result = json.loads(_fixed)   # re-raise if still invalid → caught below

        u_scores = result.get("universal_scores", universal_heuristic)
        u_scores = {k: int(min(100, max(0, float(v)))) for k, v in u_scores.items()} if u_scores else dict(universal_heuristic)

        u_total      = result.get("universal_total", universal_total_h)
        u_total      = int(u_total) if 0 < int(u_total) <= 100 else universal_total_h
        anchor_passed = result.get("anchor_passed", anchor_passed_h)
        anchor_score  = int(result.get("anchor_score", anchor_score_h))

        u_scores[f"Anchor — {anchor_type}"] = anchor_score
        combined_total = int(round(u_total * 0.75 + anchor_score * 0.25))
        dim_names = _DIM_NAMES_BASE + [f"Anchor — {anchor_type}"]

        good_items = result.get("what_i_did_good") or [f"[Universal] Communication recorded."]
        fix_items  = result.get("areas_for_improvement") or []

        cqr = result.get("communication_quality_report") or {}
        pqa = cqr.get("per_question_analysis")
        if not pqa or not isinstance(pqa, list):
            pqa = _fallback_per_question_analysis(free_transcripts + anchor_transcripts, slides=slides, scene_slug=scene_slug)
        pqa = _dedupe_pqa(pqa)

        return {
            "has_data":              True,
            "scene_slug":            scene_slug,
            "scene_label":           scene_label,
            "cq_total":              combined_total,
            "cq_scores":             u_scores,
            "dim_names":             dim_names,
            "weights":               [0.30, 0.225, 0.225, 0.25],
            "what_i_did_good":       good_items,
            "areas_for_improvement": fix_items,
            "exchange_count":        exchange_count,
            "dual_track":            True,
            "anchor_passed":         anchor_passed,
            "anchor_score":          anchor_score,
            "communication_quality_report": {
                "overall_cq_score":      int(cqr.get("overall_cq_score", combined_total) or combined_total),
                "per_question_analysis": pqa,
            },
        }

    except json.JSONDecodeError as e:
        app.logger.error(f"[CQ DUAL-TRACK] JSONDecodeError: {str(e)[:200]} — attempting repair")
        try:
            repaired = _repair_truncated_json(raw)
            if repaired and isinstance(repaired, dict):
                u_sc = repaired.get("universal_scores", universal_heuristic)
                u_sc = {k: int(min(100, max(0, float(v)))) for k, v in u_sc.items()} if u_sc else dict(universal_heuristic)
                a_p  = repaired.get("anchor_passed", anchor_passed_h)
                a_s  = int(repaired.get("anchor_score", anchor_score_h))
                u_sc[f"Anchor — {anchor_type}"] = a_s
                u_t  = int(repaired.get("universal_total", universal_total_h))
                comb = int(round(u_t * 0.75 + a_s * 0.25)) if 0 < u_t <= 100 else int(round(universal_total_h * 0.75 + a_s * 0.25))
                good = repaired.get("what_i_did_good") or []
                fix  = repaired.get("areas_for_improvement") or []
                # Coaching arrays may be truncated — fill from mock fallback
                if not good or not fix:
                    _fb = _build_dual_track_mock_result(
                        universal_heuristic, a_p, a_s,
                        anchor_type, anchor_text, scene_slug, scene_label,
                        universal_total_h, exchange_count, free_texts, target_dim
                    )
                    good = good or _fb["what_i_did_good"]
                    fix  = fix  or _fb["areas_for_improvement"]
                cqr = repaired.get("communication_quality_report") or {}
                pqa = cqr.get("per_question_analysis")
                if not pqa or not isinstance(pqa, list):
                    pqa = _fallback_per_question_analysis(free_transcripts + anchor_transcripts, slides=slides, scene_slug=scene_slug)
                pqa = _dedupe_pqa(pqa)
                return {
                    "has_data": True, "scene_slug": scene_slug, "scene_label": scene_label,
                    "cq_total": comb, "cq_scores": u_sc,
                    "dim_names": _DIM_NAMES_BASE + [f"Anchor — {anchor_type}"],
                    "weights": [], "what_i_did_good": good,
                    "areas_for_improvement": fix,
                    "exchange_count": exchange_count, "dual_track": True,
                    "anchor_passed": a_p, "anchor_score": a_s,
                    "communication_quality_report": {
                        "overall_cq_score":      int(cqr.get("overall_cq_score", comb) or comb),
                        "per_question_analysis": pqa,
                    },
                }
        except Exception:
            pass
        return _build_dual_track_mock_result(
            universal_heuristic, anchor_passed_h, anchor_score_h,
            anchor_type, anchor_text, scene_slug, scene_label,
            universal_total_h, exchange_count, free_texts, target_dim,
            per_question_override=_fallback_per_question_analysis(free_transcripts + anchor_transcripts, slides=slides, scene_slug=scene_slug),
        )

    except Exception as e:
        app.logger.error(f"[CQ DUAL-TRACK] {type(e).__name__}: {str(e)[:300]}")
        traceback.print_exc()
        return _build_dual_track_mock_result(
            universal_heuristic, anchor_passed_h, anchor_score_h,
            anchor_type, anchor_text, scene_slug, scene_label,
            universal_total_h, exchange_count, free_texts, target_dim,
            per_question_override=_fallback_per_question_analysis(free_transcripts + anchor_transcripts, slides=slides, scene_slug=scene_slug),
        )


# ─────────────────────────────────────────────
def _generate_pitch_data(delivery_score, wpm=0):
    """Return 20 simulated pitch values (Hz) for the Delivery chart.
    Higher delivery score → more expressive variance; lower → flatter/monotone."""
    import math
    base = 120 + max(0, wpm - 130) * 0.5 if wpm > 0 else 125
    base = max(90, min(180, float(base)))
    amplitude = max(5.0, (delivery_score - 40) * 1.2)
    data = []
    for i in range(20):
        wave = math.sin(i * 0.6) * amplitude + math.sin(i * 1.3) * amplitude * 0.4
        v = int(base + wave + random.randint(-5, 5))
        data.append(max(50, min(250, v)))
    return data


# GEMINI: 4-PILLAR EVALUATION ENGINE (Step 8)
# ─────────────────────────────────────────────
def run_pillar_evaluation(slides, answers, config, challenge_seed,
                          fe_qa_history=None, total_time_seconds=0):
    """
    4-Pillar precision evaluation via Gemini 1.5 Pro.
    Rubric: Structure / Fluency / Relevance / Delivery (each 0-100).

    New in this version:
    - Accepts fe_qa_history (frontend Q&A chat log) for richer QA context.
    - Accepts total_time_seconds for real WPM calculation.
    - Pre-computes word count, WPM, and filler density from actual transcript text.
    - Detects completely empty presentations and applies mandatory score penalties.
    Falls back to _mock_pillar_evaluation on error or when AI is off.
    """
    audience       = config.get("audience",   "Professor")
    scenario       = config.get("scenario",   "Class Presentation")
    difficulty     = config.get("difficulty", "Medium")
    challenge_type = (challenge_seed or {}).get("challenge_type", "Unknown")

    narration_entries = [a for a in answers if a.get("type") == "narration"]
    qa_entries        = [a for a in answers
                         if a.get("type") in ("qa_answer", "academic_qa")]

    # ── Pre-compute real metrics from the actual transcript text ───────────────
    all_narration_text = " ".join(
        (a.get("text") or "").strip() for a in narration_entries
    ).strip()
    total_words = len(all_narration_text.split()) if all_narration_text else 0

    # Real WPM (requires timer data sent from frontend)
    if total_words > 0 and total_time_seconds > 30:
        wpm_estimate = round(total_words / (total_time_seconds / 60))
    else:
        wpm_estimate = 0

    # Actual filler word scan from transcript
    FILLER_RE = r'\b(um+|uh+|like|you know|basically|sort of|kind of|i mean|literally|right\?)\b'
    filler_matches  = re.findall(FILLER_RE, all_narration_text.lower())
    filler_count    = len(filler_matches)
    filler_density  = round(filler_count / max(total_words, 1) * 100, 1)

    # ── TED Talk Like TED: Rule of Three detection ─────────────────────────────
    # Check opening (~first 120 words) for explicit 3-point framing
    opening_text = " ".join(all_narration_text.split()[:120]).lower()
    RULE_OF_THREE_RE = (
        r'\b(three\s+(?:things|points|key\s+points|main\s+points|areas|reasons|ideas|'
        r'steps|sections|aspects|topics|parts|takeaways|pillars))'
        r'|first[,\.].*second[,\.].*third'
        r'|firstly.*secondly.*thirdly'
        r'|number\s+one.*number\s+two.*number\s+three'
    )
    rule_of_three_hit = bool(re.search(RULE_OF_THREE_RE, opening_text, re.DOTALL))
    rule_of_three_note = (
        "RULE OF THREE DETECTED in opening — award +10 Structure bonus (cap at 100)."
        if rule_of_three_hit else
        "No Rule of Three framing detected in the opening."
    )

    # ── TED Conversational Dialogue vs Script ──────────────────────────────────
    CONV_RE = (
        r"\b(let'?s|imagine|think about|picture this|consider this|what if|"
        r"here'?s (?:the thing|why|what)|in other words|for example|for instance|"
        r"what does this mean|the key takeaway|the bottom line|dive into|"
        r"let me show you|the truth is|in fact|actually|right\?|you see|"
        r"now[,\s]|so[,\s](?:what|here|the|let'?s)|and here'?s|"
        r"think of it this way|the reason is|here'?s the (?:key|point|thing))\b"
    )
    conv_marker_count = len(re.findall(CONV_RE, all_narration_text.lower()))
    if conv_marker_count >= 5:
        conv_note = (
            f"{conv_marker_count} conversational markers found — "
            "speaker is engaging and dialogue-style. Boost Fluency slightly."
        )
    elif conv_marker_count >= 2:
        conv_note = (
            f"{conv_marker_count} conversational markers found — "
            "some dialogue style, but could be more engaging."
        )
    else:
        conv_note = (
            f"Only {conv_marker_count} conversational markers — "
            "speech sounds like script-reading. Penalize Fluency (−5 to −10 pts)."
        )

    # ── TED Picture Superiority Effect (PSE) per slide ────────────────────────
    # Fires when a slide has heavy text AND the user mirrors it word-for-word
    narration_map = {a["page"]: (a.get("text") or "") for a in narration_entries if "page" in a}
    pse_triggered_pages = []
    for s in slides:
        pg = s["page"]
        slide_words = re.findall(r'\b\w{4,}\b', (s.get("content", "") + " " + s.get("title", "")).lower())
        user_words  = re.findall(r'\b\w{4,}\b', narration_map.get(pg, "").lower())
        if len(slide_words) > 25 and user_words:
            overlap = len(set(user_words) & set(slide_words)) / len(set(slide_words))
            if overlap > 0.55:
                pse_triggered_pages.append(pg)
    if pse_triggered_pages:
        pse_note = (
            f"PSE TRIGGERED on slide(s) {pse_triggered_pages}: heavy bullet-point slides "
            "with >55% word-mirror narration. Cap Relevance score at 70 for those slides. "
            "Deduct −10 from overall Relevance and note it in areas_for_improvement."
        )
    else:
        pse_note = "No PSE triggered — user added oral explanation beyond slide text."

    # ── TED Jaw-Dropping Moment heuristic ─────────────────────────────────────
    full_qa_text = " ".join(
        (a.get("text") or "") for a in answers
        if a.get("type") in ("qa_answer", "academic_qa")
    )
    JAW_DROP_RE = (
        r'\b(\d+(?:\.\d+)?%|'
        r'\d[\d,]*\s*(?:million|billion|trillion)|'
        r'every\s+\d+\s+(?:second|minute|hour|day|person|people)|'
        r'one\s+in\s+\d+|'
        r'(?:shocking|remarkable|unprecedented|transformative|revolutionary|groundbreaking)|'
        r'\d{4,}\s*(?:people|students|countries|cases|deaths|lives))\b'
    )
    jaw_drop_heuristic = bool(
        re.search(JAW_DROP_RE, all_narration_text, re.IGNORECASE) or
        re.search(JAW_DROP_RE, full_qa_text, re.IGNORECASE)
    )
    jaw_drop_note = (
        "JAW-DROPPING MOMENT CANDIDATE: user cited a striking stat/data point or "
        "used powerful language. Set jaw_dropping_moment=true if you verify it was "
        "impactful AND accompanied by a pause or emphasis in context."
        if jaw_drop_heuristic else
        "No obvious jaw-dropping stat or story detected. Set jaw_dropping_moment=false "
        "unless the speech content itself contains a compelling case or conclusion."
    )

    # ── BRANCH A: Zero / near-zero speech — fast return, skip Gemini entirely ──
    # This fires when the user never spoke (< 5 real words across all slides).
    # Calling Gemini with an empty transcript produces hallucinated, generic praise
    # that is completely wrong. We return a fixed penalty result immediately.
    if total_words < 5:
        app.logger.warning(
            f"[EvalGuard] Zero-speech detected ({total_words} words). "
            "Skipping Gemini call. Returning penalty result."
        )
        return {
            "scores": {"structure": 0, "fluency": 0, "relevance": 0, "delivery": 0},
            "dimensions_info": {
                "structure": {"explanation": "No speech was recorded.", "calculation": "Score is 0 — no narration to evaluate."},
                "fluency":   {"explanation": "No speech was recorded.", "calculation": "Score is 0 — no narration to evaluate."},
                "relevance": {"explanation": "No speech was recorded.", "calculation": "Score is 0 — no narration to evaluate."},
                "delivery":  {"explanation": "No speech was recorded.", "calculation": "Score is 0 — no narration to evaluate."},
            },
            "filler_log": [],
            "what_i_did_good": [
                "No silver lining found. You did not speak during the presentation."
            ],
            "areas_for_improvement": [
                {
                    "issue": "[Delivery] No speech recorded.",
                    "example": "We found no spoken words in your entire presentation rehearsal.",
                    "how_to_fix": "Please turn on your microphone and speak clearly for each slide before clicking Next Slide. Say this instead: \"Good morning. Today I will talk about [your topic].\"",
                }
            ],
        }

    # ── BRANCH B: Real speech present — proceed to AI evaluation ─────────────

    # ── Build per-slide narration map ──────────────────────────────────────────
    narration_map = {a["page"]: (a.get("text") or "").strip() for a in narration_entries}
    slide_narration_text = "\n\n".join(
        f"[Slide {s['page']} — \"{s.get('title','')}\"]:\n"
        f"{narration_map.get(s['page']) or '(No narration recorded for this slide)'}"
        for s in slides
    )

    # ── Build QA exchange text ─────────────────────────────────────────────────
    if fe_qa_history:
        qa_history_text = "\n".join(
            f"[{'AI Examiner' if h.get('role')=='ai' else 'Presenter'}] "
            f"({h.get('type','')}) {h.get('text','')}"
            for h in fe_qa_history
        )
    else:
        qa_parts = []
        for a in qa_entries:
            q = a.get("question", "")
            t = a.get("text", "")
            qa_parts.append(f"Q: {q}\nA: {t}" if q else f"[Answer]: {t}")
        qa_history_text = "\n\n".join(qa_parts) or "(No Q&A answers recorded)"

    # ── No AI: return generic mock (no specific slide content known) ───────────
    if not AI_ENABLED:
        return _mock_pillar_evaluation(difficulty)

    # ── Compose WPM / filler notes ─────────────────────────────────────────────
    # ── WPM calibration: Talk Like TED conversational golden zone = 160-190 WPM ──
    if wpm_estimate > 0:
        if 160 <= wpm_estimate <= 190:
            wpm_note = (
                f"{wpm_estimate} WPM — TED GOLDEN ZONE (160-190 WPM). "
                "Score Delivery 95-100."
            )
        elif 145 <= wpm_estimate < 160:
            wpm_note = (
                f"{wpm_estimate} WPM — slightly below TED golden (145-159 WPM). "
                "Score Delivery 80-94."
            )
        elif 190 < wpm_estimate <= 215:
            wpm_note = (
                f"{wpm_estimate} WPM — slightly fast (191-215 WPM). "
                "Score Delivery 72-85."
            )
        elif 120 <= wpm_estimate < 145:
            wpm_note = (
                f"{wpm_estimate} WPM — noticeably slow (120-144 WPM). "
                "Score Delivery 60-78."
            )
        elif 110 <= wpm_estimate < 120:
            wpm_note = (
                f"{wpm_estimate} WPM — HESITANT (110-119 WPM). "
                "Score Delivery 50-65."
            )
        elif wpm_estimate > 215:
            wpm_note = (
                f"{wpm_estimate} WPM — RUSHING (>215 WPM). "
                "Score Delivery 45-65."
            )
        else:
            wpm_note = (
                f"{wpm_estimate} WPM — TOO SLOW (<110 WPM). "
                "Score Delivery 35-55."
            )
    else:
        wpm_note = "WPM unknown (no timer data). Estimate from text density and content coverage."

    if filler_count > 10:
        filler_note = (
            f"{filler_count} filler words ({filler_density}% of words). "
            "HIGH density — penalise Fluency significantly (deduct ≥20 pts)."
        )
    elif filler_count > 5:
        filler_note = (
            f"{filler_count} filler words ({filler_density}% of words). "
            "Moderate density — penalise Fluency mildly (deduct 5-15 pts)."
        )
    else:
        filler_note = (
            f"{filler_count} filler words ({filler_density}% of words). "
            "Low density — do NOT penalise Fluency for fillers."
        )

    slide_content_text = "\n\n".join(
        f"[Slide {s['page']}] {s.get('title', '')}\n{s.get('content', '')}"
        for s in slides
    )

    prompt = f"""You are a world-class Presentation Coach inspired by the TED Talk methodology from "Talk Like TED" by Carmine Gallo. Evaluate the presenter using the 4-Pillar rubric below. Use ONLY the real data supplied. Every quote MUST come from the actual narration text.

══════════════════════════════════════════════
TALK LIKE TED — SCORING RUBRIC (0-100 per pillar)
══════════════════════════════════════════════

1. [Structure] — Logic chain, Rule of Three, Message Map
   90-100 (TED Master):   Perfect intro/body/conclusion; Rule of Three framing; Twitter-style headline
   75-89  (Competent):    Logic complete; 1-2 abrupt transitions; no explicit 3-point frame
   60-74  (Developing):   No global frame; reading titles only; <20% transitions used
   0-59   (Novice):       Wrong order or mostly silent

   BONUS RULE — Rule of Three (+10 pts, cap 100):
   If the user explicitly states "I will cover three things / three points / three key areas" in the opening AND delivers on it, award a +10 bonus to Structure score.
   PRE-COMPUTED: {rule_of_three_note}

   PENALTY — Information Overload (−10 pts):
   If the presenter introduces >5 separate themes with no unifying frame, deduct 10 pts.

2. [Fluency] — Filler words, Dialogue vs Script reading, Conversational markers
   90-100 (Expert):     <2 filler words/min; natural conversational flow; uses "Let's", "Imagine", "Now"
   75-89  (Competent):  3-5 filler words/min; some connectors; not robotic
   60-74  (Developing): Sounds like reading a paper; zero conversational markers; OR filler >3/min
   0-59   (Novice):     Broken sentences; fillers make meaning unclear

   DIALOGUE ASSESSMENT:
   PRE-COMPUTED: {conv_note}
   Apply the penalty or boost indicated above.

3. [Content Relevance] — Slide coverage + Picture Superiority Effect (PSE)
   90-100 (Expert):    >80% slide keywords covered; added vivid descriptions beyond slide text; metaphors used
   75-89  (Competent): Core topic covered; 1-2 key facts skipped; some original explanation
   60-74  (Developing): <40% keyword match; speaker mirrors slides word-for-word (PSE violation)
   0-59   (Novice):    Completely off-topic or random words

   PSE PENALTY — Picture Superiority Effect:
   If the presenter only reads bullet points off a text-heavy slide word-for-word, cap Relevance at 70.
   PRE-COMPUTED: {pse_note}

4. [Delivery] — TED Conversational Golden Zone WPM + Punching Key Words
   95-100 (TED Golden):    160-190 WPM — TED conversational zone; varied pace; pauses for emphasis
   80-94  (Near Golden):   145-159 WPM — slightly below golden; mostly smooth
   72-85  (Slightly Fast): 191-215 WPM — could slow down at key moments
   60-78  (Noticeably Slow): 120-144 WPM — sounds hesitant; needs more confidence
   50-65  (Hesitant):      110-119 WPM — significant hesitation
   0-59   (Out of Zone):   <110 or >215 WPM — hard to follow

   PUNCHING KEY WORDS BONUS (+5 pts, cap 100):
   If the narrator shows clear emphasis before/after a key stat or conclusion
   (evidenced by natural phrasing like "and the number is...", "the answer is...",
   "here is the key...", or a dramatic single-word sentence), award +5.

   PRE-COMPUTED: {wpm_note}

══════════════════════════════════════════════
JAW-DROPPING MOMENT DETECTION
══════════════════════════════════════════════
A "Jaw-Dropping Moment" (from Talk Like TED) is when the speaker delivers:
  (a) A shocking statistic with clear emphasis
  (b) A vivid real-world story or case
  (c) A powerful one-sentence conclusion after a pause

PRE-COMPUTED HEURISTIC: {jaw_drop_note}

Set "jaw_dropping_moment": true ONLY if you can identify a specific quote in the narration
that is genuinely impactful, surprising, or emotionally resonant. Otherwise false.

══════════════════════════════════════════════
PRE-COMPUTED METRICS — USE AS HARD ANCHORS
══════════════════════════════════════════════
Total words spoken : {total_words}
WPM assessment     : {wpm_note}
Filler assessment  : {filler_note}

══════════════════════════════════════════════
SESSION CONTEXT
══════════════════════════════════════════════
Audience: {audience} | Scenario: {scenario} | Difficulty: {difficulty} | Challenge: {challenge_type}

══════════════════════════════════════════════
SLIDE CONTENT (on-screen text, what the audience sees)
══════════════════════════════════════════════
{slide_content_text}

══════════════════════════════════════════════
PRESENTER NARRATION (actual speech, slide by slide)
══════════════════════════════════════════════
{slide_narration_text}

══════════════════════════════════════════════
Q&A TRANSCRIPT
══════════════════════════════════════════════
{qa_history_text}

══════════════════════════════════════════════
MANDATORY OUTPUT RULES — READ ALL BEFORE WRITING
══════════════════════════════════════════════

⚡ SPEED RULE: Max 2 sentences per feedback item. Be concise. No essay text.

RULE 1 — LANGUAGE
Use IELTS 5.5-6.0 vocabulary only. Short, clear sentences. Write for a university student who is NOT a native English speaker. Say "show" not "demonstrate". Say "fix" not "mitigate". Say "use" not "utilise".

RULE 2 — QUOTE RULE (MOST IMPORTANT)
Every item MUST be grounded in the ACTUAL SLIDE CONTENT and ACTUAL PRESENTER NARRATION above.
NEVER invent quotes. NEVER write feedback that applies to any presenter generically.
If the user said something specific, quote EXACT words in double quotes. If they skipped a fact, name THAT EXACT FACT from the slide.

RULE 3 — 4×2 MATRIX (STRICTLY ENFORCED)
what_i_did_good      → EXACTLY 4 items, one per pillar: Structure, Fluency, Content Relevance, Delivery
areas_for_improvement → EXACTLY 4 items, one per pillar: Structure, Fluency, Content Relevance, Delivery
NO pillar may be skipped. NO extra items. Write in that exact order.

RULE 4 — what_i_did_good FORMAT — USE TED SENTENCE TEMPLATES
Each item must follow ONE of these TED-style templates:

  [Structure] Rule of Three / Message Map:
  "[Structure] The Rule of Three: You opened with a clear 3-point frame, saying: '[exact opening quote]'. This kept your audience focused and prevented information overload."

  [Fluency] Dialogue Style:
  "[Fluency] Conversational Delivery: Instead of reading a script, you used natural connectors like '[exact quote with connector]', which gave your talk a TED-style conversational feel."

  [Content Relevance] Picture Superiority:
  "[Content Relevance] Picture Superiority: Instead of just reading the slides, you used vivid language: '[exact descriptive quote]', which painted a mental image for your audience."

  [Delivery] Punching Key Words:
  "[Delivery] Punching Key Words: You emphasized a key point when you said: '[exact quote near stat or conclusion]', creating a natural pause effect that boosted impact."

  If the above templates don't match what the user DID (e.g., they didn't use Rule of Three),
  write an honest positive observation using: "[Pillar] Short Title: [exact quote]. [1 sentence why it works]."

RULE 5 — areas_for_improvement FORMAT — USE TED IMPROVEMENT TEMPLATES
Each object has EXACTLY these 4 keys. Use TED-style fix language:

  Structure weakness template:
  "issue": "[Structure] Weak Message Map: Your headline was not clear enough when transitioning between ideas."
  "example": "You said: '[exact quote at a weak transition]', which lost the logical focus for the audience."
  "how_to_fix": "Make it a Twitter-style headline (under 140 characters). Say this instead: 'Let's connect this to our main point: [rewritten example at IELTS 5.5 level].'"

  Fluency weakness template:
  "issue": "[Fluency] Script Reading over Dialogue: Your tone sounded like reading a paper rather than a TED-style talk."
  "example": "You literally read: '[exact mechanical quote from narration]', with no natural connectors."
  "how_to_fix": "Internalize your content and speak like a conversation. Say this instead: '[more natural, shorter version at IELTS 5.5 level].'"

  Relevance weakness template:
  "issue": "[Content Relevance] Picture Superiority Violation: You echoed the slide bullets instead of painting a picture."
  "example": "On Slide [N], the slide listed '[exact bullet from slide]', and you said almost the same words: '[exact mirrored quote]'."
  "how_to_fix": "Add a visual metaphor or real example. Say this instead: 'Think of it like [vivid analogy].'"

  Delivery weakness template:
  "issue": "[Delivery] [Speed verdict]: Your WPM of [X] is [outside/below/above] the TED golden zone of 160-190 WPM."
  "example": "You said: '[exact quote where speed was most noticeable]', which sounded [rushed/hesitant]."
  "how_to_fix": "Record yourself and aim for 160-190 WPM. Say this instead: '[same sentence at correct pacing with pauses marked as / ]'."

  If the above doesn't match, use the standard format:
  "dimension": "...", "issue": "...", "example": "You said: '...'" or "On Slide N, you skipped [fact].", "how_to_fix": "Say this instead: '...'"

RULE 6 — pitch_data
Return exactly 20 integers (range 50-250) simulating pitch variance for the Delivery chart.
TED golden zone WPM → high variance (180-230 range with swings). Low WPM → flat (90-120).

══════════════════════════════════════════════
Return ONLY valid JSON. No markdown fences. No text outside the JSON object.
══════════════════════════════════════════════
{{
  "radar_scores": {{
    "structure": <int 0-100>,
    "fluency":   <int 0-100>,
    "relevance": <int 0-100>,
    "delivery":  <int 0-100>
  }},
  "dimensions_info": {{
    "structure": {{"explanation": "<1 sentence: their actual transitions or Rule of Three usage>",
                   "calculation": "<1 sentence: how many slides had linking phrases + Rule of Three bonus if applied>"}},
    "fluency":   {{"explanation": "<1 sentence: filler count + dialogue vs script verdict>",
                   "calculation": "<1 sentence: filler rate + conversational markers count>"}},
    "relevance": {{"explanation": "<1 sentence: which slides had good or poor coverage + PSE verdict>",
                   "calculation": "<1 sentence: keyword match + PSE penalty if applied>"}},
    "delivery":  {{"explanation": "<1 sentence: exact WPM + TED golden zone verdict>",
                   "calculation": "<1 sentence: WPM vs 160-190 TED golden zone + punching bonus if applied>"}}
  }},
  "filler_log": [
    {{"word": "<exact filler from transcript>", "timestamp": "Slide N", "type": "Assistive or Disruptive"}}
  ],
  "what_i_did_good": [
    "[Structure] ...",
    "[Fluency] ...",
    "[Content Relevance] ...",
    "[Delivery] ..."
  ],
  "areas_for_improvement": [
    {{"dimension": "Structure",         "issue": "...", "example": "...", "how_to_fix": "Say this instead: \\"...\\""}},
    {{"dimension": "Fluency",           "issue": "...", "example": "...", "how_to_fix": "Say this instead: \\"...\\""}},
    {{"dimension": "Content Relevance", "issue": "...", "example": "...", "how_to_fix": "Say this instead: \\"...\\""}},
    {{"dimension": "Delivery",          "issue": "...", "example": "...", "how_to_fix": "Say this instead: \\"...\\""}}
  ],
  "pitch_data": [<20 integers each 50-250>],
  "jaw_dropping_moment": <true or false>
}}"""

    # ── Audit log: confirm real user data is being sent ───────────────────────
    app.logger.info(
        f"=== GEMINI PILLAR INPUT === "
        f"total_words={total_words} | wpm={wpm_note} | fillers={filler_count} | "
        f"slides={len(slides)} | qa_exchanges={len(fe_qa_history) if fe_qa_history else 0}"
    )
    app.logger.info(f"=== NARRATION PREVIEW === {all_narration_text[:300]!r}")

    try:
        response = _ai_client.chat.completions.create(
            model=EVAL_MODEL,
            max_tokens=MAX_TOKENS,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
        result = json.loads(raw)

        # Normalise: accept radar_scores (new schema) or scores (old schema)
        raw_scores = result.get("radar_scores") or result.get("scores") or {}
        result["scores"] = {
            k: int(min(100, max(0, float(v))))
            for k, v in raw_scores.items()
        }
        if not result["scores"]:
            result["scores"] = {"structure": 70, "fluency": 70, "relevance": 70, "delivery": 70}

        # Ensure dimensions_info exists for Jinja2 template
        if "dimensions_info" not in result:
            result["dimensions_info"] = {
                k: {"explanation": f"Score: {result['scores'].get(k, 70)}/100",
                    "calculation": "Based on 4-pillar rubric."}
                for k in ("structure", "fluency", "relevance", "delivery")
            }

        # Ensure pitch_data is valid (20 points)
        pd = result.get("pitch_data", [])
        if not isinstance(pd, list) or len(pd) < 5:
            result["pitch_data"] = _generate_pitch_data(
                result["scores"].get("delivery", 70), wpm_estimate
            )

        # Ensure filler_log exists
        if "filler_log" not in result:
            result["filler_log"] = []

        # Normalize what_i_did_good: Gemini sometimes returns dicts instead of the
        # requested strings. Convert any dict item to a formatted string so the
        # Jinja2 template only ever sees plain strings.
        def _norm_good_item(item):
            if not isinstance(item, dict):
                return str(item) if item else ""
            # Try every key name Gemini has been observed to use
            pillar = (item.get("pillar") or item.get("Pillar") or
                      item.get("dimension") or item.get("Dimension") or
                      item.get("category") or "")
            title  = (item.get("title") or item.get("Title") or
                      item.get("name") or item.get("Name") or "")
            obs    = (item.get("observation") or item.get("Observation") or
                      item.get("description") or item.get("Description") or
                      item.get("feedback") or item.get("text") or
                      item.get("content") or item.get("detail") or "")
            if not obs:
                # Last resort: join all non-empty string values
                obs = " | ".join(str(v) for v in item.values() if v and isinstance(v, str))
            if pillar and obs:
                prefix = f"[{pillar}]"
                if title:
                    return f"{prefix} {title}: {obs}"
                return f"{prefix} {obs}"
            return obs or str(item)

        raw_good = result.get("what_i_did_good") or []
        result["what_i_did_good"] = [_norm_good_item(x) for x in raw_good if x]

        # Normalise jaw_dropping_moment — Gemini may return it or we fall back to heuristic
        if "jaw_dropping_moment" in result:
            result["jaw_dropping_moment"] = bool(result["jaw_dropping_moment"])
        else:
            result["jaw_dropping_moment"] = jaw_drop_heuristic

        return result

    except Exception as e:
        err_type = type(e).__name__
        err_msg  = str(e)
        app.logger.error(
            f"[GEMINI EVAL ERROR] type={err_type} | msg={err_msg[:300]} | "
            f"model={EVAL_MODEL} | total_words={total_words}"
        )
        # Classify root cause for easier diagnosis in console
        if "timeout" in err_msg.lower() or "timed out" in err_msg.lower():
            app.logger.error("[GEMINI EVAL ERROR] Root cause: API TIMEOUT — relay took >120s")
        elif "json" in err_msg.lower() or isinstance(e, (ValueError, json.JSONDecodeError)):
            app.logger.error("[GEMINI EVAL ERROR] Root cause: JSON PARSE FAILURE — model returned non-JSON")
        elif "connection" in err_msg.lower() or "network" in err_msg.lower():
            app.logger.error("[GEMINI EVAL ERROR] Root cause: NETWORK / CONNECTION ERROR")
        else:
            app.logger.error(f"[GEMINI EVAL ERROR] Root cause: UNKNOWN — {err_type}")
        traceback.print_exc()
        # Return generic mock — zero-speech already handled by BRANCH A above
        return _mock_pillar_evaluation(difficulty)


def _mock_pillar_evaluation(difficulty):
    """Fallback mock when AI is unavailable. Returns the same schema as run_pillar_evaluation."""
    base = {"Easy": 74, "Medium": 61, "Hard": 49}.get(difficulty, 61)
    r = random.randint
    d_score = min(100, base + r(-6, 9))

    return {
        "scores": {
            "structure": min(100, base + r(-4, 10)),
            "fluency":   min(100, base + r(-10, 5)),
            "relevance": min(100, base + r(-8, 8)),
            "delivery":  d_score,
        },
        "dimensions_info": {
            "structure": {
                "explanation": "Measures how smoothly you connect different slides.",
                "calculation": "Computed by checking for transition words between slides.",
            },
            "fluency": {
                "explanation": "Measures filler words like 'um' or 'uh' and long pauses.",
                "calculation": "Computed by counting filler words per minute of speech.",
            },
            "relevance": {
                "explanation": "Measures how closely your words match the facts on the slide.",
                "calculation": "Computed by matching slide keywords with your speech text.",
            },
            "delivery": {
                "explanation": "Speaking speed in golden zone (120-150 WPM) gets the highest score.",
                "calculation": "Computed by dividing total words by your presentation time (WPM).",
            },
        },
        "filler_log": [],
        "jaw_dropping_moment": False,
        "pitch_data": _generate_pitch_data(d_score, 130),
        "what_i_did_good": [
            "[Structure] Rehearsal Completed: You moved through all slides in order and finished the session.",
            "[Fluency] Steady Pace: Your speech kept a consistent flow without major stops.",
            "[Content Relevance] Topic Covered: You addressed the main subject shown on the slides.",
            "[Delivery] Session Finished: You reached the end of the presentation without stopping early.",
        ],
        "areas_for_improvement": [
            {
                "dimension": "Structure",
                "issue": "No linking words detected between slides.",
                "example": "You moved between slides without using any transition phrases.",
                "how_to_fix": "Use a short bridge before each slide. Say this instead: \"Now let's look at the next point.\"",
            },
            {
                "dimension": "Fluency",
                "issue": "Some hesitations or filler words may have slowed your delivery.",
                "example": "Try to replace 'um' or 'uh' with a short silent pause.",
                "how_to_fix": "Close your mouth and breathe for one second. Say this instead: \"[pause] The next key idea is...\"",
            },
            {
                "dimension": "Content Relevance",
                "issue": "Slide content may not have been fully covered in your speech.",
                "example": "Some slides may have had key facts that were not mentioned in your narration.",
                "how_to_fix": "Before clicking Next Slide, check the screen. Say this instead: \"As you can see on the slide, [key fact].\"",
            },
            {
                "dimension": "Delivery",
                "issue": "Speaking speed may need adjustment for the golden zone.",
                "example": "Aim for 120-150 WPM for the clearest international presentation style.",
                "how_to_fix": "Practice reading your script aloud at a steady pace. Say this instead: \"[speak at 130 WPM — one word every 0.46 seconds].\"",
            },
        ],
    }


# ─────────────────────────────────────────────
# GEMINI: GENERATE TRAINING PLAN (Step 9)
# ─────────────────────────────────────────────
def generate_training_plan(pillar_eval, config, challenge_type):
    """Generate a personalized training plan based on 4-pillar scores."""
    difficulty = config.get("difficulty", "Medium")
    audience   = config.get("audience",   "Professor")
    scenario   = config.get("scenario",   "Class Presentation")

    scores = pillar_eval.get("scores", {})
    weakest = min(scores, key=scores.get) if scores else "structure"
    worst_score = scores.get(weakest, 60)
    areas = pillar_eval.get("areas_for_improvement", [])
    top_issues = " | ".join(a["issue"] for a in areas[:3]) if areas else "see score breakdown"

    if not AI_ENABLED:
        return _template_training_plan_v2(difficulty, audience, challenge_type, weakest, worst_score)

    prompt = f"""You are a world-class presentation coach. Write a short, practical training plan.

SESSION:
- Audience: {audience} | Scenario: {scenario} | Difficulty: {difficulty}
- Challenge Type: {challenge_type}

4-PILLAR SCORES (0-100):
{json.dumps(scores, indent=2)}
Weakest pillar: {weakest} ({worst_score}/100)

TOP ISSUES FOUND:
{top_issues}

Write the plan using these exact sections (markdown, simple IELTS 5.5 English):
## Training Plan

**Profile:** [one line summary]

### Fix First — {weakest.title()} ({worst_score}/100)
[2 specific drills. Reference {challenge_type} where useful.]

### Quick Wins This Week
[3 bullet points — easy actions to do right now]

### 30-Day Schedule
- Week 1: [focus]
- Week 2: [focus]
- Week 3: [focus]
- Week 4: [focus]

Keep every sentence short and simple. No jargon."""

    try:
        response = _ai_client.chat.completions.create(
            model=EVAL_MODEL,
            max_tokens=2048,
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        app.logger.error(f"Training plan generation failed: {e}")
        return _template_training_plan_v2(difficulty, audience, challenge_type, weakest, worst_score)


def _template_training_plan_v2(difficulty, audience, challenge_type, weakest, worst_score):
    return f"""## Training Plan

**Profile:** {difficulty} difficulty · {audience} audience · Weakest area: {weakest.title()} ({worst_score}/100)

### Fix First — {weakest.title()} ({worst_score}/100)
Record yourself and play it back. Look for the moments where you lose the point. Practice the Pyramid Principle: say your main idea first, then give 2-3 supporting facts. When you face a **{challenge_type}** question, answer in one clear sentence first, then explain.

### Quick Wins This Week
- Read your slide title out loud before you start talking about it.
- Use a linking phrase between each slide. Example: "Now let us move to the next point."
- Practice in front of a mirror for 10 minutes every day.

### 30-Day Schedule
- Week 1: Daily 10-min recording sessions. Focus on transitions and pace.
- Week 2: Two full rehearsals at {difficulty} difficulty. Ask a friend to give feedback.
- Week 3: Two Hard difficulty sessions. Practice handling tough questions.
- Week 4: One full mock session from start to finish. Record and review."""


# ─────────────────────────────────────────────
# WHISPER STUB: AUDIO TRANSCRIPTION (Step 4)
# ─────────────────────────────────────────────
def transcribe_audio(audio_file):
    """
    Standard transcription interface — ready to wire to OpenAI/Groq Whisper.
    audio_file: a file-like object (wav/mp3/webm)
    Returns: str transcript in English
    Usage (OpenAI Whisper):
        import openai
        client = openai.OpenAI(api_key=os.environ["OPENAI_API_KEY"])
        transcript = client.audio.transcriptions.create(
            model="whisper-1", file=audio_file, language="en"
        )
        return transcript.text
    Usage (Groq Whisper):
        import groq
        client = groq.Groq(api_key=os.environ["GROQ_API_KEY"])
        transcript = client.audio.transcriptions.create(
            model="whisper-large-v3", file=audio_file, language="en"
        )
        return transcript.text
    """
    raise NotImplementedError(
        "transcribe_audio: Connect to OpenAI or Groq Whisper API. See docstring for wiring instructions."
    )


# ─────────────────────────────────────────────
# PAGE ROUTES
# ─────────────────────────────────────────────

@app.route("/")
def index():
    session.clear()
    return render_template("index.html", ai_enabled=AI_ENABLED)


@app.route("/config")
def config_page():
    if "slide_key" not in session and "slides" not in session:
        return redirect(url_for("index"))
    return render_template("config.html", slides=_load_slides(session), ai_enabled=AI_ENABLED)


@app.route("/sandbox")
def sandbox():
    if "state" not in session:
        return redirect(url_for("index"))
    state = session["state"]
    slides = _load_slides(session)
    cfg = session.get("config", {})
    current_slide = next((s for s in slides if s["page"] == state["current_page"]), slides[0])
    has_file = bool(session.get("filepath") and os.path.exists(session.get("filepath", "")))
    return render_template(
        "sandbox.html",
        state=state,
        slides=slides,
        config=cfg,
        current_slide=current_slide,
        total_pages=len(slides),
        ai_enabled=AI_ENABLED,
        has_file=has_file,
    )


def _extract_pptx_shape_texts(shapes, out=None):
    """
    Recursively collect (top_emu, text) pairs from PPTX shapes.
    Handles groups (type 6), tables (type 19), and regular text shapes.
    """
    if out is None:
        out = []
    for shape in shapes:
        try:
            stype = int(shape.shape_type)
        except Exception:
            stype = -1
        if stype == 6:                          # GROUP — recurse
            try:
                _extract_pptx_shape_texts(shape.shapes, out)
            except Exception:
                pass
        elif stype == 19:                       # TABLE
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            out.append((shape.top or 0, t))
            except Exception:
                pass
        elif hasattr(shape, "text") and shape.text.strip():
            out.append((shape.top or 0, shape.text.strip()))
    return out


def _render_pptx_page(filepath, page_num):
    """
    Render a PPTX slide as an SVG (bytes, mimetype) tuple.
    SVG is served directly to the browser — uses system fonts so CJK text
    (Chinese / Japanese / Korean) renders correctly without needing embedded fonts.
    Returns (None, None) on failure.
    """
    try:
        from pptx import Presentation

        prs         = Presentation(filepath)
        slides_list = list(prs.slides)
        if page_num < 1 or page_num > len(slides_list):
            return None, None

        slide = slides_list[page_num - 1]
        total = len(slides_list)
        W, H  = 1280, 720

        CJK_FONTS = ("'PingFang SC','Microsoft YaHei','Noto Sans CJK SC',"
                     "Arial,Helvetica,sans-serif")

        def esc(s: str) -> str:
            return (s.replace("&", "&amp;").replace("<", "&lt;")
                     .replace(">", "&gt;").replace('"', "&quot;")
                     .replace("'", "&apos;"))

        # ── Title ────────────────────────────────────────────────────────────
        title_text  = ""
        title_shape = None
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text  = slide.shapes.title.text.strip()
                title_shape = slide.shapes.title
        except Exception:
            pass

        # ── Body text (all shapes, sorted top→bottom) ────────────────────────
        raw = _extract_pptx_shape_texts(slide.shapes)
        raw = [(top, txt) for (top, txt) in raw if txt != title_text]
        raw.sort(key=lambda x: x[0])

        # ── Build body <text> elements ───────────────────────────────────────
        body_els = []
        y = 148
        for _, block in raw:
            lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
            for line in lines:
                if y > H - 72:
                    break
                body_els.append(
                    f'<text x="68" y="{y}" font-size="19" fill="#acb2d8" '
                    f'font-family={CJK_FONTS!r}>'
                    f'&#9656;&#160;&#160;{esc(line[:115])}</text>'
                )
                y += 34
            y += 10

        badge = f"Slide {page_num} / {total}"

        svg = (
            f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{W}" height="{H}" viewBox="0 0 {W} {H}">'
            f'<defs>'
            f'  <linearGradient id="bg" x1="0" y1="0" x2="0.4" y2="1">'
            f'    <stop offset="0%" stop-color="#0c0c1c"/>'
            f'    <stop offset="100%" stop-color="#10102a"/>'
            f'  </linearGradient>'
            f'</defs>'
            f'<rect width="{W}" height="{H}" fill="url(#bg)"/>'
            f'<rect width="{W}" height="114" fill="#161638"/>'
            f'<rect x="0" y="112" width="{W}" height="2" fill="#4848c8"/>'
            f'<text x="52" y="76" font-size="32" font-weight="700" fill="#d8dcff" '
            f'font-family={CJK_FONTS!r}>{esc(title_text[:90])}</text>'
            + "".join(body_els)
            + f'<rect x="{W - 188}" y="{H - 48}" width="166" height="28" rx="5" fill="#1c1c44"/>'
            f'<text x="{W - 180}" y="{H - 28}" font-size="13" fill="#6870aa" '
            f'font-family="monospace,sans-serif">{esc(badge)}</text>'
            f"</svg>"
        )
        return svg.encode("utf-8"), "image/svg+xml"
    except Exception as e:
        app.logger.error(f"_render_pptx_page page={page_num} failed: {e}")
        return None, None


# ─────────────────────────────────────────────────────────────────────────────
# CONTENT QUALITY EVALUATION (Step 8c) — Per-slide 5-dimension transcript audit
# ─────────────────────────────────────────────────────────────────────────────
def run_content_quality_evaluation(fe_transcripts, slides, scene_slug, config):
    """
    Per-slide Content Quality audit — 5 fixed-order dimensions + optimized script.
    Returns: {"has_data": bool, "slide_transcripts_report": [...], "scene_label": str}
    """
    slide_map = {
        s.get("page", s.get("slide_id", 0)): s
        for s in (slides or [])
    }
    valid = [t for t in (fe_transcripts or []) if (t.get("text") or "").strip()]

    scene_label = {
        "thesis_defense":     "Thesis Defense",
        "case_pitch":         "MBA Case Pitch",
        "class_presentation": "Class Presentation",
    }.get(scene_slug or "", "Class Presentation")

    register_guide_en = {
        "thesis_defense":     (
            "Strict Academic Prose. Replace all informal words "
            "(get/a lot of/things/stuff) with academic alternatives "
            "(obtain/substantial/components/material). Zero tolerance for casual register."
        ),
        "case_pitch":         (
            "Executive Presence. Use precise finance and strategy vocabulary "
            "(ROI/synergy/scalability/market penetration/KPI-driven). "
            "Crisp, results-oriented, data-backed."
        ),
        "class_presentation": (
            "Vivid and confident tone. Optimise for clarity and engagement. "
            "Avoid excessive casual filler (like/you know/kind of)."
        ),
    }.get(scene_slug or "", "Balanced register, clear and fluent.")

    register_guide_zh = {
        "thesis_defense":     "学术论文答辩语体。严格Academic Prose，零容忍非正式词汇（get/a lot of/things），必须替换为学术词汇（obtain/substantial/components）。",
        "case_pitch":         "商业高管气场（Executive Presence）。使用精准财务/战略驱动词（ROI/synergy/scalability/KPI-driven），语言crisp、results-oriented。",
        "class_presentation": "允许适度生动表达，优化重点在生动自信。避免过度口语化（like/you know/kind of）。",
    }.get(scene_slug or "", "平衡语体，清晰流利即可。")

    if not valid:
        return {"has_data": False, "slide_transcripts_report": [], "scene_label": scene_label}

    # ── MOCK FALLBACK ────────────────────────────────────────────────────────
    if not AI_ENABLED:
        mock = []
        for idx, t in enumerate(valid):
            page = t.get("page", idx + 1)
            s    = slide_map.get(page, {})
            mock.append({
                "slide_id":    page,
                "slide_title": s.get("title", f"Slide {page}"),
                "raw_transcript": t["text"].strip(),
                "script_analysis": {
                    "grammar_fluency":     {"score": 68, "feedback_zh": "主谓一致基本正确，但出现了若干中式英语结构（Chinglish），如\"according to the data shows\"。建议修正为\"as the data shows\"。"},
                    "structural_logic":    {"score": 72, "feedback_zh": "有基本逻辑顺序，但未采用结论先行（BLUF）结构。建议先抛出核心论点，再展开论据支撑，而非直接罗列背景信息。"},
                    "transition_hook":     {"score": 65, "feedback_zh": "开头缺乏承接上一页内容的逻辑钩子，结尾也未为下一张PPT埋下伏笔。建议在开头加入\"Building on this point...\"等过渡表达。"},
                    "slide_alignment":     {"score": 75, "feedback_zh": "整体内容与PPT基本对应，但演讲时漏掉了PPT上的核心数据指标（Key Metrics），应主动在演讲中点出关键数字。"},
                    "vocabulary_register": {"score": 70, "feedback_zh": f"语体适配度中等。当前场景为{scene_label}：{register_guide_zh}"},
                },
                "optimized_script": (
                    "Building upon the foundation established in the previous section, "
                    "this slide presents critical findings that directly substantiate our central argument. "
                    "The evidence demonstrates a statistically significant correlation, "
                    "with key metrics revealing a 23% improvement over the established baseline. "
                    "In conclusion, these results validate our core hypothesis and naturally "
                    "lead us to the next dimension of our analysis."
                ),
            })
        return {"has_data": True, "slide_transcripts_report": mock, "scene_label": scene_label}

    # ── REAL AI PATH ─────────────────────────────────────────────────────────
    slides_text = ""
    for t in valid:
        page = t.get("page", 0)
        s    = slide_map.get(page, {})
        slides_text += (
            f"\n\n--- SLIDE {page}: {s.get('title', '(untitled)')} ---\n"
            f"PPT Elements (Key Points on Slide): {s.get('content', '(no content recorded)')}\n"
            f"User Transcript: {t['text'].strip()}"
        )

    prompt = f"""You are an expert English presentation coach for non-native speakers.

SCENE: {scene_label}
VOCABULARY REGISTER RULE: {register_guide_en}

Analyze EACH slide transcript below on exactly 5 dimensions IN THIS STRICT ORDER:

1. grammar_fluency — Diagnose tense errors, subject-verb agreement, singular/plural issues, and Chinglish expressions. Be specific.
2. structural_logic — Does the slide use BLUF (Bottom Line Up Front)? Is there a clear conclusion first, then supporting evidence? Identify any logic gaps.
3. transition_hook — Does the opening smoothly link to the previous slide? Does the closing hook into the next? MUST reference findings from structural_logic — if there is a logic gap, specify what "logic hook" phrase would smooth the transition.
4. slide_alignment — Cross-check against the PPT Elements. Did the user miss key metrics, go off-topic, or cover elements not on the slide?
5. vocabulary_register — Scene-aware audit. Apply the register rule strictly. Name specific words that should be upgraded and provide replacements.

For each slide, also write an `optimized_script` that:
- Fixes ALL grammar issues identified
- Uses BLUF structure (lead with the conclusion)
- Opens with a smooth transition hook from the previous slide's content
- Closes with a hook that naturally sets up the next slide
- Uses vocabulary precisely calibrated for {scene_label}

SLIDES TO ANALYZE:
{slides_text}

Return ONLY a valid JSON object — no markdown, no code fences, no comments:
{{
  "slides": [
    {{
      "slide_id": <integer page number>,
      "slide_title": "<string>",
      "script_analysis": {{
        "grammar_fluency":     {{"score": <0-100>, "feedback_zh": "<Chinese feedback 2-4 sentences>"}},
        "structural_logic":    {{"score": <0-100>, "feedback_zh": "<Chinese feedback 2-4 sentences>"}},
        "transition_hook":     {{"score": <0-100>, "feedback_zh": "<Chinese feedback 2-4 sentences>"}},
        "slide_alignment":     {{"score": <0-100>, "feedback_zh": "<Chinese feedback 2-4 sentences>"}},
        "vocabulary_register": {{"score": <0-100>, "feedback_zh": "<Chinese feedback 2-4 sentences>"}}
      }},
      "optimized_script": "<complete polished English script for this slide, 3-6 sentences>"
    }}
  ]
}}"""

    try:
        response = _ai_client.chat.completions.create(
            model=EVAL_MODEL,
            max_tokens=MAX_TOKENS,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$",           "", raw)

        # Multi-layer JSON repair — handles Chinese text with embedded quotes/newlines
        data = None
        parse_err = None
        for attempt, candidate in enumerate([
            raw,
            re.sub(r"(?<![\\])\\(?![\"\\\/bfnrtu])", r"\\\\", raw),  # fix bad escapes
        ]):
            try:
                data = json.loads(candidate)
                break
            except json.JSONDecodeError as e:
                parse_err = e
        if data is None and _HAS_JSON_REPAIR:
            try:
                repaired = _repair_json(raw, return_objects=True)
                if isinstance(repaired, dict):
                    data = repaired
            except Exception:
                pass
        if data is None:
            raise parse_err or ValueError("All JSON repair attempts failed")

        transcript_map = {t.get("page", 0): t.get("text", "").strip() for t in valid}
        result = []
        for sd in data.get("slides", []):
            sid = sd.get("slide_id") or sd.get("page") or 0
            try:
                sid = int(sid)
            except Exception:
                sid = 0
            result.append({
                "slide_id":         sid,
                "slide_title":      sd.get("slide_title", f"Slide {sid}"),
                "raw_transcript":   transcript_map.get(sid, ""),
                "script_analysis":  sd.get("script_analysis", {}),
                "optimized_script": sd.get("optimized_script", ""),
            })

        app.logger.info(f"[ContentQuality] scene={scene_slug} | slides_audited={len(result)}")
        return {"has_data": True, "slide_transcripts_report": result, "scene_label": scene_label}

    except Exception as e:
        app.logger.error(f"[ContentQuality] Error: {type(e).__name__}: {str(e)[:200]}")
        traceback.print_exc()
        return {"has_data": False, "slide_transcripts_report": [], "scene_label": scene_label, "error": str(e)}


@app.route("/x/slide-image/<int:page>")
def slide_image(page):
    """Serve a specific slide page as PNG — PDF via fitz, PPTX via Pillow."""
    from flask import Response as _R
    filepath = session.get("filepath", "")
    if not filepath or not os.path.exists(filepath):
        svg = (
            f'<svg xmlns="http://www.w3.org/2000/svg" width="800" height="500">'
            f'<rect width="800" height="500" fill="#12121f"/>'
            f'<text x="400" y="240" text-anchor="middle" fill="#6366f1" '
            f'font-size="22" font-family="monospace">Slide {page}</text>'
            f'<text x="400" y="275" text-anchor="middle" fill="#374151" '
            f'font-size="14" font-family="sans-serif">No file uploaded — mock mode</text>'
            f'</svg>'
        )
        return _R(svg, mimetype="image/svg+xml")

    ext = os.path.splitext(filepath)[1].lstrip(".").lower()

    # ── PPTX / PPT → SVG renderer ───────────────────────────────────────────
    if ext in ("pptx", "ppt"):
        svg_bytes, mime = _render_pptx_page(filepath, page)
        if svg_bytes is None:
            return "Page out of range or render failed", 404
        return _R(svg_bytes, mimetype=mime,
                  headers={"Cache-Control": "max-age=3600",
                           "X-Content-Type-Options": "nosniff"})

    # ── PDF → fitz renderer ─────────────────────────────────────────────────
    try:
        import fitz
        doc = fitz.open(filepath)
        if page < 1 or page > len(doc):
            doc.close()
            return "Page out of range", 404
        p   = doc[page - 1]
        mat = fitz.Matrix(2.0, 2.0)
        pix = p.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")
        doc.close()
        return _R(img_bytes, mimetype="image/png",
                  headers={"Cache-Control": "max-age=3600"})
    except Exception as e:
        app.logger.error(f"slide_image page={page} failed: {e}")
        return "Error generating image", 500


@app.route("/report")
def report():
    report_key = session.get("report_key")
    report_data = _load_report(report_key)
    if not report_data:
        # Cookie lost or file missing — restart from home
        app.logger.warning(f"[Report] No report data found for key={report_key!r} — redirecting home")
        return redirect(url_for("index"))
    return render_template(
        "report.html",
        evaluation=report_data.get("evaluation"),
        qa_bank=report_data.get("qa_bank", []),
        config=session.get("config", {}),
        answers=report_data.get("answers", []),
        ai_enabled=AI_ENABLED,
    )


# ─────────────────────────────────────────────
# API ROUTES
# ─────────────────────────────────────────────

@app.route("/x/upload", methods=["POST"])
def api_upload():
    """
    Step 1: Accept file → extract page images → Vision AI analysis.
    Full try/except with traceback so errors are always visible in logs.
    """
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400
        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400
        if not allowed_file(file.filename):
            return jsonify({"error": "Invalid file type. Please upload PDF, PPT, or PPTX."}), 400

        # Extract extension from the ORIGINAL filename BEFORE secure_filename()
        # strips non-ASCII chars (e.g. Chinese filenames lose their dot+ext).
        orig_ext = os.path.splitext(file.filename)[1].lstrip(".").lower()
        safe_stem = secure_filename(file.filename)
        # If secure_filename wiped everything (non-ASCII name), fall back to "upload"
        if not safe_stem or safe_stem == orig_ext:
            safe_stem = "upload"
        # Only append extension if secure_filename didn't already preserve it
        if orig_ext and not safe_stem.lower().endswith(f".{orig_ext}"):
            filename = f"{safe_stem}.{orig_ext}"
        else:
            filename = safe_stem

        os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
        save_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(save_path)
        app.logger.info(f"File saved: {save_path}")

        ext = orig_ext

        # Step 1a: PDF image extraction happens later in /x/start-session for Vision.
        # We only do text extraction here so the config page shows the real page count fast.

        # Build real (text-extracted) slide previews immediately — no AI call.
        # These give the config page the REAL page count + titles.
        # Vision analysis in /x/start-session will REPLACE these with richer AI output.
        # We NEVER store base64 images in the Flask session (4 KB cookie limit).
        if ext == "pdf":
            slides_preview = extract_pdf_text_slides(save_path)
        elif ext in ("ppt", "pptx"):
            ppt_data = extract_ppt_images_as_base64(save_path)
            slides_preview = pptx_to_slides(ppt_data) if ppt_data else MOCK_SLIDES
        else:
            slides_preview = MOCK_SLIDES

        session["slide_key"] = _save_slides(slides_preview)   # real page count, real text
        session["filename"] = filename
        session["filepath"] = save_path        # start-session re-reads from disk
        session["answers"]  = []
        session["qa_bank"]  = []

        page_count = len(slides_preview)
        return jsonify({
            "success":    True,
            "filename":   filename,
            "page_count": page_count,
            "ai_pending": AI_ENABLED and ext == "pdf",
            "message":    f"File uploaded. {page_count} page(s) ready for AI analysis.",
        })

    except Exception as e:
        traceback.print_exc()   # Full stack trace visible in Replit workflow logs
        app.logger.error(f"Upload failed: {e}")
        return jsonify({"error": str(e), "status": "error"}), 500


@app.route("/x/start-session", methods=["POST"])
def api_start_session():
    """
    Steps 2 → 3: Receive config → Vision analysis (if PDF) → Master Engine → sandbox.
    AI Vision is deferred here (not in /x/upload) to avoid proxy timeout on upload.
    """
    data = request.get_json() or {}
    audience   = data.get("audience",   "Professor")
    scenario   = data.get("scenario",   "Class Presentation")
    difficulty = data.get("difficulty", "Medium")

    # ── Step 2: Vision analysis — re-read file from disk (images never stored in cookie) ──
    filepath = session.get("filepath", "")
    filename = session.get("filename", "")
    # Start with whatever text-extracted slides were stored during upload (real page count)
    slides   = _load_slides(session)
    try:
        if filepath and os.path.exists(filepath) and AI_ENABLED:
            ext = filepath.rsplit(".", 1)[-1].lower()
            if ext == "pdf":
                app.logger.info(f"Re-extracting PDF images from {filepath}…")
                images_b64 = extract_pdf_images_as_base64(filepath)
                app.logger.info(f"Vision: analysing {len(images_b64)} page(s)…")
                ai_slides, ok = analyze_slides_with_claude(images_b64, filename)
                if ok and ai_slides:
                    # Vision may have analysed fewer pages than text extraction (e.g. max_pages cap
                    # or a short Vision response). Never let Vision shrink the deck — merge so that
                    # every page the user uploaded is represented in the session.
                    if len(ai_slides) < len(slides):
                        seen_pages = {s["page"] for s in ai_slides}
                        extra = [s for s in slides if s["page"] not in seen_pages]
                        slides = sorted(ai_slides + extra, key=lambda s: s["page"])
                        app.logger.info(
                            f"Vision enhanced {len(ai_slides)} page(s); "
                            f"merged with text-extracted to {len(slides)} total"
                        )
                    else:
                        slides = ai_slides
                        app.logger.info(f"Vision done: {len(slides)} slide(s) — AI-enhanced")
                else:
                    app.logger.info("Vision returned empty — keeping text-extracted slides")
            elif ext in ("ppt", "pptx"):
                ppt_data = extract_ppt_images_as_base64(filepath)
                if ppt_data:
                    slides = pptx_to_slides(ppt_data)
            session["slide_key"] = _save_slides(slides)
        else:
            app.logger.info(f"Skipping Vision (filepath={filepath!r}, AI_ENABLED={AI_ENABLED})")
            # slides already loaded from session above — re-save to ensure disk key is fresh
            if "slide_key" not in session:
                session["slide_key"] = _save_slides(slides)
    except Exception as e:
        traceback.print_exc()
        app.logger.error(f"Vision failed — keeping text-extracted slides: {e}")
        slides = _load_slides(session)
        session["slide_key"] = _save_slides(slides)

    # ── Step 3: Master Engine (Challenge seed + QA bank) ──────────────────────────
    try:
        challenge_seed, static_qa_bank = build_master_engine(
            slides, audience, scenario, difficulty
        )
    except Exception as e:
        traceback.print_exc()
        app.logger.error(f"Master engine failed, using mock: {e}")
        challenge_seed, static_qa_bank = MOCK_CHALLENGE, MOCK_QA_BANK

    session["config"] = {
        "audience":  audience,
        "scenario":  scenario,
        "difficulty": difficulty,
        "persona":   AUDIENCE_PERSONA.get(audience, AUDIENCE_PERSONA["Professor"]),
    }
    session["challenge_seed"] = challenge_seed
    if scenario in ("Class Presentation", "Academic Presentation"):
        # Dual-track: Q1 = AI-free, Q2 = anchor (when difficulty >= Medium).
        # "Academic Presentation" kept for backward compat with old sessions.
        _scene_for_qa = "class_presentation"
        session["qa_bank"] = build_dual_track_qa(slides, audience, _scene_for_qa, difficulty)
    elif scenario == "Thesis Defense":
        # Use DEFENSE_QUESTION_BANK; question count aligned with difficulty selection.
        # Each sampled template question is then fused with the user's actual thesis
        # topic/slide content via generate_custom_defense_question() — the AI examiner
        # never reads the raw template verbatim. answering_strategy stays anchored.
        import random as _rand
        _Q_COUNT = {"Easy": 3, "Medium": 5, "Hard": 8}
        _pool = list(DEFENSE_QUESTION_BANK)
        _rand.shuffle(_pool)
        _sampled = _pool[:_Q_COUNT.get(difficulty, 5)]
        session["qa_bank"] = customize_defense_qa_bank(_sampled, slides)
    else:
        session["qa_bank"] = []

    max_rounds = MAX_FOLLOWUP_ROUNDS.get(difficulty, 2)
    session["state"] = {
        "current_page":        1,
        "total_interruptions": 0,
        "in_qa_mode":          False,
        "follow_up_round":     0,
        "chat_history":        [],
        "max_follow_up_rounds": max_rounds,
    }

    return jsonify({"success": True, "redirect": "/sandbox"})


@app.route("/x/session-state", methods=["GET"])
def api_session_state():
    return jsonify({
        "state": session.get("state", {}),
        "config": session.get("config", {}),
        "slides": _load_slides(session),
    })


@app.route("/x/check-slide", methods=["POST"])
def api_check_slide():
    """
    Step 5: Core State Machine — user clicked 'Finish current slide'.
    Interrupt logic driven by difficulty-aware max_follow_up_rounds.
    """
    if "state" not in session:
        return jsonify({"error": "No active session"}), 400

    state = dict(session["state"])
    config = session.get("config", {})
    challenge_seed = session.get("challenge_seed") or MOCK_CHALLENGE
    slides = _load_slides(session)

    scenario = config.get("scenario", "Class Presentation")
    trigger_page = (challenge_seed or {}).get("trigger_page", 2)
    current_page = state["current_page"]

    req_data = request.get_json() or {}
    narration = req_data.get("narration", "")
    if narration:
        answers_list = list(session.get("answers", []))
        answers_list.append({"type": "narration", "page": current_page, "text": narration})
        session["answers"] = answers_list

    should_interrupt = (
        scenario in INTERRUPTABLE_SCENARIOS
        and current_page == trigger_page
        and state["total_interruptions"] < 1
        and not state["in_qa_mode"]
    )

    if should_interrupt:
        state["in_qa_mode"] = True
        state["follow_up_round"] = 1
        state["chat_history"] = []
        session["state"] = state
        c_type = challenge_seed.get("challenge_type", "")
        strategy = INTERRUPT_STRATEGIES.get(c_type, "")
        return jsonify({
            "action": "INTERRUPT",
            "question": challenge_seed["initial_challenge"],
            "challenge_type": c_type,
            "answering_strategy": strategy,
            "page": current_page,
        })

    next_page = current_page + 1
    if next_page > len(slides):
        # ── Academic / Class Presentation: trigger post-session Q&A before report ──
        if scenario in ("Academic Presentation", "Class Presentation", "Thesis Defense"):
            qa_bank = session.get("qa_bank", [])
            difficulty = config.get("difficulty", "Medium")
            if scenario == "Thesis Defense":
                qa_count = {"Easy": 3, "Medium": 5, "Hard": 8}.get(difficulty, 5)
            else:
                qa_count = {"Easy": 1, "Medium": 2, "Hard": 3}.get(difficulty, 2)
            questions = [q for q in qa_bank if isinstance(q, dict)][:qa_count]
            if questions:
                state["academic_qa_mode"] = True
                state["academic_qa_index"] = 0
                state["academic_qa_total"] = len(questions)
                session["state"] = state
                session["academic_qa_questions"] = questions
                first_q = questions[0]
                return jsonify({
                    "action":              "ACADEMIC_QA_START",
                    "question":            first_q.get("question", ""),
                    "question_num":        1,
                    "total_questions":     len(questions),
                    "challenge_type":      first_q.get("challenge_type", ""),
                    "category":            first_q.get("category", ""),
                    "answering_strategy":  first_q.get("answering_strategy", ""),
                })
        session["state"] = state
        return jsonify({"action": "PRESENTATION_DONE"})

    state["current_page"] = next_page
    state["in_qa_mode"] = False
    state["follow_up_round"] = 0
    state["chat_history"] = []
    session["state"] = state

    current_slide = next((s for s in slides if s["page"] == next_page), slides[0])
    return jsonify({"action": "NEXT_SLIDE", "page": next_page, "slide": current_slide})


@app.route("/x/submit-answer", methods=["POST"])
def api_submit_answer():
    """
    Step 6: Dynamic multi-round follow-up with Claude AI.
    Respects difficulty-controlled max_follow_up_rounds.
    """
    if "state" not in session:
        return jsonify({"error": "No active session"}), 400

    data = request.get_json()
    user_answer = data.get("answer", "")
    state = dict(session["state"])
    config = session.get("config", {})
    slides = _load_slides(session)
    challenge_seed = session.get("challenge_seed") or MOCK_CHALLENGE

    follow_up_round = state.get("follow_up_round", 1)
    max_rounds = state.get("max_follow_up_rounds", 2)

    chat_history = list(state.get("chat_history", []))
    chat_history.append({"role": "user", "content": user_answer})
    state["chat_history"] = chat_history

    answers_list = list(session.get("answers", []))
    answers_list.append({
        "type": "qa_answer",
        "page": state["current_page"],
        "round": follow_up_round,
        "text": user_answer,
        "question_id": f"interrupt_p{state['current_page']}_r{follow_up_round}",
    })
    session["answers"] = answers_list

    if follow_up_round < max_rounds:
        # Generate AI follow-up
        follow_up = generate_followup_question(
            slides=slides,
            audience=config.get("audience", "Professor"),
            difficulty=config.get("difficulty", "Medium"),
            challenge_type=challenge_seed.get("challenge_type", ""),
            chat_history=chat_history,
            user_answer=user_answer,
            current_page=state["current_page"],
        )

        chat_history.append({"role": "assistant", "content": follow_up})
        state["follow_up_round"] = follow_up_round + 1
        state["chat_history"] = chat_history
        session["state"] = state

        return jsonify({
            "status": "CONTINUE_QA",
            "next_question": follow_up,
            "round": state["follow_up_round"],
        })
    else:
        # QA finished
        state["in_qa_mode"] = False
        state["chat_history"] = []
        state["total_interruptions"] = state.get("total_interruptions", 0) + 1
        slides_list = _load_slides(session)
        next_page = state["current_page"] + 1
        if next_page <= len(slides_list):
            state["current_page"] = next_page
        state["follow_up_round"] = 0
        session["state"] = state

        return jsonify({
            "status": "QA_FINISHED",
            "message": "质询结束，请继续您的 Presentation。",
            "next_page": state["current_page"],
        })


@app.route("/x/submit-academic-qa", methods=["POST"])
def api_submit_academic_qa():
    """
    Step 7: Handle post-presentation Q&A for Class Presentation scenario.
    Cycles through qa_bank questions (count = 1/2/3 based on Easy/Medium/Hard).
    Returns ACADEMIC_QA_NEXT (more questions) or ACADEMIC_QA_DONE (all answered).
    """
    if "state" not in session:
        return jsonify({"error": "No active session"}), 400

    data     = request.get_json() or {}
    answer   = data.get("answer", "").strip()
    state    = dict(session["state"])
    config   = session.get("config", {})
    questions = session.get("academic_qa_questions", [])

    current_idx = state.get("academic_qa_index", 0)

    # Persist the answer
    current_q = questions[current_idx] if current_idx < len(questions) else {}
    answers_list = list(session.get("answers", []))
    answers_list.append({
        "type":            "academic_qa",
        "question_idx":    current_idx,
        "question_id":     current_q.get("id", f"q{current_idx + 1}"),
        "question":        current_q.get("question", ""),
        "question_type":   current_q.get("question_type",   "free"),   # "free" | "anchor"
        "anchor_type":     current_q.get("anchor_type",     ""),
        "target_dim":      current_q.get("target_dim",      ""),
        "scaffold_signal": current_q.get("scaffold_signal", ""),
        "answering_strategy": current_q.get("answering_strategy", ""),
        "text":            answer,
    })
    session["answers"] = answers_list

    next_idx = current_idx + 1
    if next_idx < len(questions):
        state["academic_qa_index"] = next_idx
        session["state"] = state
        next_q = questions[next_idx]
        return jsonify({
            "status":              "ACADEMIC_QA_NEXT",
            "question":            next_q.get("question", ""),
            "question_num":        next_idx + 1,
            "total_questions":     len(questions),
            "challenge_type":      next_q.get("challenge_type", ""),
            "category":            next_q.get("category", ""),
            "answering_strategy":  next_q.get("answering_strategy", ""),
        })
    else:
        state["academic_qa_mode"]  = False
        state["academic_qa_index"] = 0
        session["state"] = state
        return jsonify({"status": "ACADEMIC_QA_DONE"})


@app.route("/x/speechace-score", methods=["POST"])
def api_speechace_score():
    """
    Accepts a multipart POST with:
      - audio: audio blob (webm/ogg)
      - text:  user narration fallback (only used if no slide content found)
      - page:  slide page number (int)

    Reference text priority:
      1. slide['content'] from session (what the user should say) — authoritative
      2. user-submitted 'text' field (Web Speech transcript) — fallback
    This ensures Speechace scores pronunciation against the intended script,
    not the user's own speech (which would make all scores trivially high).
    """
    audio_file   = request.files.get("audio")
    user_text    = (request.form.get("text") or "").strip()
    page         = int(request.form.get("page") or 0)

    if not audio_file:
        return jsonify({"error": "No audio provided"}), 400

    # ── Resolve authoritative reference text from slide content ──────────────
    slides = _load_slides()
    slide_ref = ""
    for s in slides:
        if s.get("page") == page:
            title   = s.get("title", "")
            content = s.get("content", "")
            slide_ref = f"{title}. {content}".strip(" .") if title else content
            break

    ref_text = slide_ref or user_text   # prefer slide content; fall back to narration

    audio_bytes = audio_file.read()
    app.logger.info(
        f"[Speechace] Scoring slide {page} | audio={len(audio_bytes)}B | "
        f"ref_src={'slide' if slide_ref else 'user'} | ref_len={len(ref_text)} chars"
    )

    result = audio_engine.recognize_and_diagnose(
        audio_bytes, reference_text=ref_text, filename=audio_file.filename or "audio.webm"
    )
    result["page"] = page
    return jsonify(result)


@app.route("/x/tts-demo")
def api_tts_demo():
    """
    Returns MP3 audio for a pronunciation demo of the given word.
    Query param: word=<string>
    """
    word  = (request.args.get("word") or "").strip()
    voice = request.args.get("voice", "Stella")
    if not word:
        return jsonify({"error": "No word provided"}), 400

    app.logger.info(f"[CosyVoice] TTS demo: word={word!r} voice={voice}")
    mp3_bytes = audio_engine.generate_pronunciation_demo(word, voice=voice)
    if not mp3_bytes:
        return jsonify({"error": "TTS generation failed"}), 502

    from flask import Response
    return Response(mp3_bytes, mimetype="audio/mpeg",
                    headers={"Content-Disposition": f'inline; filename="{word}.mp3"'})


@app.route("/x/finish-presentation", methods=["POST"])
def api_finish_presentation():
    """
    Steps 7 & 8: Generate QA bank if needed + run 4-pillar AI evaluation + training plan.
    Accepts JSON body with frontend-collected real transcripts and QA history.
    """
    config         = session.get("config", {})
    answers        = list(session.get("answers", []))
    slides         = _load_slides(session)
    challenge_seed = session.get("challenge_seed") or MOCK_CHALLENGE
    scenario       = config.get("scenario", "Class Presentation")

    # ── Read real performance data sent by the frontend ────────────────────────
    req_data           = request.get_json(silent=True) or {}
    fe_transcripts     = req_data.get("presentation_transcripts", [])   # [{page,text,words}]
    fe_qa_history      = req_data.get("qa_chat_history", [])            # [{role,text,type}]
    total_time_seconds = int(req_data.get("total_time_seconds", 0) or 0)
    scene_slug         = req_data.get("scene",    None)                  # 'thesis_defense' | 'case_pitch' | 'class_presentation'
    pronunciation_data = req_data.get("pronunciation_data", {})          # {str(page): diagnostic_dict}

    # Merge frontend transcripts into session answers (fill gaps from voice/type)
    if fe_transcripts:
        existing_pages = {a["page"] for a in answers if a.get("type") == "narration"}
        for t in fe_transcripts:
            pg  = t.get("page", 0)
            txt = (t.get("text") or "").strip()
            if pg not in existing_pages and txt:
                answers.append({"type": "narration", "page": pg, "text": txt})
        session["answers"] = answers

    # Merge frontend QA history if it has more data than what was saved to session
    if fe_qa_history:
        user_qa = [h for h in fe_qa_history if h.get("role") == "user"]
        ai_qa   = [h for h in fe_qa_history if h.get("role") == "ai"]
        existing_qa_count = sum(
            1 for a in answers if a.get("type") in ("qa_answer", "academic_qa")
        )
        if len(user_qa) > existing_qa_count:
            for i, ua in enumerate(user_qa):
                ai_q = ai_qa[i]["text"] if i < len(ai_qa) else ""
                answers.append({
                    "type":     ua.get("type", "qa_answer"),
                    "question": ai_q,
                    "text":     ua.get("text", ""),
                })
            session["answers"] = answers

    # Step 7: QA bank — Thesis Defense generates it post-presentation
    qa_bank = session.get("qa_bank", [])
    if scenario == "Thesis Defense" and not qa_bank:
        _, qa_bank = build_master_engine(
            slides,
            config.get("audience", "Professor"),
            "Class Presentation",
            config.get("difficulty", "Medium"),
        )
        session["qa_bank"] = qa_bank

    # Steps 8 / 8b / 8c — run all three evaluations IN PARALLEL to cut latency
    # ── Log pronunciation data received from Speechace ─────────────────────────
    app.logger.info(
        f"[Pronunciation] Received data for {len(pronunciation_data)} slide(s)"
    )

    app.logger.info("[Eval] Launching pillar + CQ + ContentQuality in parallel…")
    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as _pool:
        _fut_pillar = _pool.submit(
            run_pillar_evaluation,
            slides, answers, config, challenge_seed,
            fe_qa_history, total_time_seconds,
        )
        _fut_cq = _pool.submit(
            run_communication_quality_evaluation,
            answers, config, fe_qa_history, scene_slug, total_time_seconds,
            slides, qa_bank,
        )
        _fut_cqual = _pool.submit(
            run_content_quality_evaluation,
            fe_transcripts, slides, scene_slug, config,
        )
        pillar_eval          = _fut_pillar.result()
        cq_eval              = _fut_cq.result()
        content_quality_eval = _fut_cqual.result()

    app.logger.info(
        f"[Eval] All done | "
        f"CQ has_data={cq_eval.get('has_data')} cq_total={cq_eval.get('cq_total')} | "
        f"CQual has_data={content_quality_eval.get('has_data')} "
        f"slides={len(content_quality_eval.get('slide_transcripts_report', []))}"
    )

    # ── Merge Speechace pronunciation_diagnostic into slide_transcripts_report ─
    # pronunciation_data keys are str(page_int) because JSON serialises integer
    # object-keys to strings: {1: diag} → {"1": diag}.
    if pronunciation_data:
        report = content_quality_eval.get("slide_transcripts_report") or []

        # If the report is empty (no transcripts submitted), build minimal stubs
        # so the pronunciation panel still appears in the report.
        if not report:
            for page_str in sorted(pronunciation_data.keys(),
                                   key=lambda x: int(x) if x.isdigit() else 0):
                diag = pronunciation_data[page_str]
                if not isinstance(diag, dict):
                    continue
                try:
                    page_int = int(page_str)
                except ValueError:
                    continue
                s = next((sl for sl in slides if sl.get("page") == page_int), {})
                report.append({
                    "slide_id":    page_int,
                    "slide_title": s.get("title", f"Slide {page_int}"),
                    "raw_transcript": diag.get("transcript", "(no transcript recorded)"),
                    "script_analysis": {},
                    "optimized_script": "",
                    "pronunciation_diagnostic": diag,
                })
                app.logger.info(
                    f"[Pronunciation] Stub slide {page_int} | "
                    f"score={diag.get('overall_score')} errors={len(diag.get('error_list', []))}"
                )
            if report:
                content_quality_eval["slide_transcripts_report"] = report
                content_quality_eval["has_data"] = True
        else:
            # Report already has slides — merge pronunciation into matching slides
            for slide in report:
                page_key = str(slide.get("slide_id", ""))
                diag = pronunciation_data.get(page_key)
                if diag and isinstance(diag, dict):
                    slide["pronunciation_diagnostic"] = diag
                    app.logger.info(
                        f"[Pronunciation] Merged diag for slide {page_key} | "
                        f"score={diag.get('overall_score')} errors={len(diag.get('error_list', []))}"
                    )

    challenge_type = challenge_seed.get("challenge_type", "General")

    # Step 9: Training plan
    training_plan = generate_training_plan(pillar_eval, config, challenge_type)

    evaluation = {
        "pillar":                pillar_eval,
        "communication_quality": cq_eval,
        "content_quality":       content_quality_eval,
        "training_plan":         training_plan,
        "scenario":              scenario,
        "audience":              config.get("audience",   "Professor"),
        "difficulty":            config.get("difficulty", "Medium"),
        "challenge_type":        challenge_type,
        "ai_powered":            AI_ENABLED,
    }

    # ── Persist large blobs to disk; store only UUID key in cookie ────────────
    # Prevents Flask session cookie overflow (4KB limit).
    report_key = _save_report(evaluation, qa_bank, answers)
    session["report_key"] = report_key
    # Remove stale large keys that would push the cookie over limit
    session.pop("evaluation", None)
    session.pop("qa_bank", None)
    session.pop("answers", None)
    app.logger.info(f"[Report] Saved report data → key={report_key}")
    return jsonify({"success": True, "redirect": "/report"})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 8000))
    app.run(host="0.0.0.0", port=port, debug=True)
